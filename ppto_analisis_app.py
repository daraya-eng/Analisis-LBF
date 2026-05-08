"""
App Análisis Presupuesto 2026 — Seguimiento PPTO vs Real
Ejecutar: python ppto_analisis_app.py
Abrir:    http://localhost:8052
"""

import math
import os
import pickle
import datetime
import pandas as pd
import dash
from dash import dcc, html, dash_table, Input, Output, State
from dash.exceptions import PreventUpdate
import plotly.graph_objects as go

try:
    import pyodbc
except ImportError:
    pyodbc = None


# ─── Fecha dinámica ────────────────────────────────────────────────────────────
_HOY = datetime.date.today()
_MES_ACT = _HOY.month
_ANO_ACT = _HOY.year

_MESES_NOMBRE = {
    1: 'Enero', 2: 'Febrero', 3: 'Marzo', 4: 'Abril',
    5: 'Mayo', 6: 'Junio', 7: 'Julio', 8: 'Agosto',
    9: 'Septiembre', 10: 'Octubre', 11: 'Noviembre', 12: 'Diciembre'
}

# ─── Conexión ──────────────────────────────────────────────────────────────────
try:
    from db_config import CONN_STR
except ImportError:
    CONN_STR = os.environ.get("DB_CONN_STR") or None  # Fallback a variable de entorno

def get_conn():
    if not CONN_STR or pyodbc is None:
        raise RuntimeError("Sin acceso a BD. Requiere db_config.py y conexión VPN.")
    return pyodbc.connect(CONN_STR, timeout=30)

# ─── Caché local ───────────────────────────────────────────────────────────────
DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data")
os.makedirs(DATA_DIR, exist_ok=True)

_CACHE_NAMES = ["categoria", "zona", "zona_cat", "cliente", "producto", "desalineacion", "clientes_caida", "precios", "incremental", "resumen_ud_grupo", "ppto_vs_venta"]

def _save_cache(name: str, df: pd.DataFrame):
    path = os.path.join(DATA_DIR, f"df_{name}.pkl")
    with open(path, "wb") as f:
        pickle.dump(df, f)

def _load_cache(name: str) -> pd.DataFrame:
    path = os.path.join(DATA_DIR, f"df_{name}.pkl")
    if os.path.exists(path):
        try:
            with open(path, "rb") as f:
                return pickle.load(f)
        except Exception as e:
            print(f"[WARN] Caché {name} corrupto: {e}")
    return pd.DataFrame()

def _save_kpis_cache(kpis: dict):
    import json
    path = os.path.join(DATA_DIR, "pvv_kpis.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump(kpis, f)

def _load_kpis_cache() -> dict:
    import json
    path = os.path.join(DATA_DIR, "pvv_kpis.json")
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            print(f"[WARN] pvv_kpis.json corrupto: {e}")
    return {}


def _ventas_cte():
    """CTE 'ventas_completas': DW_TOTAL_FACTURA + guías pendientes no facturadas.
    Las guías se deduplicam (1 por día por guia_num+part_code) y se excluyen
    si su guia_num ya aparece en DW_TOTAL_FACTURA (ya fue facturada).
    """
    return """
    ventas_completas AS (
        SELECT VENDEDOR, RUT, CODIGO, CATEGORIA, ANO, MES,
               CAST(VENTA AS float) AS VENTA, CAST(CANT AS float) AS CANT
        FROM ventas_completas
        UNION ALL
        SELECT v.VENDEDOR,
               g.cust_code  AS RUT,
               g.part_code  AS CODIGO,
               g.categoria  AS CATEGORIA,
               YEAR(g.fecha) AS ANO,
               MONTH(g.fecha) AS MES,
               CAST(g.ext_price_amt AS float) AS VENTA,
               0             AS CANT
        FROM (
            SELECT guia_num, part_code, cust_code, categoria, vendedor,
                   MAX(ext_price_amt) AS ext_price_amt,
                   MAX(fecha)         AS fecha
            FROM vw_guias_por_facturar
            WHERE ext_price_amt > 0 AND part_code <> 'SIN'
            GROUP BY guia_num, part_code, cust_code, categoria, vendedor
        ) g
        JOIN (SELECT DISTINCT VENDEDOR FROM ventas_completas) v
             ON SUBSTRING(v.VENDEDOR, CHARINDEX('-', v.VENDEDOR) + 1,
                          LEN(v.VENDEDOR)) = g.vendedor
        WHERE NOT EXISTS (
            SELECT 1 FROM ventas_completas f WHERE f.GUIA_NUM = g.guia_num
        )
    )"""

# ─── Filtros de ventas (vendedores/códigos que ensucian la data) ───────────────
_VEND_EXCLUIR = (
    "89-FACTURACION MUESTRA Y U OBSEQU",
    "90-FACTURACION USO INTERNO",
    "96-FACTURACION FALTANTES",
    "97-DONACIONES",
    "98-FACTURACION OTROS CONCEPTOS",
    "99-FACTURACION MERMAS",
)
_COD_EXCLUIR = ("FLETE", "NINV", "SIN")

# Fragmento WHERE para agregar a SELECT FROM DW_TOTAL_FACTURA (prefijo AND)
_DW_FILTRO = (
    "VENDEDOR NOT IN ("
    + ",".join(f"'{v}'" for v in _VEND_EXCLUIR)
    + ") AND ISNULL(CODIGO,'') NOT IN ("
    + ",".join(f"'{c}'" for c in _COD_EXCLUIR)
    + ",'') "
)

# ─── Remapeo de categorías ─────────────────────────────────────────────────────
# Categorías que se fusionan en otra: Servicios → EQM
_CAT_ALIAS = {"Servicios": "EQM"}

# Expresión SQL para normalizar categoría en [PPTO 2026] (alias p)
_SQL_CAT_PPTO = (
    "CASE WHEN LTRIM(RTRIM(ISNULL(p.[CATEGORÍA 2026],''))) = 'Servicios' "
    "THEN 'EQM' ELSE ISNULL(LTRIM(RTRIM(p.[CATEGORÍA 2026])),'(sin cat)') END"
)
# Expresión SQL para normalizar CATEGORIA en DW_TOTAL_FACTURA
_SQL_CAT_DW = (
    "CASE WHEN ISNULL(CATEGORIA,'') = 'Servicios' "
    "THEN 'EQM' ELSE ISNULL(CATEGORIA,'(sin cat)') END"
)
# Versión sin fallback '(sin cat)' para CTEs intermedias
_SQL_CAT_DW_RAW = (
    "CASE WHEN ISNULL(CATEGORIA,'') = 'Servicios' THEN 'EQM' ELSE ISNULL(CATEGORIA,'') END"
)

def _cat_in(cat: str) -> str:
    """Genera cláusula IN para filtrar por categoría incluyendo aliases.
    Ej: 'EQM' → \"IN ('EQM','Servicios')\"
    """
    cats = [cat] + [src for src, dst in _CAT_ALIAS.items() if dst == cat and src != cat]
    return "IN (" + ",".join(f"'{c}'" for c in cats) + ")"

def _sql_ppto_ytd(p_alias='p', zp_alias='zp'):
    """Expresión SQL para ppto_ytd respetando que 'PRODUCTOS NUEVOS' arranca en Mayo."""
    base = f"TRY_CAST(REPLACE(ISNULL({p_alias}.[PPTO 2026],'0'),',','.') AS float)"
    factor_nuevos = f"({_MES_ACT} - 4.0) / 8.0" if _MES_ACT >= 5 else "0.0"
    return (
        f"CASE WHEN LTRIM(RTRIM(ISNULL({p_alias}.DESCRIPCION,''))) = 'PRODUCTOS NUEVOS' "
        f"THEN {base} * {factor_nuevos} "
        f"ELSE {base} * ISNULL({zp_alias}.peso_ytd, {_MES_ACT}/12.0) "
        f"END"
    )

# ─── Helpers de formato ────────────────────────────────────────────────────────
def _fmt(n):
    if n is None or (isinstance(n, float) and math.isnan(n)):
        return "—"
    n = float(n)
    if abs(n) >= 1e9:
        return f"${n/1e9:.1f}MM"
    if abs(n) >= 1e6:
        return f"${n/1e6:.0f}M"
    return f"${n:,.0f}"

def _fmt_abs(n):
    if n is None or (isinstance(n, float) and math.isnan(n)):
        return "—"
    v = float(n)
    if v < 0:
        return f"-${abs(v):,.0f}".replace(",", ".")
    return f"${v:,.0f}".replace(",", ".")

def _fmt_pct(n):
    if n is None or (isinstance(n, float) and math.isnan(n)):
        return "—"
    return f"{float(n):.1f}%"

# ─── Estilos ───────────────────────────────────────────────────────────────────
CARD_STYLE = {
    "background": "white",
    "borderRadius": "8px",
    "padding": "16px",
    "marginBottom": "16px",
    "boxShadow": "0 1px 4px rgba(0,0,0,0.08)",
}

# ─── Helpers UI ───────────────────────────────────────────────────────────────
def _kpi_card(titulo, valor, color="#1F3864", subtitulo=""):
    return html.Div(style={
        "background": color, "color": "white", "borderRadius": "8px",
        "padding": "12px 16px", "minWidth": "150px", "flex": "1",
    }, children=[
        html.Div(titulo, style={"fontSize": "11px", "opacity": "0.85", "marginBottom": "4px"}),
        html.Div(valor, style={"fontSize": "22px", "fontWeight": "700"}),
        html.Div(subtitulo, style={"fontSize": "10px", "opacity": "0.75", "marginTop": "2px"}),
    ])

def _seccion(titulo, color="#1F3864"):
    return html.H3(titulo, style={
        "color": color, "fontSize": "14px", "fontWeight": "700",
        "marginBottom": "12px", "marginTop": "0"
    })

def _sem(v, umbrales=(80, 50)):
    """Semáforo: verde >= u[0], amarillo >= u[1], rojo < u[1]."""
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return "⚫"
    if v >= umbrales[0]:
        return "🟢"
    if v >= umbrales[1]:
        return "🟡"
    return "🔴"

def _sem_precio(v, umbrales=(5, -5)):
    """Semáforo precio: verde si cerca del PPTO, amarillo si >+5%, rojo si < -5%."""
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return "⚫"
    if v >= umbrales[0]:
        return "🟢"
    if v >= umbrales[1]:
        return "🟡"
    return "🔴"

def _kpi_row(cards):
    return html.Div(style={"display": "flex", "gap": "12px", "marginBottom": "16px",
                            "flexWrap": "wrap"}, children=cards)

def _empty_msg(msg="Sin datos disponibles"):
    return html.Div(msg, style={"color": "#888", "padding": "32px", "textAlign": "center",
                                 "fontSize": "14px"})

# ─── Globals de datos ──────────────────────────────────────────────────────────
_df_categoria      = pd.DataFrame()
_df_zona           = pd.DataFrame()
_df_zona_cat       = pd.DataFrame()
_df_cliente        = pd.DataFrame()
_df_desalineacion  = pd.DataFrame()
_df_producto       = pd.DataFrame()
_df_clientes_caida   = pd.DataFrame()
_df_precios          = pd.DataFrame()
_df_incremental      = pd.DataFrame()
_df_resumen_ud_grupo = pd.DataFrame()
_df_ppto_vs_venta    = pd.DataFrame()
_pvv_kpis: dict      = {}
_last_update         = "—"

# ─── Funciones de carga ────────────────────────────────────────────────────────

def _load_categoria():
    """Carga análisis por categoría: PPTO vs real."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            fact_cat AS (
                SELECT {_SQL_CAT_DW} AS categoria,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_facturas,
                       SUM(CASE WHEN ANO = 2025 AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(CANT AS float) ELSE 0 END) AS cant_2026_ytd
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO} AND CATEGORIA IS NOT NULL AND CATEGORIA <> ''
                GROUP BY {_SQL_CAT_DW}
            ),
            guias_cat AS (
                SELECT CASE WHEN ISNULL(g.categoria,'') = 'Servicios' THEN 'EQM'
                            ELSE ISNULL(g.categoria,'(sin cat)') END AS categoria,
                       SUM(CASE WHEN YEAR(g.fecha) = {_ANO_ACT} AND MONTH(g.fecha) <= {_MES_ACT}
                                THEN CAST(g.ext_price_amt AS float) ELSE 0 END) AS venta_guias
                FROM (
                    SELECT guia_num, part_code, cust_code, categoria, vendedor,
                           MAX(ext_price_amt) AS ext_price_amt, MAX(fecha) AS fecha
                    FROM vw_guias_por_facturar
                    WHERE ext_price_amt > 0 AND part_code NOT IN ('SIN','FLETE','NINV')
                    GROUP BY guia_num, part_code, cust_code, categoria, vendedor
                ) g
                WHERE NOT EXISTS (SELECT 1 FROM DW_TOTAL_FACTURA f2 WHERE f2.GUIA_NUM = g.guia_num)
                  AND g.categoria IS NOT NULL AND g.categoria <> ''
                GROUP BY CASE WHEN ISNULL(g.categoria,'') = 'Servicios' THEN 'EQM'
                              ELSE ISNULL(g.categoria,'(sin cat)') END
            ),
            ppto_cat AS (
                SELECT {_SQL_CAT_PPTO} AS categoria,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                           * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0)) AS ppto_ytd,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[CANT 2026],'0'),',','.') AS float)) AS cant_ppto,
                       COUNT(DISTINCT NULLIF(LTRIM(RTRIM(ISNULL(p.RUT,''))), '')) AS clientes_ppto,
                       COUNT(DISTINCT NULLIF(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '')) AS productos_ppto
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') NOT IN ('','SIN')
                GROUP BY {_SQL_CAT_PPTO}
            ),
            ppto_incr_cat AS (
                SELECT {_SQL_CAT_PPTO} AS categoria,
                       SUM(CASE WHEN LTRIM(RTRIM(ISNULL(p.DESCRIPCION,''))) <> 'PRODUCTOS NUEVOS'
                                THEN TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) ELSE 0 END) AS ppto_incr_anual,
                       SUM(CASE WHEN LTRIM(RTRIM(ISNULL(p.DESCRIPCION,''))) = 'PRODUCTOS NUEVOS'
                                THEN TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) ELSE 0 END) AS ppto_nuevos_anual,
                       SUM(CASE WHEN LTRIM(RTRIM(ISNULL(p.DESCRIPCION,''))) <> 'PRODUCTOS NUEVOS'
                                THEN TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                                     * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0) ELSE 0 END) AS ppto_incr_ytd,
                       SUM(CASE WHEN LTRIM(RTRIM(ISNULL(p.DESCRIPCION,''))) = 'PRODUCTOS NUEVOS'
                                THEN TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                                     * {("({} - 4.0) / 8.0".format(_MES_ACT)) if _MES_ACT >= 5 else "0.0"} ELSE 0 END) AS ppto_nuevos_ytd
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') IN ('','SIN')
                GROUP BY {_SQL_CAT_PPTO}
            )
            SELECT p.categoria,
                   p.ppto_anual, p.ppto_ytd, p.cant_ppto, p.clientes_ppto, p.productos_ppto,
                   ISNULL(f.venta_facturas, 0)  AS venta_facturas,
                   ISNULL(g.venta_guias, 0)     AS venta_guias,
                   ISNULL(f.venta_facturas, 0) + ISNULL(g.venta_guias, 0) AS venta_2026_ytd,
                   ISNULL(f.venta_2025_ytd, 0) AS venta_2025_ytd,
                   ISNULL(f.cant_2026_ytd, 0)  AS cant_2026_ytd,
                   ISNULL(ic.ppto_incr_anual, 0)  AS ppto_incr_anual,
                   ISNULL(ic.ppto_nuevos_anual, 0) AS ppto_nuevos_anual,
                   ISNULL(ic.ppto_incr_ytd, 0)    AS ppto_incr_ytd,
                   ISNULL(ic.ppto_nuevos_ytd, 0)  AS ppto_nuevos_ytd
            FROM ppto_cat p
            LEFT JOIN fact_cat f ON f.categoria = p.categoria
            LEFT JOIN guias_cat g ON g.categoria = p.categoria
            LEFT JOIN ppto_incr_cat ic ON ic.categoria = p.categoria
            ORDER BY p.ppto_anual DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "cant_ppto", "venta_facturas", "venta_guias",
                  "venta_2026_ytd", "venta_2025_ytd", "cant_2026_ytd", "ppto_ytd",
                  "ppto_incr_anual", "ppto_nuevos_anual", "ppto_incr_ytd", "ppto_nuevos_ytd"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        # ppto_total = c/cód + incr + nuevos (total presupuesto completo)
        df["ppto_total"] = df["ppto_anual"] + df["ppto_incr_anual"] + df["ppto_nuevos_anual"]
        # total_ppto_ytd = c/cód YTD + incr YTD + nuevos YTD (nuevos = 0 antes de Mayo)
        df["total_ppto_ytd"] = df["ppto_ytd"] + df["ppto_incr_ytd"] + df["ppto_nuevos_ytd"]
        df["cumpl_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] / r["total_ppto_ytd"] * 100) if r["total_ppto_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["gap"] = df["venta_2026_ytd"] - df["total_ppto_ytd"]
        df["crec_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] - r["venta_2025_ytd"]) / r["venta_2025_ytd"] * 100
            if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_categoria] {e}")
        return pd.DataFrame()


def _load_zona():
    """Carga análisis por zona."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            guias_zona AS (
                SELECT v.VENDEDOR AS zona,
                       SUM(CASE WHEN YEAR(g.fecha) = {_ANO_ACT} AND MONTH(g.fecha) <= {_MES_ACT}
                                THEN CAST(g.ext_price_amt AS float) ELSE 0 END) AS venta_guias
                FROM (
                    SELECT guia_num, part_code, cust_code, categoria, vendedor,
                           MAX(ext_price_amt) AS ext_price_amt, MAX(fecha) AS fecha
                    FROM vw_guias_por_facturar
                    WHERE ext_price_amt > 0 AND part_code NOT IN ('SIN','FLETE','NINV')
                    GROUP BY guia_num, part_code, cust_code, categoria, vendedor
                ) g
                JOIN (SELECT DISTINCT VENDEDOR FROM DW_TOTAL_FACTURA) v
                     ON SUBSTRING(v.VENDEDOR, CHARINDEX('-', v.VENDEDOR)+1, LEN(v.VENDEDOR)) = g.vendedor
                WHERE NOT EXISTS (SELECT 1 FROM DW_TOTAL_FACTURA f2 WHERE f2.GUIA_NUM = g.guia_num)
                GROUP BY v.VENDEDOR
            ),
            ppto_zona AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                           * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0)) AS ppto_ytd,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[CANT 2026],'0'),',','.') AS float)) AS cant_ppto
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL))
            ),
            clientes_zona_26 AS (
                SELECT DISTINCT RUT, VENDEDOR AS zona
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                  AND VENDEDOR NOT IN ({",".join(f"'{v}'" for v in _VEND_EXCLUIR)})
            ),
            venta_26_zona AS (
                SELECT VENDEDOR AS zona,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_facturas,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(CANT  AS float) ELSE 0 END) AS cant_2026_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT} AND {_DW_FILTRO}
                GROUP BY VENDEDOR
            ),
            venta_25_zona AS (
                SELECT cz.zona,
                       SUM(CAST(f.VENTA AS float)) AS venta_2025_full,
                       SUM(CASE WHEN f.MES <= {_MES_ACT} THEN CAST(f.VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA f
                INNER JOIN clientes_zona_26 cz ON cz.RUT = f.RUT
                WHERE f.ANO = 2025
                  AND ISNULL(f.CODIGO,'') NOT IN ({",".join(f"'{c}'" for c in _COD_EXCLUIR)},'')
                GROUP BY cz.zona
            )
            SELECT p.zona, p.ppto_anual, p.ppto_ytd, p.cant_ppto,
                   ISNULL(v26.venta_facturas, 0) AS venta_facturas,
                   ISNULL(gz.venta_guias, 0)     AS venta_guias,
                   ISNULL(v26.venta_facturas, 0) + ISNULL(gz.venta_guias, 0) AS venta_2026_ytd,
                   ISNULL(v25.venta_2025_full, 0) AS venta_2025_full,
                   ISNULL(v25.venta_2025_ytd, 0)  AS venta_2025_ytd,
                   ISNULL(v26.cant_2026_ytd, 0)   AS cant_2026_ytd
            FROM ppto_zona p
            LEFT JOIN venta_26_zona v26 ON v26.zona = p.zona
            LEFT JOIN venta_25_zona v25 ON v25.zona = p.zona
            LEFT JOIN guias_zona gz     ON gz.zona  = p.zona
            ORDER BY p.ppto_anual DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "cant_ppto", "venta_facturas", "venta_guias",
                  "venta_2026_ytd", "venta_2025_full", "venta_2025_ytd", "cant_2026_ytd", "ppto_ytd"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["cumpl_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] / r["ppto_ytd"] * 100) if r["ppto_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["gap"] = df["venta_2026_ytd"] - df["ppto_ytd"]
        df["crec_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] - r["venta_2025_ytd"]) / r["venta_2025_ytd"] * 100
            if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["proyeccion_anual"] = df.apply(
            lambda r: (r["venta_2026_ytd"] / _MES_ACT * 12) if _MES_ACT > 0 else 0,
            axis=1
        )
        df["gap_vs_ppto_anual"] = df["proyeccion_anual"] - df["ppto_anual"]
        return df
    except Exception as e:
        print(f"[ERROR _load_zona] {e}")
        return pd.DataFrame()


def _load_zona_cat():
    """Carga análisis por zona × categoría."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            ventas_completas AS (
                SELECT VENDEDOR, CATEGORIA, ANO, MES,
                       CAST(VENTA AS float) AS VENTA
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO}
            ),
            guias_pendientes AS (
                SELECT v.VENDEDOR, g.categoria AS CATEGORIA, YEAR(g.fecha) AS ANO, MONTH(g.fecha) AS MES,
                       CAST(g.ext_price_amt AS float) AS VENTA
                FROM (
                    SELECT guia_num, part_code, cust_code, categoria, vendedor,
                           MAX(ext_price_amt) AS ext_price_amt, MAX(fecha) AS fecha
                    FROM vw_guias_por_facturar
                    WHERE ext_price_amt > 0 AND part_code NOT IN ('SIN','FLETE','NINV')
                    GROUP BY guia_num, part_code, cust_code, categoria, vendedor
                ) g
                JOIN (SELECT DISTINCT VENDEDOR FROM DW_TOTAL_FACTURA) v
                     ON SUBSTRING(v.VENDEDOR, CHARINDEX('-', v.VENDEDOR)+1, LEN(v.VENDEDOR)) = g.vendedor
                WHERE NOT EXISTS (
                    SELECT 1 FROM DW_TOTAL_FACTURA f2 WHERE f2.GUIA_NUM = g.guia_num
                )
            ),
            ppto_cat AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       {_SQL_CAT_PPTO} AS categoria,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                           * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0)) AS ppto_ytd
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),
                         {_SQL_CAT_PPTO}
            ),
            clientes_zona_26 AS (
                SELECT DISTINCT RUT, VENDEDOR AS zona
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                  AND VENDEDOR NOT IN ({",".join(f"'{v}'" for v in _VEND_EXCLUIR)})
            ),
            real_cat_26 AS (
                SELECT VENDEDOR AS zona,
                       {_SQL_CAT_DW_RAW} AS categoria,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN VENTA ELSE 0 END) AS venta_2026_ytd
                FROM ventas_completas
                WHERE ANO = {_ANO_ACT}
                GROUP BY VENDEDOR, {_SQL_CAT_DW_RAW}
            ),
            real_cat_25 AS (
                SELECT cz.zona,
                       CASE WHEN ISNULL(f.CATEGORIA,'') = 'Servicios' THEN 'EQM'
                            ELSE ISNULL(f.CATEGORIA,'') END AS categoria,
                       SUM(CASE WHEN f.MES <= {_MES_ACT} THEN CAST(f.VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA f
                INNER JOIN clientes_zona_26 cz ON cz.RUT = f.RUT
                WHERE f.ANO = 2025
                  AND ISNULL(f.CODIGO,'') NOT IN ({",".join(f"'{c}'" for c in _COD_EXCLUIR)},'')
                GROUP BY cz.zona, CASE WHEN ISNULL(f.CATEGORIA,'') = 'Servicios' THEN 'EQM'
                                       ELSE ISNULL(f.CATEGORIA,'') END
            )
            SELECT COALESCE(p.zona, r26.zona) AS zona,
                   COALESCE(p.categoria, r26.categoria) AS categoria,
                   ISNULL(p.ppto_anual, 0) AS ppto_anual,
                   ISNULL(p.ppto_ytd, 0) AS ppto_ytd,
                   ISNULL(r26.venta_2026_ytd, 0) AS venta_2026_ytd,
                   ISNULL(r25.venta_2025_ytd, 0) AS venta_2025_ytd
            FROM ppto_cat p
            FULL OUTER JOIN real_cat_26 r26 ON r26.zona = p.zona AND r26.categoria = p.categoria
            LEFT JOIN real_cat_25 r25 ON r25.zona = COALESCE(p.zona, r26.zona)
                                     AND r25.categoria = COALESCE(p.categoria, r26.categoria)
            ORDER BY ISNULL(p.ppto_anual, 0) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "ppto_ytd", "venta_2026_ytd", "venta_2025_ytd"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        return df
    except Exception as e:
        print(f"[ERROR _load_zona_cat] {e}")
        return pd.DataFrame()


def _load_cliente():
    """Carga análisis por cliente."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            ventas_completas AS (
                SELECT VENDEDOR, RUT, CODIGO, CATEGORIA, ANO, MES,
                       CAST(VENTA AS float) AS VENTA, CAST(CANT AS float) AS CANT
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO}
            ),
            ppto_cli AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT,''))) IN ('','0')
                            THEN '(sin rut)' ELSE p.RUT END AS RUT,
                       MAX(CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT,''))) IN ('','0')
                                THEN 'Presupuesto sin Cliente asignado'
                                ELSE p.NOMBRE END) AS nombre,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                           * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0)) AS ppto_ytd,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[CANT 2026],'0'),',','.') AS float)) AS cant_ppto,
                       COUNT(DISTINCT NULLIF(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '')) AS productos_ppto
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),
                         CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT,''))) IN ('','0')
                              THEN '(sin rut)' ELSE p.RUT END
            ),
            real_cli AS (
                SELECT RUT,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026_ytd,
                       SUM(CASE WHEN ANO = 2025 AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd,
                       SUM(CASE WHEN ANO = 2025 THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_full,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(CANT AS float) ELSE 0 END) AS cant_2026_ytd
                FROM ventas_completas
                GROUP BY RUT
            )
            SELECT p.zona, p.RUT, LEFT(p.nombre, 50) AS nombre,
                   p.ppto_anual, p.ppto_ytd, p.cant_ppto, p.productos_ppto,
                   ISNULL(r.venta_2026_ytd, 0) AS venta_2026_ytd,
                   ISNULL(r.venta_2025_ytd, 0) AS venta_2025_ytd,
                   ISNULL(r.venta_2025_full, 0) AS venta_2025_full,
                   ISNULL(r.cant_2026_ytd, 0) AS cant_2026_ytd
            FROM ppto_cli p
            LEFT JOIN real_cli r ON r.RUT = p.RUT
            ORDER BY p.ppto_anual DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "cant_ppto", "venta_2026_ytd", "venta_2025_ytd",
                  "venta_2025_full", "cant_2026_ytd", "ppto_ytd"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["cumpl_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] / r["ppto_ytd"] * 100) if r["ppto_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["gap"] = df["venta_2026_ytd"] - df["ppto_ytd"]
        df["crec_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] - r["venta_2025_ytd"]) / r["venta_2025_ytd"] * 100
            if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_cliente] {e}")
        return pd.DataFrame()


def _load_clientes_caida():
    """Clientes cuya venta 2026 YTD es menor que 2025 YTD, agrupados por zona 2026."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH clientes_26 AS (
                SELECT f.RUT,
                       f.VENDEDOR AS zona,
                       MAX(LEFT(ISNULL(f.NOMBRE, f.RUT), 50)) AS nombre,
                       SUM(CASE WHEN f.MES <= {_MES_ACT} THEN CAST(f.VENTA AS float) ELSE 0 END) AS venta_2026_ytd
                FROM DW_TOTAL_FACTURA f
                WHERE f.ANO = {_ANO_ACT}
                  AND {_DW_FILTRO}
                  AND ISNULL(f.RUT, '') NOT IN ('', '0')
                GROUP BY f.RUT, f.VENDEDOR
            ),
            clientes_25 AS (
                SELECT f.RUT,
                       SUM(CASE WHEN f.MES <= {_MES_ACT} THEN CAST(f.VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA f
                WHERE f.ANO = 2025
                  AND {_DW_FILTRO}
                  AND ISNULL(f.RUT, '') NOT IN ('', '0')
                GROUP BY f.RUT
            ),
            ppto_cli AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT, ''))) IN ('', '0')
                            THEN '(sin rut)' ELSE LTRIM(RTRIM(p.RUT)) END AS RUT,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026], '0'), ',', '.') AS float)) AS ppto_anual
                FROM [PPTO 2026] p
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026], '0'), ',', '.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),
                         CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT, ''))) IN ('', '0')
                              THEN '(sin rut)' ELSE LTRIM(RTRIM(p.RUT)) END
            )
            SELECT c26.zona,
                   c26.RUT,
                   c26.nombre,
                   ISNULL(pp.ppto_anual, 0)         AS ppto_anual,
                   c26.venta_2026_ytd,
                   ISNULL(c25.venta_2025_ytd, 0)    AS venta_2025_ytd,
                   c26.venta_2026_ytd
                       - ISNULL(c25.venta_2025_ytd, 0) AS diferencia
            FROM clientes_26 c26
            LEFT JOIN clientes_25 c25 ON c25.RUT = c26.RUT
            LEFT JOIN ppto_cli    pp  ON pp.RUT  = c26.RUT AND pp.zona = c26.zona
            WHERE c26.venta_2026_ytd < ISNULL(c25.venta_2025_ytd, 0)
              AND ISNULL(c25.venta_2025_ytd, 0) > 0
            ORDER BY diferencia ASC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "venta_2026_ytd", "venta_2025_ytd", "diferencia"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["caida_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] - r["venta_2025_ytd"]) / r["venta_2025_ytd"] * 100
            if r["venta_2025_ytd"] > 0 else float("nan"), axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_clientes_caida] {e}")
        return pd.DataFrame()


def _load_producto():
    """Carga análisis por producto × cliente."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH
            ventas_completas AS (
                SELECT VENDEDOR, RUT, CODIGO, CATEGORIA, ANO, MES,
                       CAST(VENTA AS float) AS VENTA, CAST(CANT AS float) AS CANT
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO}
            ),
            ppto_prod AS (
                SELECT ISNULL(LTRIM(RTRIM(VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       RUT,
                       CODIGO,
                       MAX(DESCRIPCION) AS descripcion,
                       ISNULL(LTRIM(RTRIM([CATEGORÍA 2026])), '') AS categoria,
                       SUM(TRY_CAST(REPLACE(ISNULL([CANT 2026],'0'),',','.') AS float)) AS cant_ppto,
                       SUM(TRY_CAST(REPLACE(ISNULL([PRECIO 2026],'0'),',','.') AS float) * TRY_CAST(REPLACE(ISNULL([CANT 2026],'0'),',','.') AS float)) AS ppto_prod,
                       AVG(TRY_CAST(REPLACE(ISNULL([PRECIO 2026],'0'),',','.') AS float)) AS precio_ppto
                FROM [PPTO 2026]
                WHERE TRY_CAST(REPLACE(ISNULL([PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(CODIGO, '') <> '' AND ISNULL(CODIGO, '') <> 'SIN'
                GROUP BY LTRIM(RTRIM(VENDEDOR_ACTUAL)), RUT, CODIGO, LTRIM(RTRIM([CATEGORÍA 2026]))
            ),
            real_prod AS (
                SELECT VENDEDOR AS zona, CODIGO,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026_ytd,
                       SUM(CASE WHEN ANO = 2025 AND MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT} THEN CAST(CANT AS float) ELSE 0 END) AS cant_2026_ytd,
                       SUM(CASE WHEN ANO = 2025 AND MES <= {_MES_ACT} THEN CAST(CANT AS float) ELSE 0 END) AS cant_2025_ytd
                FROM ventas_completas
                WHERE ISNULL(CODIGO,'') <> '' AND ISNULL(CODIGO,'') <> 'SIN'
                GROUP BY VENDEDOR, CODIGO
            )
            SELECT TOP 500
                   p.zona, p.RUT, p.CODIGO, LEFT(p.descripcion, 60) AS descripcion, p.categoria,
                   p.cant_ppto, p.precio_ppto, p.ppto_prod,
                   ISNULL(r.venta_2026_ytd, 0) AS venta_2026_ytd,
                   ISNULL(r.venta_2025_ytd, 0) AS venta_2025_ytd,
                   ISNULL(r.cant_2026_ytd, 0) AS cant_2026_ytd,
                   ISNULL(r.cant_2025_ytd, 0) AS cant_2025_ytd
            FROM ppto_prod p
            LEFT JOIN real_prod r ON r.CODIGO = p.CODIGO AND r.zona = p.zona
            ORDER BY p.ppto_prod DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["cant_ppto", "precio_ppto", "ppto_prod", "venta_2026_ytd",
                  "venta_2025_ytd", "cant_2026_ytd", "cant_2025_ytd"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["precio_real_2026"] = df.apply(
            lambda r: r["venta_2026_ytd"] / r["cant_2026_ytd"] if r["cant_2026_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["delta_precio_pct"] = df.apply(
            lambda r: (r["precio_real_2026"] - r["precio_ppto"]) / r["precio_ppto"] * 100
            if (not math.isnan(r["precio_real_2026"]) and r["precio_ppto"] > 0) else float("nan"),
            axis=1
        )
        df["gap_cant"] = df["cant_2026_ytd"] - df["cant_ppto"] / 12 * _MES_ACT
        df["gap_venta"] = df["venta_2026_ytd"] - df["ppto_prod"] / 12 * _MES_ACT
        return df
    except Exception as e:
        print(f"[ERROR _load_producto] {e}")
        return pd.DataFrame()


_last_desa_error = ""

def _load_desalineacion():
    """Detalle del presupuesto 2026 por Categoría × Nombre × Situación (solo [PPTO 2026])."""
    global _last_desa_error
    _last_desa_error = ""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            base AS (
                SELECT {_SQL_CAT_PPTO}                                                 AS categoria,
                       ISNULL(LTRIM(RTRIM(p.[SITUACI\u00d3N ])),'(sin situaci\u00f3n)') AS situacion,
                       ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),'(sin zona)')             AS zona,
                       CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT,''))) IN ('','0')
                            THEN NULL ELSE p.RUT END AS rut,
                       NULLIF(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') AS codigo,
                       TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) AS ppto_linea,
                       {_sql_ppto_ytd()} AS ppto_ytd_linea
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp
                       ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
            )
            SELECT categoria, situacion,
                   SUM(ppto_linea)        AS ppto_anual,
                   SUM(ppto_ytd_linea)    AS ppto_ytd,
                   COUNT(DISTINCT zona)   AS n_zonas,
                   COUNT(DISTINCT rut)    AS n_clientes,
                   COUNT(DISTINCT codigo) AS n_productos
            FROM base
            GROUP BY categoria, situacion
            ORDER BY categoria, SUM(ppto_linea) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "ppto_ytd", "n_zonas", "n_clientes", "n_productos"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        tot = df["ppto_anual"].sum()
        df["pct_total"]     = df["ppto_anual"] / tot * 100 if tot > 0 else 0.0
        df["ppto_anual_cat"] = df.groupby("categoria")["ppto_anual"].transform("sum")
        df["ppto_ytd_cat"]   = df.groupby("categoria")["ppto_ytd"].transform("sum")
        df["pct_cat"] = df.apply(
            lambda r: r["ppto_anual"] / r["ppto_anual_cat"] * 100 if r["ppto_anual_cat"] > 0 else 0.0,
            axis=1
        )
        print(f"[INFO] _load_desalineacion OK — {len(df)} filas")
        return df
    except Exception as e:
        _last_desa_error = str(e)
        print(f"[ERROR _load_desalineacion] {e}")
        return pd.DataFrame()


def _load_precios_analisis():
    """Carga análisis de precio 2026 vs 2025 a nivel producto × zona(2026) × cliente.
    La zona es siempre la estructura 2026. La venta 2025 se trae por RUT+CODIGO
    sin importar en qué zona estaba el cliente en 2025 (estructura diferente).
    Para productos sin historial 2025, usa precio PPTO 2026 como referencia.
    """
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH
            -- Ventas 2026: zona según estructura 2026 (VENDEDOR actual)
            v2026 AS (
                SELECT
                    LTRIM(RTRIM(ISNULL(VENDEDOR,''))) AS zona,
                    LTRIM(RTRIM(ISNULL(RUT,'')))      AS rut,
                    LTRIM(RTRIM(ISNULL(CODIGO,'')))   AS codigo,
                    {_SQL_CAT_DW_RAW}                  AS categoria,
                    SUM(CAST(VENTA AS float))          AS venta_2026,
                    SUM(CAST(CANT  AS float))          AS cant_2026
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO}
                  AND ANO = {_ANO_ACT} AND MES <= {_MES_ACT}
                  AND TRY_CAST(CANT  AS float) > 0
                  AND TRY_CAST(VENTA AS float) > 0
                  AND ISNULL(RUT,'')    NOT IN ('','0')
                  AND ISNULL(CODIGO,'') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(ISNULL(VENDEDOR,''))),
                         LTRIM(RTRIM(ISNULL(RUT,''))),
                         LTRIM(RTRIM(ISNULL(CODIGO,''))),
                         {_SQL_CAT_DW_RAW}
            ),
            -- Ventas 2025: agrupadas solo por RUT+CODIGO, sin zona
            -- (la zona 2025 no coincide con la estructura 2026)
            v2025 AS (
                SELECT
                    LTRIM(RTRIM(ISNULL(RUT,'')))    AS rut,
                    LTRIM(RTRIM(ISNULL(CODIGO,''))) AS codigo,
                    SUM(CAST(VENTA AS float))        AS venta_2025,
                    SUM(CAST(CANT  AS float))        AS cant_2025
                FROM DW_TOTAL_FACTURA
                WHERE {_DW_FILTRO}
                  AND ANO = 2025 AND MES <= {_MES_ACT}
                  AND TRY_CAST(CANT  AS float) > 0
                  AND TRY_CAST(VENTA AS float) > 0
                  AND ISNULL(RUT,'')    NOT IN ('','0')
                  AND ISNULL(CODIGO,'') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(ISNULL(RUT,''))),
                         LTRIM(RTRIM(ISNULL(CODIGO,'')))
            ),
            -- Referencia PPTO 2026: descripción, nombre cliente y precio presupuestado
            ppto_ref AS (
                SELECT
                    ISNULL(LTRIM(RTRIM(VENDEDOR_ACTUAL)), '') AS zona,
                    ISNULL(LTRIM(RTRIM(RUT)), '')             AS rut,
                    ISNULL(LTRIM(RTRIM(CODIGO)), '')          AS codigo,
                    MAX(LEFT(ISNULL(DESCRIPCION,''), 60))     AS descripcion,
                    MAX(LEFT(ISNULL(NOMBRE,''), 50))          AS nombre_cliente,
                    AVG(TRY_CAST(REPLACE(ISNULL([PRECIO 2026],'0'),',','.') AS float)) AS precio_ppto
                FROM [PPTO 2026]
                WHERE ISNULL(LTRIM(RTRIM(CODIGO)),'') NOT IN ('','SIN')
                  AND TRY_CAST(REPLACE(ISNULL([PRECIO 2026],'0'),',','.') AS float) > 0
                GROUP BY ISNULL(LTRIM(RTRIM(VENDEDOR_ACTUAL)),''),
                         ISNULL(LTRIM(RTRIM(RUT)),''),
                         ISNULL(LTRIM(RTRIM(CODIGO)),'')
            )
            SELECT TOP 3000
                a.zona,
                a.rut,
                ISNULL(p.nombre_cliente, a.rut)  AS nombre_cliente,
                a.codigo,
                ISNULL(p.descripcion, a.codigo)  AS descripcion,
                a.categoria,
                a.venta_2026,
                a.cant_2026,
                ISNULL(h.venta_2025, 0)          AS venta_2025,
                ISNULL(h.cant_2025,  0)          AS cant_2025,
                CASE WHEN a.cant_2026 > 0
                     THEN a.venta_2026 / a.cant_2026 ELSE NULL END AS precio_2026,
                -- Referencia de precio: real 2025 si existe, sino PPTO 2026
                CASE
                    WHEN ISNULL(h.cant_2025, 0) > 0 THEN h.venta_2025 / h.cant_2025
                    WHEN ISNULL(p.precio_ppto, 0) > 0 THEN p.precio_ppto
                    ELSE NULL
                END AS precio_2025,
                CASE
                    WHEN ISNULL(h.cant_2025, 0) > 0 THEN 'real'
                    WHEN ISNULL(p.precio_ppto, 0) > 0 THEN 'ppto'
                    ELSE 'sin_ref'
                END AS ref_precio
            FROM v2026 a
            -- venta 2025 se une solo por RUT+CODIGO (ignora zona 2025)
            LEFT JOIN v2025 h  ON h.rut = a.rut AND h.codigo = a.codigo
            -- PPTO se une por zona+RUT+CODIGO para traer descripción y precio ppto
            LEFT JOIN ppto_ref p ON p.zona = a.zona AND p.rut = a.rut AND p.codigo = a.codigo
            ORDER BY a.venta_2026 DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["venta_2026", "cant_2026", "venta_2025", "cant_2025"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        for c in ["precio_2026", "precio_2025"]:
            df[c] = pd.to_numeric(df[c], errors="coerce")

        # Efecto precio: cuánto cambia la venta 2026 por solo el cambio de precio
        # precio_2025 puede ser precio real 2025 O precio PPTO 2026 (para productos nuevos)
        df["delta_precio"]     = df["precio_2026"] - df["precio_2025"]
        df["delta_precio_pct"] = df.apply(
            lambda r: r["delta_precio"] / r["precio_2025"] * 100
            if (pd.notna(r["precio_2025"]) and r["precio_2025"] > 0 and pd.notna(r["delta_precio"]))
            else float("nan"), axis=1
        )
        # Efecto precio = Δprecio × cant_2026
        df["efecto_precio"]    = df.apply(
            lambda r: r["delta_precio"] * r["cant_2026"]
            if pd.notna(r["delta_precio"]) and r["cant_2026"] > 0
            else float("nan"), axis=1
        )
        # Efecto volumen: solo aplica cuando hay historial real 2025 (ref_precio == 'real')
        df["efecto_volumen"]   = df.apply(
            lambda r: (r["cant_2026"] - r["cant_2025"]) * r["precio_2025"]
            if (r.get("ref_precio") == "real" and pd.notna(r["precio_2025"]) and r["precio_2025"] > 0)
            else float("nan"), axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_precios_analisis] {e}")
        return pd.DataFrame()


def _load_ppto_incremental():
    """Carga presupuesto incremental: filas de [PPTO 2026] sin CODIGO o sin RUT.
    Estos son presupuestos adicionales no asociados a producto/cliente específico.
    Se agrupan por categoría + zona + descripción para análisis aparte.
    """
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            SELECT
                {_SQL_CAT_PPTO}                                      AS categoria,
                ISNULL(LTRIM(RTRIM(VENDEDOR_ACTUAL)), '(sin zona)')  AS zona,
                ISNULL(LTRIM(RTRIM(DESCRIPCION)), '(sin descripción)') AS descripcion,
                ISNULL(LTRIM(RTRIM(NOMBRE)), '')                      AS nombre_cliente,
                ISNULL(LTRIM(RTRIM(RUT)), '')                         AS rut,
                SUM(TRY_CAST(REPLACE(ISNULL([PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                SUM(TRY_CAST(REPLACE(ISNULL([CANT 2026],'0'),',','.') AS float)) AS cant_ppto,
                AVG(TRY_CAST(REPLACE(ISNULL([PRECIO 2026],'0'),',','.') AS float)) AS precio_ppto
            FROM [PPTO 2026] p
            WHERE TRY_CAST(REPLACE(ISNULL([PPTO 2026],'0'),',','.') AS float) > 0
              AND (
                  ISNULL(LTRIM(RTRIM(ISNULL(CODIGO,''))), '') IN ('','SIN')
                  OR ISNULL(LTRIM(RTRIM(ISNULL(RUT,''))), '') IN ('','0')
              )
            GROUP BY {_SQL_CAT_PPTO},
                     ISNULL(LTRIM(RTRIM(VENDEDOR_ACTUAL)), '(sin zona)'),
                     ISNULL(LTRIM(RTRIM(DESCRIPCION)), '(sin descripción)'),
                     ISNULL(LTRIM(RTRIM(NOMBRE)), ''),
                     ISNULL(LTRIM(RTRIM(RUT)), '')
            ORDER BY SUM(TRY_CAST(REPLACE(ISNULL([PPTO 2026],'0'),',','.') AS float)) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "cant_ppto", "precio_ppto"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        return df
    except Exception as e:
        print(f"[ERROR _load_ppto_incremental] {e}")
        return pd.DataFrame()


def _load_ppto_vs_venta():
    """Lee PPTO_VS_VENTA y KPIs desde VW_RESUMEN_KPIS_DASHBOARD."""
    global _pvv_kpis
    try:
        conn = get_conn()
        cur = conn.cursor()

        # KPIs globales desde la vista (pre-calculados igual que Power BI)
        cur.execute("SELECT * FROM VW_RESUMEN_KPIS_DASHBOARD")
        kpi_cols  = [d[0].strip() for d in cur.description]
        kpi_row   = cur.fetchone()
        if kpi_row:
            kd = dict(zip(kpi_cols, kpi_row))
            _pvv_kpis = {
                "ppto_total":      float(kd.get("PPTO_TOTAL_2026", 0) or 0),
                "ppto_trazable":   float(kd.get("PPTO_TRAZABLE_2026", 0) or 0),
                "ppto_sin_cliente":float(kd.get("PPTO_SIN_CLIENTE_2026", 0) or 0),
                "ppto_incr_sin_pn":float(kd.get("PPTO_INCREMENTAL_SIN_PN_2026", 0) or 0),
                "ppto_incr_con_pn":float(kd.get("PPTO_INCREMENTAL_CON_PN_2026", 0) or 0),
                "venta_ytd":       float(kd.get("VENTA_YTD_2026", 0) or 0),
                "venta_ytd_25":    float(kd.get("VENTA_YTD_2025", 0) or 0),
                "meta_ytd":        float(kd.get("META_YTD_2026", 0) or 0),
                "alcance_ytd":     float(kd.get("ALCANCE_YTD", 0) or 0),
                "cumpl_ppto":      float(kd.get("CUMPLIMIENTO_PPTO_TOTAL", 0) or 0),
                "cumpl_trazable":  float(kd.get("CUMPLIMIENTO_TRAZABLE", 0) or 0),
                "gap_meta":        float(kd.get("GAP_META_YTD", 0) or 0),
                "gap_ppto":        float(kd.get("GAP_PPTO_TOTAL", 0) or 0),
                "var_abs_25":      float(kd.get("VAR_ABS_VS_2025", 0) or 0),
            }

        # Leer PPTO_VS_VENTA completa (VENTA_2026 ya incluye guías vía BI_TOTAL_FACTURA)
        cur.execute("SELECT * FROM PPTO_VS_VENTA")
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()

        df = pd.DataFrame.from_records(rows, columns=cols)

        # Normalizar columnas numéricas
        for c in cols:
            if any(k in c.upper() for k in ["PPTO", "VENTA", "CANT", "PRECIO", "MONTO"]):
                df[c] = pd.to_numeric(
                    df[c].astype(str).str.replace(",", ".", regex=False) if df[c].dtype == object else df[c],
                    errors="coerce"
                ).fillna(0)

        # Limpiar columnas de texto clave
        for c in ["VENDEDOR_ACTUAL", "CATEGORIA_2026", "TIPO_ANALISIS", "ESTADO_ANALISIS"]:
            if c in df.columns:
                df[c] = df[c].astype(str).str.strip()

        print(f"[INFO] _load_ppto_vs_venta OK — {len(df)} filas | PPTO={_pvv_kpis.get('ppto_total',0):,.0f}")
        return df
    except Exception as e:
        print(f"[ERROR _load_ppto_vs_venta] {e}")
        return pd.DataFrame()


def _load_resumen_ud_grupo():
    """Carga venta 2026 vs 2025 YTD agrupada por UD_GRUPO desde DW_TOTAL_FACTURA."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH v26 AS (
                SELECT ISNULL(LTRIM(RTRIM(UD_GRUPO)), '(sin grupo)') AS ud_grupo,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT} AND {_DW_FILTRO}
                GROUP BY LTRIM(RTRIM(UD_GRUPO))
            ),
            v25 AS (
                SELECT ISNULL(LTRIM(RTRIM(UD_GRUPO)), '(sin grupo)') AS ud_grupo,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = 2025 AND {_DW_FILTRO}
                GROUP BY LTRIM(RTRIM(UD_GRUPO))
            )
            SELECT COALESCE(a.ud_grupo, b.ud_grupo) AS ud_grupo,
                   ISNULL(a.venta_2026_ytd, 0) AS venta_2026_ytd,
                   ISNULL(b.venta_2025_ytd, 0) AS venta_2025_ytd,
                   ISNULL(a.venta_2026_ytd, 0) - ISNULL(b.venta_2025_ytd, 0) AS diferencia
            FROM v26 a
            FULL OUTER JOIN v25 b ON b.ud_grupo = a.ud_grupo
            ORDER BY ISNULL(a.venta_2026_ytd, 0) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["venta_2026_ytd", "venta_2025_ytd", "diferencia"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["var_pct"] = df.apply(
            lambda r: (r["diferencia"] / r["venta_2025_ytd"] * 100) if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_resumen_ud_grupo] {e}")
        return pd.DataFrame()


def _load_cat_detalle(categoria):
    """Carga detalle de productos para una categoría.
    - PPTO: de [PPTO 2026], agrupado por zona + CODIGO (suma todos los clientes de la zona)
    - Venta: de DW_TOTAL_FACTURA, columna VENTA, por CATEGORIA + VENDEDOR + CODIGO
    - Muestra productos con PPTO sin venta Y productos con venta sin PPTO
    """
    try:
        conn = get_conn()
        cur = conn.cursor()
        cat_safe = categoria.replace("'", "''")
        cat_in_ppto = _cat_in(cat_safe)   # IN ('EQM','Servicios') si cat=='EQM'
        cat_in_dw   = _cat_in(cat_safe)
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            -- PPTO con código válido: productos reales comparables con venta
            ppto_prod AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       ISNULL(p.RUT, '') AS rut,
                       MAX(ISNULL(p.NOMBRE, '(sin nombre)')) AS nombre_cliente,
                       LTRIM(RTRIM(p.CODIGO)) AS CODIGO,
                       MAX(ISNULL(p.DESCRIPCION, '(sin descripción)')) AS descripcion,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)
                           * ISNULL(zp.peso_ytd, {_MES_ACT}/12.0)) AS ppto_ytd,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[CANT 2026],'0'),',','.') AS float)) AS cant_ppto
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND LTRIM(RTRIM(ISNULL(p.[CATEGOR\u00cdA 2026],''))) {cat_in_ppto}
                  AND ISNULL(LTRIM(RTRIM(p.CODIGO)),'') NOT IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), ISNULL(p.RUT,''),
                         LTRIM(RTRIM(p.CODIGO))
            ),
            -- PPTO incremental: sin código válido (no comparables con venta real)
            ppto_incr AS (
                SELECT ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)), '(sin zona)') AS zona,
                       MAX(ISNULL(p.DESCRIPCION, '(sin descripción)')) AS descripcion,
                       MAX(ISNULL(p.NOMBRE,'')) AS nombre_cliente,
                       SUM(TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float)) AS ppto_anual,
                       SUM({_sql_ppto_ytd()}) AS ppto_ytd
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND LTRIM(RTRIM(ISNULL(p.[CATEGOR\u00cdA 2026],''))) {cat_in_ppto}
                  AND ISNULL(LTRIM(RTRIM(p.CODIGO)),'') IN ('','SIN')
                GROUP BY LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),
                         ISNULL(LTRIM(RTRIM(p.DESCRIPCION)),'(sin descripción)'),
                         ISNULL(LTRIM(RTRIM(p.NOMBRE)),'')
            ),
            -- Venta real: DW_TOTAL_FACTURA por CATEGORIA + VENDEDOR + RUT + CODIGO
            real_prod AS (
                SELECT VENDEDOR AS zona,
                       RUT,
                       MAX(NOMBRE)      AS nombre_cliente,
                       CODIGO,
                       MAX(DESCRIPCION) AS descripcion_dw,
                       SUM(CASE WHEN ANO = {_ANO_ACT} AND MES <= {_MES_ACT}
                                THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ISNULL(CODIGO,'') <> '' AND ISNULL(CODIGO,'') <> 'SIN'
                  AND ISNULL(CATEGORIA,'') {cat_in_dw}
                  AND {_DW_FILTRO}
                GROUP BY VENDEDOR, RUT, CODIGO
            ),
            -- Venta total por zona+codigo (sin filtro RUT) para casos sin RUT en PPTO
            real_zona AS (
                SELECT zona, CODIGO,
                       SUM(venta_2026_ytd) AS venta_2026_ytd
                FROM real_prod
                GROUP BY zona, CODIGO
            )
            -- 1. Productos con PPTO y código válido
            SELECT p.zona, p.rut, LEFT(p.nombre_cliente, 50) AS nombre_cliente,
                   p.CODIGO, LEFT(ISNULL(p.descripcion, ''), 70) AS descripcion,
                   p.ppto_anual, p.ppto_ytd, p.cant_ppto,
                   CASE
                       WHEN p.rut <> '' THEN ISNULL(r.venta_2026_ytd, 0)
                       ELSE ISNULL(rz.venta_2026_ytd, 0)
                   END AS venta_2026_ytd,
                   'PPTO' AS origen,
                   0 AS es_incremental
            FROM ppto_prod p
            LEFT JOIN real_prod r  ON r.CODIGO = p.CODIGO AND r.zona = p.zona AND r.RUT = p.rut
            LEFT JOIN real_zona rz ON rz.CODIGO = p.CODIGO AND rz.zona = p.zona

            UNION ALL

            -- 2. Productos con venta en DW pero sin PPTO
            SELECT r.zona, r.RUT, LEFT(ISNULL(r.nombre_cliente, r.RUT), 50) AS nombre_cliente,
                   r.CODIGO, LEFT(ISNULL(r.descripcion_dw, r.CODIGO), 70) AS descripcion,
                   0 AS ppto_anual, 0 AS ppto_ytd, 0 AS cant_ppto,
                   r.venta_2026_ytd,
                   'SIN_PPTO' AS origen,
                   0 AS es_incremental
            FROM real_prod r
            WHERE r.venta_2026_ytd > 0
              AND NOT EXISTS (
                  SELECT 1 FROM ppto_prod p
                  WHERE p.CODIGO = r.CODIGO AND p.zona = r.zona AND p.rut = r.RUT
              )
              AND NOT EXISTS (
                  SELECT 1 FROM ppto_prod p
                  WHERE p.CODIGO = r.CODIGO AND p.zona = r.zona AND p.rut = ''
              )

            UNION ALL

            -- 3. PPTO incremental (sin código): solo PPTO, sin venta comparable
            SELECT i.zona, '' AS rut,
                   LEFT(ISNULL(i.nombre_cliente,'INCREMENTALES'), 50) AS nombre_cliente,
                   '(sin código)' AS CODIGO,
                   LEFT(i.descripcion, 70) AS descripcion,
                   i.ppto_anual, i.ppto_ytd, 0 AS cant_ppto,
                   0 AS venta_2026_ytd,
                   'INCR' AS origen,
                   1 AS es_incremental
            FROM ppto_incr i

            ORDER BY es_incremental, ppto_anual DESC, venta_2026_ytd DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["ppto_anual", "ppto_ytd", "cant_ppto", "venta_2026_ytd", "es_incremental"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        return df
    except Exception as e:
        print(f"[ERROR _load_cat_detalle] {e}")
        return pd.DataFrame()


def _build_cat_detalle(categoria):
    """Construye el panel de detalle de productos para una categoría."""
    df = _load_cat_detalle(categoria)
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion(f"📦 Detalle productos: {categoria}"),
            _empty_msg("Sin datos para esta categoría."),
        ])

    zonas = sorted(df["zona"].dropna().unique().tolist())
    zona_opts = [{"label": "Todas las zonas", "value": "TODAS"}] + \
                [{"label": z, "value": z} for z in zonas]

    clientes_ppto = df[df["origen"] == "PPTO"][["rut", "nombre_cliente"]].drop_duplicates()
    cli_opts = [{"label": "Todos los clientes", "value": "TODOS"}] + \
               [{"label": f"{row['nombre_cliente']} ({row['rut']})", "value": row["rut"]}
                for _, row in clientes_ppto.sort_values("nombre_cliente").iterrows() if row["rut"]]

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"📦 Detalle productos — {categoria} (YTD {mes_nombre} {_ANO_ACT})", "#2E75B6"),
        dcc.Store(id="cat-det-store", data=df.to_dict("records")),
        dcc.Store(id="cat-det-categoria", data=categoria),
        html.P("Venta: DW_TOTAL_FACTURA por Categoría + Zona + Código. PPTO: por cliente presupuestado.",
               style={"fontSize": "11px", "color": "#888", "marginBottom": "8px"}),
        # Filtros integrados en la tabla de productos
        html.Div(style={"display": "flex", "gap": "16px", "marginBottom": "12px", "flexWrap": "wrap"}, children=[
            html.Div(children=[
                html.Label("Zona:", style={"fontSize": "12px", "fontWeight": "600", "marginRight": "6px"}),
                dcc.Dropdown(id="cat-det-zona", options=zona_opts, value="TODAS", clearable=False,
                             style={"width": "260px", "fontSize": "12px"}),
            ]),
            html.Div(children=[
                html.Label("Cliente:", style={"fontSize": "12px", "fontWeight": "600", "marginRight": "6px"}),
                dcc.Dropdown(id="cat-det-cli", options=cli_opts, value="TODOS", clearable=False,
                             style={"width": "340px", "fontSize": "12px"}),
            ]),
            html.Div(children=[
                html.Label("Buscar:", style={"fontSize": "12px", "fontWeight": "600", "marginRight": "6px"}),
                dcc.Input(id="cat-det-buscar", type="text", placeholder="Código o descripción...",
                          debounce=True,
                          style={"width": "220px", "fontSize": "12px", "padding": "5px 8px",
                                 "border": "1px solid #ccc", "borderRadius": "4px"}),
            ]),
        ]),
        html.Div(id="cat-det-tabla"),
        html.Hr(style={"margin": "20px 0", "borderColor": "#ddd"}),
        # Top 25 por cliente para esta categoría
        html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr", "gap": "24px"}, children=[
            html.Div(children=[
                html.H4("🔴 Top 25 Clientes — Mayor Caída en esta Categoría",
                        style={"color": "#c0392b", "fontSize": "13px", "fontWeight": "700", "marginBottom": "8px"}),
                html.Div(id="cat-top25-caida"),
            ]),
            html.Div(children=[
                html.H4("🟢 Top 25 Clientes — Mayor Crecimiento en esta Categoría",
                        style={"color": "#27ae60", "fontSize": "13px", "fontWeight": "700", "marginBottom": "8px"}),
                html.Div(id="cat-top25-alza"),
            ]),
        ]),
        html.Div(id="cat-top25-cli-detalle"),
    ])


def _reload_all_data():
    """Recarga todos los globals desde la BD y guarda caché local."""
    global _df_categoria, _df_zona, _df_zona_cat, _df_cliente
    global _df_producto, _df_desalineacion, _df_clientes_caida, _df_precios, _df_incremental
    global _df_resumen_ud_grupo, _df_ppto_vs_venta, _pvv_kpis, _last_update
    _df_categoria        = _load_categoria()
    _df_zona             = _load_zona()
    _df_zona_cat         = _load_zona_cat()
    _df_cliente          = _load_cliente()
    _df_producto         = _load_producto()
    _df_desalineacion    = _load_desalineacion()
    _df_clientes_caida   = _load_clientes_caida()
    _df_precios          = _load_precios_analisis()
    _df_incremental      = _load_ppto_incremental()
    _df_resumen_ud_grupo = _load_resumen_ud_grupo()
    _df_ppto_vs_venta    = _load_ppto_vs_venta()
    _last_update         = datetime.datetime.now().strftime("%d/%m/%Y %H:%M:%S")
    pairs = [
        ("categoria",        _df_categoria),
        ("zona",             _df_zona),
        ("zona_cat",         _df_zona_cat),
        ("cliente",          _df_cliente),
        ("producto",         _df_producto),
        ("desalineacion",    _df_desalineacion),
        ("clientes_caida",   _df_clientes_caida),
        ("precios",          _df_precios),
        ("incremental",      _df_incremental),
        ("resumen_ud_grupo", _df_resumen_ud_grupo),
        ("ppto_vs_venta",    _df_ppto_vs_venta),
    ]
    for name, df in pairs:
        if not df.empty:
            _save_cache(name, df)
    if _pvv_kpis:
        _save_kpis_cache(_pvv_kpis)
    with open(os.path.join(DATA_DIR, "last_update.txt"), "w") as f:
        f.write(_last_update)
    print(f"[INFO] Datos actualizados y caché guardado: {_last_update}")


def _load_from_cache():
    """Carga todos los globals desde archivos de caché locales."""
    global _df_categoria, _df_zona, _df_zona_cat, _df_cliente
    global _df_producto, _df_desalineacion, _df_clientes_caida, _df_precios, _df_incremental
    global _df_resumen_ud_grupo, _df_ppto_vs_venta, _pvv_kpis, _last_update
    _df_categoria        = _load_cache("categoria")
    _df_zona             = _load_cache("zona")
    _df_zona_cat         = _load_cache("zona_cat")
    _df_cliente          = _load_cache("cliente")
    _df_producto         = _load_cache("producto")
    _df_desalineacion    = _load_cache("desalineacion")
    _df_clientes_caida   = _load_cache("clientes_caida")
    _df_precios          = _load_cache("precios")
    _df_incremental      = _load_cache("incremental")
    _df_resumen_ud_grupo = _load_cache("resumen_ud_grupo")
    _df_ppto_vs_venta    = _load_cache("ppto_vs_venta")
    _pvv_kpis            = _load_kpis_cache()
    ts_path = os.path.join(DATA_DIR, "last_update.txt")
    if os.path.exists(ts_path):
        with open(ts_path) as f:
            _last_update = f.read().strip()
    loaded = [n for n in _CACHE_NAMES
              if not globals().get(f"_df_{n.replace('_cat','_cat') }", pd.DataFrame()).empty]
    print(f"[INFO] Caché cargado ({_last_update})")


# ─── Layout Tab 1: Validación ──────────────────────────────────────────────────

# ─── Layout Tab 1: Por Categoría ──────────────────────────────────────────────

def layout_categoria():
    df = _df_categoria
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("🗂️ Por Categoría"),
            _empty_msg("No se pudieron cargar datos por categoría."),
        ])

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    total_ppto_anual  = df["ppto_anual"].sum()          # solo c/cód
    total_ppto_incr   = df["ppto_incr_anual"].sum()     if "ppto_incr_anual"   in df.columns else 0
    total_ppto_nuevos = df["ppto_nuevos_anual"].sum()   if "ppto_nuevos_anual" in df.columns else 0
    total_ppto_total  = df["ppto_total"].sum()           if "ppto_total"        in df.columns else total_ppto_anual
    total_ppto_ytd    = df["total_ppto_ytd"].sum()       if "total_ppto_ytd"    in df.columns else df["ppto_ytd"].sum()
    total_venta_ytd   = df["venta_2026_ytd"].sum()
    cumpl_total       = (total_venta_ytd / total_ppto_ytd * 100) if total_ppto_ytd > 0 else 0
    gap_total         = total_venta_ytd - total_ppto_ytd

    # Tabla
    total_facturas = df["venta_facturas"].sum() if "venta_facturas" in df.columns else total_venta_ytd
    total_guias    = df["venta_guias"].sum()    if "venta_guias"    in df.columns else 0

    table_data = []
    for _, r in df.iterrows():
        table_data.append({
            "cat": r["categoria"],
            "sem": _sem(r["cumpl_pct"]),
            "PPTO 2026": _fmt_abs(r.get("ppto_total", r["ppto_anual"])),
            "PPTO c/Cód": _fmt_abs(r["ppto_anual"]),
            "PPTO Incr.": _fmt_abs(r.get("ppto_incr_anual", 0)),
            "PPTO Prod. Nuevos": _fmt_abs(r.get("ppto_nuevos_anual", 0)),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(r.get("total_ppto_ytd", r["ppto_ytd"])),
            "Facturas YTD": _fmt_abs(r.get("venta_facturas", r["venta_2026_ytd"])),
            "Guías Pend.": _fmt_abs(r.get("venta_guias", 0)),
            "Venta Total YTD": _fmt_abs(r["venta_2026_ytd"]),
            "Alcance %": _fmt_pct(r["cumpl_pct"]),
            "Gap $": _fmt_abs(r["gap"]),
            "Venta 2025 YTD": _fmt_abs(r["venta_2025_ytd"]),
            "Crec vs 2025": _fmt_pct(r["crec_pct"]),
        })
    total_v25 = df["venta_2025_ytd"].sum()
    crec_total_pct = (total_venta_ytd - total_v25) / total_v25 * 100 if total_v25 > 0 else float("nan")
    table_data.append({
        "cat": "── TOTAL",
        "sem": "",
        "PPTO 2026": _fmt_abs(total_ppto_total),
        "PPTO c/Cód": _fmt_abs(total_ppto_anual),
        "PPTO Incr.": _fmt_abs(total_ppto_incr),
        "PPTO Prod. Nuevos": _fmt_abs(total_ppto_nuevos),
        f"PPTO YTD ({mes_nombre})": _fmt_abs(total_ppto_ytd),
        "Facturas YTD": _fmt_abs(total_facturas),
        "Guías Pend.": _fmt_abs(total_guias),
        "Venta Total YTD": _fmt_abs(total_venta_ytd),
        "Alcance %": _fmt_pct(cumpl_total),
        "Gap $": _fmt_abs(gap_total),
        "Venta 2025 YTD": _fmt_abs(total_v25),
        "Crec vs 2025": _fmt_pct(crec_total_pct),
    })

    return html.Div(children=[
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"🗂️ Por Categoría — YTD {mes_nombre} {_ANO_ACT}"),
            _kpi_row([
                _kpi_card("PPTO 2026 Total", _fmt(total_ppto_total), "#1F3864"),
                _kpi_card(f"PPTO YTD ({mes_nombre})", _fmt(total_ppto_ytd), "#2E75B6",
                          f"c/Cód + Incr. (Prod.Nuevos $0 antes Mayo)"),
                _kpi_card("Venta 2026 YTD", _fmt(total_venta_ytd), "#27ae60"),
                _kpi_card("Alcance %", _fmt_pct(cumpl_total),
                          "#27ae60" if cumpl_total >= 80 else "#e67e22" if cumpl_total >= 50 else "#c0392b"),
                _kpi_card("Gap vs PPTO YTD", _fmt(gap_total),
                          "#c0392b" if gap_total < 0 else "#27ae60"),
            ]),
        ]),
        html.Div(style=CARD_STYLE, children=[
            _seccion("Detalle por Categoría"),
            html.P("💡 Haz clic en una fila para ver el detalle de productos.",
                   style={"fontSize": "11px", "color": "#666", "marginBottom": "8px"}),
            dash_table.DataTable(
                id="cat-tabla",
                data=table_data,
                columns=[{"name": c, "id": c} for c in list(table_data[0].keys())] if table_data else [],
                style_table={"overflowX": "auto"},
                style_header={"backgroundColor": "#1F3864", "color": "white",
                              "fontWeight": "bold", "fontSize": "12px"},
                style_cell={"fontSize": "12px", "padding": "6px 10px",
                            "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
                style_cell_conditional=[
                    {"if": {"column_id": "cat"}, "textAlign": "left"},
                    {"if": {"column_id": "sem"}, "textAlign": "center", "width": "40px"},
                ],
                style_data_conditional=[
                    {"if": {"filter_query": '{cat} = "── TOTAL"'}, "fontWeight": "700", "backgroundColor": "#e8edf7"},
                ],
                page_size=30,
                sort_action="native",
                row_selectable="single",
            ),
        ]),
        html.Div(id="cat-detalle-container", style={"marginTop": "0"}),
    ])


# ─── Layout Tab 3: Por Zona ────────────────────────────────────────────────────

_CATS_VALIDAS = ["SQ", "EVA", "MAH", "EQM"]


def _load_top25_clientes():
    """Carga Top 25 clientes con mayor caída y mayor alza YTD 2026 vs 2025, desde DW_TOTAL_FACTURA."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH
            v26 AS (
                SELECT RUT,
                       MAX(LEFT(ISNULL(NOMBRE, RUT), 50)) AS nombre,
                       MAX(ISNULL(VENDEDOR,'')) AS zona,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT} AND {_DW_FILTRO}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                GROUP BY RUT
            ),
            v25 AS (
                SELECT RUT,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = 2025 AND {_DW_FILTRO}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                GROUP BY RUT
            )
            SELECT a.RUT, a.nombre, a.zona,
                   a.venta_2026_ytd,
                   ISNULL(b.venta_2025_ytd, 0) AS venta_2025_ytd,
                   a.venta_2026_ytd - ISNULL(b.venta_2025_ytd, 0) AS diferencia
            FROM v26 a
            LEFT JOIN v25 b ON b.RUT = a.RUT
            WHERE ISNULL(b.venta_2025_ytd, 0) > 0
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["venta_2026_ytd", "venta_2025_ytd", "diferencia"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["dif_pct"] = df.apply(
            lambda r: (r["diferencia"] / r["venta_2025_ytd"] * 100) if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_top25_clientes] {e}")
        return pd.DataFrame()


def _load_detalle_cliente(rut: str):
    """Carga detalle por producto y categoría de un cliente, comparando 2026 vs 2025 YTD."""
    try:
        conn = get_conn()
        cur = conn.cursor()
        rut_safe = rut.replace("'", "''")
        cur.execute(f"""
            WITH
            v26 AS (
                SELECT CODIGO, MAX(LEFT(ISNULL(DESCRIPCION, CODIGO), 60)) AS descripcion,
                       {_SQL_CAT_DW} AS categoria,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026
                FROM DW_TOTAL_FACTURA
                WHERE ANO = {_ANO_ACT} AND RUT = '{rut_safe}'
                  AND {_DW_FILTRO}
                GROUP BY CODIGO, {_SQL_CAT_DW}
            ),
            v25 AS (
                SELECT CODIGO, MAX(LEFT(ISNULL(DESCRIPCION, CODIGO), 60)) AS descripcion,
                       {_SQL_CAT_DW} AS categoria,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025
                FROM DW_TOTAL_FACTURA
                WHERE ANO = 2025 AND RUT = '{rut_safe}'
                  AND {_DW_FILTRO}
                GROUP BY CODIGO, {_SQL_CAT_DW}
            )
            SELECT COALESCE(a.CODIGO, b.CODIGO) AS codigo,
                   COALESCE(a.descripcion, b.descripcion, COALESCE(a.CODIGO, b.CODIGO)) AS descripcion,
                   COALESCE(a.categoria, b.categoria) AS categoria,
                   ISNULL(a.venta_2026, 0) AS venta_2026,
                   ISNULL(b.venta_2025, 0) AS venta_2025,
                   ISNULL(a.venta_2026, 0) - ISNULL(b.venta_2025, 0) AS diferencia
            FROM v26 a
            FULL OUTER JOIN v25 b ON b.CODIGO = a.CODIGO AND b.categoria = a.categoria
            ORDER BY ISNULL(a.venta_2026,0) + ISNULL(b.venta_2025,0) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["venta_2026", "venta_2025", "diferencia"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["dif_pct"] = df.apply(
            lambda r: (r["diferencia"] / r["venta_2025"] * 100) if r["venta_2025"] > 0 else float("nan"),
            axis=1
        )
        return df
    except Exception as e:
        print(f"[ERROR _load_detalle_cliente] {e}")
        return pd.DataFrame()


def layout_zona():
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    cat_opts = [{"label": "Todas las categorías", "value": "TODAS"}] + \
               [{"label": c, "value": c} for c in _CATS_VALIDAS]

    zonas_disponibles = sorted(
        _df_zona[_df_zona["ppto_anual"] > 0]["zona"].dropna().unique().tolist()
    )
    zona_opts = [{"label": "— Selecciona una zona —", "value": ""}] + \
                [{"label": z, "value": z} for z in zonas_disponibles]

    df_init = _df_zona[_df_zona["ppto_anual"] > 0].copy()
    kpi_init, tabla_init = _build_zona_content(df_init)

    return html.Div(children=[
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"🏢 Por Zona — YTD {mes_nombre} {_ANO_ACT}"),
            html.Div(style={"display": "flex", "gap": "32px", "flexWrap": "wrap",
                            "alignItems": "center", "marginBottom": "12px"}, children=[
                html.Div(style={"display": "flex", "alignItems": "center"}, children=[
                    html.Label("Filtrar por Categoría:", style={"fontWeight": "600", "fontSize": "13px",
                                                                 "marginRight": "8px", "whiteSpace": "nowrap"}),
                    dcc.Dropdown(
                        id="zona-cat-filtro",
                        options=cat_opts,
                        value="TODAS",
                        clearable=False,
                        style={"width": "260px"},
                    ),
                ]),
                html.Div(style={"display": "flex", "alignItems": "center"}, children=[
                    html.Label("Top 25 clientes de zona:", style={"fontWeight": "600", "fontSize": "13px",
                                                                    "marginRight": "8px", "whiteSpace": "nowrap",
                                                                    "color": "#1F3864"}),
                    dcc.Dropdown(
                        id="zona-selector",
                        options=zona_opts,
                        value="",
                        clearable=False,
                        placeholder="Selecciona una zona...",
                        style={"width": "280px"},
                    ),
                ]),
            ]),
            html.Div(id="zona-kpi-row", children=kpi_init),
        ]),
        html.Div(id="zona-tabla-container", children=tabla_init),
        html.Div(id="zona-top25-container"),
        html.Div(id="zona-caida-container"),
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"📉 Top 25 Clientes — Variación YTD {mes_nombre} 2026 vs 2025"),
            html.P("Comparación de venta real YTD 2026 vs mismo período 2025 (DW_TOTAL_FACTURA). "
                   "💡 Haz clic en un cliente para ver el detalle por producto y categoría.",
                   style={"fontSize": "11px", "color": "#666", "marginBottom": "12px"}),
            html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr", "gap": "24px"}, children=[
                html.Div(children=[
                    html.H4("🔴 Top 25 Mayor Caída",
                            style={"color": "#c0392b", "fontSize": "13px", "fontWeight": "700", "marginBottom": "8px"}),
                    html.Div(id="zona-top25-caida"),
                ]),
                html.Div(children=[
                    html.H4("🟢 Top 25 Mayor Crecimiento",
                            style={"color": "#27ae60", "fontSize": "13px", "fontWeight": "700", "marginBottom": "8px"}),
                    html.Div(id="zona-top25-alza"),
                ]),
            ]),
        ]),
        html.Div(id="zona-cli-detalle"),
    ])


def _build_zona_content(df):
    """Construye KPI row y tabla para el tab zona dado un DataFrame filtrado."""
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    if df.empty:
        return html.Div(), _empty_msg("Sin datos.")

    total_ppto_anual = df["ppto_anual"].sum()
    total_ppto_ytd   = df["ppto_ytd"].sum()
    total_venta_ytd  = df["venta_2026_ytd"].sum()
    total_v25_ytd    = df["venta_2025_ytd"].sum()
    total_proy       = df["proyeccion_anual"].sum() if "proyeccion_anual" in df.columns else 0
    cumpl_total      = (total_venta_ytd / total_ppto_ytd * 100) if total_ppto_ytd > 0 else 0
    crec_total_pct   = (total_venta_ytd - total_v25_ytd) / total_v25_ytd * 100 if total_v25_ytd > 0 else float("nan")

    kpi_row = _kpi_row([
        _kpi_card("PPTO Anual 2026",       _fmt(total_ppto_anual), "#1F3864"),
        _kpi_card(f"PPTO YTD ({mes_nombre})", _fmt(total_ppto_ytd),   "#2E75B6"),
        _kpi_card("Venta 2026 YTD",        _fmt(total_venta_ytd),  "#27ae60"),
        _kpi_card("Cumplimiento %",         _fmt_pct(cumpl_total),
                  "#27ae60" if cumpl_total >= 80 else "#e67e22" if cumpl_total >= 50 else "#c0392b"),
    ])

    has_guias = "venta_facturas" in df.columns and "venta_guias" in df.columns
    total_facturas = df["venta_facturas"].sum() if has_guias else total_venta_ytd
    total_guias    = df["venta_guias"].sum()    if has_guias else 0

    table_data = []
    for _, r in df.iterrows():
        row = {
            "Zona":            r["zona"],
            "sem":             _sem(r["cumpl_pct"]),
            "PPTO Anual":      _fmt_abs(r["ppto_anual"]),
            "PPTO YTD":        _fmt_abs(r["ppto_ytd"]),
        }
        if has_guias:
            row["Facturas YTD"] = _fmt_abs(r["venta_facturas"])
            row["Guías Pend."]  = _fmt_abs(r["venta_guias"])
        row.update({
            "Venta Total YTD": _fmt_abs(r["venta_2026_ytd"]),
            "Cumpl %":         _fmt_pct(r["cumpl_pct"]),
            "Gap $":           _fmt_abs(r["gap"]),
            "Venta 2025 YTD":  _fmt_abs(r["venta_2025_ytd"]),
            "Crec vs 2025":    _fmt_pct(r["crec_pct"]),
            "Proyec. Anual":   _fmt_abs(r["proyeccion_anual"]) if "proyeccion_anual" in df.columns else "—",
            "Gap vs PPTO Anual": _fmt_abs(r["gap_vs_ppto_anual"]) if "gap_vs_ppto_anual" in df.columns else "—",
        })
        table_data.append(row)
    total_row = {
        "Zona":            "── TOTAL",
        "sem":             "",
        "PPTO Anual":      _fmt_abs(total_ppto_anual),
        "PPTO YTD":        _fmt_abs(total_ppto_ytd),
    }
    if has_guias:
        total_row["Facturas YTD"] = _fmt_abs(total_facturas)
        total_row["Guías Pend."]  = _fmt_abs(total_guias)
    total_row.update({
        "Venta Total YTD": _fmt_abs(total_venta_ytd),
        "Cumpl %":         _fmt_pct(cumpl_total),
        "Gap $":           _fmt_abs(total_venta_ytd - total_ppto_ytd),
        "Venta 2025 YTD":  _fmt_abs(total_v25_ytd),
        "Crec vs 2025":    _fmt_pct(crec_total_pct),
        "Proyec. Anual":   _fmt_abs(total_proy),
        "Gap vs PPTO Anual": _fmt_abs(total_ppto_anual - total_proy),
    })
    table_data.append(total_row)

    tabla = html.Div(style=CARD_STYLE, children=[
        _seccion("Detalle por Zona"),
        dash_table.DataTable(
            data=table_data,
            columns=[{"name": c, "id": c} for c in list(table_data[0].keys())],
            style_table={"overflowX": "auto"},
            style_header={"backgroundColor": "#1F3864", "color": "white",
                          "fontWeight": "bold", "fontSize": "12px"},
            style_cell={"fontSize": "12px", "padding": "6px 10px",
                        "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
            style_cell_conditional=[
                {"if": {"column_id": "Zona"}, "textAlign": "left"},
                {"if": {"column_id": "sem"}, "textAlign": "center", "width": "40px"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{Zona} = "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            page_size=30,
            sort_action="native",
        ),
    ])

    return kpi_row, tabla


def _build_caida_clientes(zona: str) -> html.Div:
    """Clientes de la zona (2026) cuya venta 2026 YTD < venta 2025 YTD, ordenados por mayor caída."""
    df = _df_clientes_caida
    if df.empty or not zona or zona == "TODAS":
        return html.Div()

    df = df[df["zona"] == zona].copy()
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion(f"📉 Clientes en Caída 2026 vs 2025 — {zona}"),
            html.Div("✅ Ningún cliente de esta zona muestra caída respecto a 2025.",
                     style={"color": "#27ae60", "fontSize": "13px", "padding": "12px 0",
                            "fontWeight": "600"}),
        ])

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    df = df.sort_values("diferencia", ascending=True)

    tot_ppto     = df["ppto_anual"].sum()
    tot_v26      = df["venta_2026_ytd"].sum()
    tot_v25      = df["venta_2025_ytd"].sum()
    tot_dif      = df["diferencia"].sum()
    tot_caida_pct = (tot_v26 - tot_v25) / tot_v25 * 100 if tot_v25 > 0 else float("nan")

    table_data = []
    for _, r in df.iterrows():
        table_data.append({
            "Cliente":         r["nombre"],
            "PPTO 2026":       _fmt_abs(r["ppto_anual"]) if r["ppto_anual"] > 0 else "—",
            "Venta 2026 YTD":  _fmt_abs(r["venta_2026_ytd"]),
            "Venta 2025 YTD":  _fmt_abs(r["venta_2025_ytd"]),
            "Diferencia $":    _fmt_abs(r["diferencia"]),
            "Caída %":         _fmt_pct(r["caida_pct"]),
            "_dif":            r["diferencia"],
        })
    table_data.append({
        "Cliente":        f"── TOTAL ({len(df)} clientes)",
        "PPTO 2026":      _fmt_abs(tot_ppto) if tot_ppto > 0 else "—",
        "Venta 2026 YTD": _fmt_abs(tot_v26),
        "Venta 2025 YTD": _fmt_abs(tot_v25),
        "Diferencia $":   _fmt_abs(tot_dif),
        "Caída %":        _fmt_pct(tot_caida_pct),
        "_dif":           tot_dif,
    })

    display_cols = ["Cliente", "PPTO 2026", "Venta 2026 YTD", "Venta 2025 YTD",
                    "Diferencia $", "Caída %"]

    kpis = _kpi_row([
        _kpi_card("Clientes en caída", str(len(df)), "#c0392b"),
        _kpi_card("Caída total $", _fmt(tot_dif), "#c0392b"),
        _kpi_card("Caída %", _fmt_pct(tot_caida_pct), "#c0392b"),
        _kpi_card("PPTO 2026 en riesgo", _fmt(tot_ppto) if tot_ppto > 0 else "—",
                  "#e67e22" if tot_ppto > 0 else "#888"),
    ])

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"📉 Clientes en Caída 2026 vs 2025 — {zona} | YTD {mes_nombre} {_ANO_ACT}"),
        kpis,
        html.P(
            "Solo clientes con actividad en ambos años. "
            "Zona asignada según actividad 2026. Ordenados por mayor caída absoluta.",
            style={"fontSize": "11px", "color": "#666", "marginBottom": "10px"},
        ),
        dash_table.DataTable(
            data=[{k: v for k, v in row.items() if not k.startswith("_")} for row in table_data],
            columns=[{"name": c, "id": c} for c in display_cols],
            style_table={"overflowX": "auto"},
            style_header={"backgroundColor": "#c0392b", "color": "white",
                          "fontWeight": "bold", "fontSize": "12px"},
            style_cell={"fontSize": "12px", "padding": "6px 10px",
                        "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
            style_cell_conditional=[
                {"if": {"column_id": "Cliente"}, "textAlign": "left",
                 "maxWidth": "240px", "overflow": "hidden", "textOverflow": "ellipsis"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": f'{{Cliente}} contains "TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#fdecea"},
                {"if": {"filter_query": '{_dif} < 0', "column_id": "Diferencia $"},
                 "color": "#c0392b", "fontWeight": "600"},
                {"if": {"filter_query": '{_dif} < 0', "column_id": "Caída %"},
                 "color": "#c0392b"},
            ],
            page_size=50,
            sort_action="native",
        ),
    ])


def _build_top25_clientes(zona: str) -> html.Div:
    """Top 25 clientes con mayor brecha en la zona seleccionada."""
    df = _df_cliente
    if df.empty or not zona or zona == "TODAS":
        return html.Div()

    df = df[df["zona"] == zona].copy()
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion(f"Top 25 Clientes — {zona}"),
            _empty_msg("Sin clientes con PPTO para esta zona."),
        ])

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    # Precio estimado = valor / cantidad (YTD)
    df["precio_ppto_est"] = df.apply(
        lambda r: r["ppto_ytd"] / r["cant_ppto"] if r["cant_ppto"] > 0 else 0, axis=1
    )
    df["precio_real_est"] = df.apply(
        lambda r: r["venta_2026_ytd"] / r["cant_2026_ytd"] if r["cant_2026_ytd"] > 0 else 0, axis=1
    )
    # Gap precio por unidad: real - ppto (negativo = vendió más barato que lo presupuestado)
    df["gap_precio_unit"] = df.apply(
        lambda r: r["precio_real_est"] - r["precio_ppto_est"]
        if r["precio_ppto_est"] > 0 and r["precio_real_est"] > 0 else float("nan"), axis=1
    )
    # Dejado de ganar por precio = (precio_ppto - precio_real) × cant_real
    df["dejado_ganar_precio"] = df.apply(
        lambda r: (r["precio_ppto_est"] - r["precio_real_est"]) * r["cant_2026_ytd"]
        if r["precio_ppto_est"] > 0 and r["precio_real_est"] > 0 and r["cant_2026_ytd"] > 0
        else 0, axis=1
    )
    df["gap_cant"] = df["cant_2026_ytd"] - df["cant_ppto"]

    # Ordenar por mayor brecha negativa (más lejos del PPTO)
    top25 = df.sort_values("gap", ascending=True).head(25)

    # Totales
    tot_ppto_ytd   = top25["ppto_ytd"].sum()
    tot_venta_26   = top25["venta_2026_ytd"].sum()
    tot_venta_25   = top25["venta_2025_ytd"].sum()
    tot_gap        = top25["gap"].sum()
    tot_cant_ppto  = top25["cant_ppto"].sum()
    tot_cant_real  = top25["cant_2026_ytd"].sum()
    tot_gap_cant   = top25["gap_cant"].sum()
    tot_dejado     = top25["dejado_ganar_precio"].sum()
    tot_cumpl      = tot_venta_26 / tot_ppto_ytd * 100 if tot_ppto_ytd > 0 else float("nan")

    table_data = []
    for _, r in top25.iterrows():
        table_data.append({
            "sem":              _sem(r["cumpl_pct"]),
            "Cliente":          r["nombre"],
            "Venta 2025 YTD":   _fmt_abs(r["venta_2025_ytd"]),
            "Venta 2026 YTD":   _fmt_abs(r["venta_2026_ytd"]),
            "PPTO YTD":         _fmt_abs(r["ppto_ytd"]),
            "Gap $":            _fmt_abs(r["gap"]),
            "Cumpl %":          _fmt_pct(r["cumpl_pct"]),
            "Cant PPTO":        f"{r['cant_ppto']:,.0f}" if r["cant_ppto"] > 0 else "—",
            "Cant Real":        f"{r['cant_2026_ytd']:,.0f}" if r["cant_2026_ytd"] > 0 else "—",
            "Gap Cant":         f"{r['gap_cant']:+,.0f}" if r["cant_ppto"] > 0 else "—",
            "Precio PPTO Est":  _fmt_abs(r["precio_ppto_est"]) if r["precio_ppto_est"] > 0 else "—",
            "Precio Real Est":  _fmt_abs(r["precio_real_est"]) if r["precio_real_est"] > 0 else "—",
            "Gap Precio/U":     _fmt_abs(r["gap_precio_unit"]) if pd.notna(r["gap_precio_unit"]) else "—",
            "Dejado Ganar x$":  _fmt_abs(r["dejado_ganar_precio"]) if r["dejado_ganar_precio"] != 0 else "—",
            "_gap_raw":         r["gap"],
            "_cumpl":           r["cumpl_pct"] if pd.notna(r["cumpl_pct"]) else -1,
        })
    table_data.append({
        "sem":             "",
        "Cliente":         "── TOTAL TOP 25",
        "Venta 2025 YTD":  _fmt_abs(tot_venta_25),
        "Venta 2026 YTD":  _fmt_abs(tot_venta_26),
        "PPTO YTD":        _fmt_abs(tot_ppto_ytd),
        "Gap $":           _fmt_abs(tot_gap),
        "Cumpl %":         _fmt_pct(tot_cumpl),
        "Cant PPTO":       f"{tot_cant_ppto:,.0f}",
        "Cant Real":       f"{tot_cant_real:,.0f}",
        "Gap Cant":        f"{tot_gap_cant:+,.0f}",
        "Precio PPTO Est": "—",
        "Precio Real Est": "—",
        "Gap Precio/U":    "—",
        "Dejado Ganar x$": _fmt_abs(tot_dejado),
        "_gap_raw":        tot_gap,
        "_cumpl":          -1,
    })

    display_cols = ["sem", "Cliente", "Venta 2025 YTD", "Venta 2026 YTD", "PPTO YTD",
                    "Gap $", "Cumpl %", "Cant PPTO", "Cant Real", "Gap Cant",
                    "Precio PPTO Est", "Precio Real Est", "Gap Precio/U", "Dejado Ganar x$"]

    style_data_cond = [
        {"if": {"filter_query": '{Cliente} = "── TOTAL TOP 25"'},
         "fontWeight": "700", "backgroundColor": "#e8edf7"},
        {"if": {"filter_query": '{_gap_raw} < 0', "column_id": "Gap $"},
         "color": "#c0392b", "fontWeight": "600"},
        {"if": {"filter_query": '{_cumpl} < 50', "column_id": "Cumpl %"},
         "color": "#c0392b"},
        {"if": {"filter_query": '{_cumpl} >= 80', "column_id": "Cumpl %"},
         "color": "#27ae60"},
    ]

    # KPIs del top 25
    kpis = _kpi_row([
        _kpi_card(f"Gap Total Top 25", _fmt(tot_gap),
                  "#c0392b" if tot_gap < 0 else "#27ae60"),
        _kpi_card("Cumpl % Top 25", _fmt_pct(tot_cumpl),
                  "#27ae60" if tot_cumpl >= 80 else "#e67e22" if tot_cumpl >= 50 else "#c0392b"),
        _kpi_card("Gap Cantidad", f"{tot_gap_cant:+,.0f}",
                  "#c0392b" if tot_gap_cant < 0 else "#27ae60"),
        _kpi_card("Dejado Ganar x Precio", _fmt(tot_dejado),
                  "#c0392b" if tot_dejado > 0 else "#27ae60"),
    ])

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"Top 25 Clientes por Brecha — {zona} | YTD {mes_nombre} {_ANO_ACT}"),
        kpis,
        html.P(
            "Ordenados por mayor brecha negativa (gap $ ascendente). "
            "Precio estimado = Valor YTD / Cantidad. "
            "Dejado Ganar = (Precio PPTO - Precio Real) × Cant. real.",
            style={"fontSize": "11px", "color": "#666", "marginBottom": "10px"},
        ),
        dash_table.DataTable(
            data=[{k: v for k, v in row.items() if not k.startswith("_")} for row in table_data],
            columns=[{"name": c, "id": c} for c in display_cols],
            style_table={"overflowX": "auto"},
            style_header={"backgroundColor": "#1F3864", "color": "white",
                          "fontWeight": "bold", "fontSize": "12px"},
            style_cell={"fontSize": "12px", "padding": "6px 10px",
                        "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
            style_cell_conditional=[
                {"if": {"column_id": "sem"},     "textAlign": "center", "width": "36px"},
                {"if": {"column_id": "Cliente"}, "textAlign": "left", "maxWidth": "220px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
            ],
            style_data_conditional=style_data_cond,
            page_size=26,
            sort_action="native",
        ),
    ])




# ─── Layout Tab 4: Por Cliente ─────────────────────────────────────────────────

def layout_cliente(zona_filtro="TODAS"):
    df = _df_cliente
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("👤 Por Cliente"),
            _empty_msg("No se pudieron cargar datos por cliente."),
        ])

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    # Opciones de zona para dropdown
    zonas = sorted(df["zona"].dropna().unique().tolist())
    dropdown_opts = [{"label": "Todas las zonas", "value": "TODAS"}] + \
                    [{"label": z, "value": z} for z in zonas]

    # KPIs globales (sin filtro)
    total_clientes   = len(df)
    clientes_80      = (df["cumpl_pct"] >= 80).sum()
    clientes_sin_vta = (df["venta_2026_ytd"] == 0).sum()
    gap_total        = df["gap"].sum()

    return html.Div(children=[
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"👤 Por Cliente — YTD {mes_nombre} {_ANO_ACT}"),
            _kpi_row([
                _kpi_card("Clientes en PPTO", str(total_clientes), "#1F3864"),
                _kpi_card("Clientes Cumpl ≥ 80%", str(clientes_80), "#27ae60"),
                _kpi_card("Clientes sin venta 2026", str(clientes_sin_vta),
                          "#c0392b" if clientes_sin_vta > 0 else "#27ae60"),
                _kpi_card("Gap Total YTD", _fmt(gap_total),
                          "#c0392b" if gap_total < 0 else "#27ae60"),
            ]),
        ]),
        html.Div(style=CARD_STYLE, children=[
            html.Div(style={"marginBottom": "12px"}, children=[
                html.Label("Filtrar por Zona:", style={"fontSize": "12px", "fontWeight": "600",
                                                        "marginRight": "8px"}),
                dcc.Dropdown(
                    id="cli-zona-filter",
                    options=dropdown_opts,
                    value="TODAS",
                    clearable=False,
                    style={"width": "300px", "fontSize": "12px", "display": "inline-block"},
                ),
            ]),
            html.Div(id="cli-tabla-container"),
        ]),
    ])


def _build_cliente_tabla(zona_filtro="TODAS"):
    df = _df_cliente
    if df.empty:
        return _empty_msg()
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    if zona_filtro and zona_filtro != "TODAS":
        df = df[df["zona"] == zona_filtro]
    if df.empty:
        return _empty_msg(f"No hay clientes para zona: {zona_filtro}")

    total_ppto_anual_cli = df["ppto_anual"].sum()
    total_ppto_ytd_cli   = df["ppto_ytd"].sum()
    total_venta_ytd_cli  = df["venta_2026_ytd"].sum()
    total_v25_ytd_cli    = df["venta_2025_ytd"].sum()
    cumpl_total_cli      = (total_venta_ytd_cli / total_ppto_ytd_cli * 100) if total_ppto_ytd_cli > 0 else float("nan")
    crec_total_cli       = (total_venta_ytd_cli - total_v25_ytd_cli) / total_v25_ytd_cli * 100 if total_v25_ytd_cli > 0 else float("nan")

    table_data = []
    for _, r in df.iterrows():
        row_style = "normal"
        if r["venta_2026_ytd"] == 0:
            row_style = "sin_venta"
        table_data.append({
            "zona":              r["zona"],
            "nombre":            r["nombre"],
            "sem":               _sem(r["cumpl_pct"]),
            "PPTO Anual":        _fmt_abs(r["ppto_anual"]),
            f"PPTO YTD":         _fmt_abs(r["ppto_ytd"]),
            "Venta 2026 YTD":    _fmt_abs(r["venta_2026_ytd"]),
            "Cumpl %":           _fmt_pct(r["cumpl_pct"]),
            "Gap $":             _fmt_abs(r["gap"]),
            "Venta 2025 YTD":    _fmt_abs(r["venta_2025_ytd"]),
            "Crec vs 2025":      _fmt_pct(r["crec_pct"]),
            "Prods PPTO":        int(r["productos_ppto"]) if pd.notna(r["productos_ppto"]) else 0,
            "_sin_venta":        row_style,
            "_cumpl":            r["cumpl_pct"] if pd.notna(r["cumpl_pct"]) else -1,
        })
    table_data.append({
        "zona":           "",
        "nombre":         "── TOTAL",
        "sem":            "",
        "PPTO Anual":     _fmt_abs(total_ppto_anual_cli),
        "PPTO YTD":       _fmt_abs(total_ppto_ytd_cli),
        "Venta 2026 YTD": _fmt_abs(total_venta_ytd_cli),
        "Cumpl %":        _fmt_pct(cumpl_total_cli),
        "Gap $":          _fmt_abs(total_venta_ytd_cli - total_ppto_ytd_cli),
        "Venta 2025 YTD": _fmt_abs(total_v25_ytd_cli),
        "Crec vs 2025":   _fmt_pct(crec_total_cli),
        "Prods PPTO":     "",
        "_sin_venta":     "normal",
        "_cumpl":         -1,
    })

    display_cols = ["zona", "nombre", "sem", "PPTO Anual", "PPTO YTD",
                    "Venta 2026 YTD", "Cumpl %", "Gap $", "Venta 2025 YTD",
                    "Crec vs 2025", "Prods PPTO"]

    style_data_cond = [
        {
            "if": {"filter_query": '{_sin_venta} = "sin_venta"'},
            "backgroundColor": "#ffeeee",
            "color": "#c0392b",
        },
        {"if": {"filter_query": '{nombre} = "── TOTAL"'}, "fontWeight": "700", "backgroundColor": "#e8edf7"},
    ]

    return dash_table.DataTable(
        data=[{k: v for k, v in row.items() if not k.startswith("_")} for row in table_data],
        columns=[{"name": c, "id": c} for c in display_cols],
        style_table={"overflowX": "auto"},
        style_header={"backgroundColor": "#1F3864", "color": "white",
                      "fontWeight": "bold", "fontSize": "12px"},
        style_cell={"fontSize": "12px", "padding": "6px 10px",
                    "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
        style_cell_conditional=[
            {"if": {"column_id": "zona"}, "textAlign": "left"},
            {"if": {"column_id": "nombre"}, "textAlign": "left", "maxWidth": "200px",
             "overflow": "hidden", "textOverflow": "ellipsis"},
            {"if": {"column_id": "sem"}, "textAlign": "center", "width": "40px"},
        ],
        style_data_conditional=style_data_cond,
        page_size=50,
        sort_action="native",
        filter_action="native",
    )


# ─── Layout Tab 5: Por Producto ────────────────────────────────────────────────

def layout_producto(zona_filtro="TODAS"):
    df = _df_producto
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("📦 Por Producto"),
            _empty_msg("No se pudieron cargar datos por producto."),
        ])

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    zonas = sorted(df["zona"].dropna().unique().tolist())
    dropdown_opts = [{"label": "Todas las zonas", "value": "TODAS"}] + \
                    [{"label": z, "value": z} for z in zonas]

    total_prods      = len(df)
    prods_sin_venta  = (df["venta_2026_ytd"] == 0).sum()
    gap_cant_total   = df["gap_cant"].sum()
    gap_venta_total  = df["gap_venta"].sum()

    return html.Div(children=[
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"📦 Por Producto — YTD {mes_nombre} {_ANO_ACT}"),
            _kpi_row([
                _kpi_card("Productos en PPTO (top 500)", str(total_prods), "#1F3864"),
                _kpi_card("Prods sin venta 2026", str(prods_sin_venta),
                          "#c0392b" if prods_sin_venta > 0 else "#27ae60"),
                _kpi_card("Gap Cantidad YTD", f"{gap_cant_total:,.0f}" if not math.isnan(gap_cant_total) else "—",
                          "#c0392b" if gap_cant_total < 0 else "#27ae60"),
                _kpi_card("Gap $ YTD", _fmt(gap_venta_total),
                          "#c0392b" if gap_venta_total < 0 else "#27ae60"),
            ]),
        ]),
        html.Div(style=CARD_STYLE, children=[
            html.Div(style={"marginBottom": "12px"}, children=[
                html.Label("Filtrar por Zona:", style={"fontSize": "12px", "fontWeight": "600",
                                                        "marginRight": "8px"}),
                dcc.Dropdown(
                    id="prod-zona-filter",
                    options=dropdown_opts,
                    value="TODAS",
                    clearable=False,
                    style={"width": "300px", "fontSize": "12px", "display": "inline-block"},
                ),
            ]),
            html.Div(id="prod-tabla-container"),
        ]),
    ])


def _build_producto_tabla(zona_filtro="TODAS"):
    df = _df_producto
    if df.empty:
        return _empty_msg()
    if zona_filtro and zona_filtro != "TODAS":
        df = df[df["zona"] == zona_filtro]
    if df.empty:
        return _empty_msg(f"No hay productos para zona: {zona_filtro}")

    total_venta_ytd_prod = df["venta_2026_ytd"].sum()
    total_gap_prod       = df["gap_venta"].sum()

    table_data = []
    for _, r in df.iterrows():
        dp = r["delta_precio_pct"]
        sem_p = _sem_precio(dp) if not (isinstance(dp, float) and math.isnan(dp)) else "⚫"
        table_data.append({
            "zona":           r["zona"],
            "codigo":         r["CODIGO"],
            "descripcion":    r["descripcion"],
            "categoria":      r["categoria"],
            "sem":            sem_p,
            "Cant PPTO":      f"{r['cant_ppto']:,.0f}",
            "P° PPTO":        _fmt(r["precio_ppto"]),
            "P° Real 2026":   _fmt(r["precio_real_2026"]) if pd.notna(r["precio_real_2026"]) else "—",
            "Δ Precio %":     _fmt_pct(dp),
            "Cant 2026 YTD":  f"{r['cant_2026_ytd']:,.0f}",
            "Venta 2026 YTD": _fmt_abs(r["venta_2026_ytd"]),
            "Gap $":          _fmt_abs(r["gap_venta"]),
        })
    table_data.append({
        "zona":           "",
        "codigo":         "",
        "descripcion":    "── TOTAL",
        "categoria":      "",
        "sem":            "",
        "Cant PPTO":      "",
        "P° PPTO":        "",
        "P° Real 2026":   "",
        "Δ Precio %":     "",
        "Cant 2026 YTD":  "",
        "Venta 2026 YTD": _fmt_abs(total_venta_ytd_prod),
        "Gap $":          _fmt_abs(total_gap_prod),
    })

    display_cols = ["zona", "codigo", "descripcion", "categoria", "sem", "Cant PPTO",
                    "P° PPTO", "P° Real 2026", "Δ Precio %", "Cant 2026 YTD",
                    "Venta 2026 YTD", "Gap $"]

    return dash_table.DataTable(
        data=table_data,
        columns=[{"name": c, "id": c} for c in display_cols],
        style_table={"overflowX": "auto"},
        style_header={"backgroundColor": "#1F3864", "color": "white",
                      "fontWeight": "bold", "fontSize": "12px"},
        style_cell={"fontSize": "12px", "padding": "6px 10px",
                    "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
        style_cell_conditional=[
            {"if": {"column_id": "zona"}, "textAlign": "left"},
            {"if": {"column_id": "codigo"}, "textAlign": "left"},
            {"if": {"column_id": "descripcion"}, "textAlign": "left", "maxWidth": "220px",
             "overflow": "hidden", "textOverflow": "ellipsis"},
            {"if": {"column_id": "categoria"}, "textAlign": "left"},
            {"if": {"column_id": "sem"}, "textAlign": "center", "width": "40px"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{descripcion} = "── TOTAL"'}, "fontWeight": "700", "backgroundColor": "#e8edf7"},
        ],
        page_size=50,
        sort_action="native",
        filter_action="native",
    )


# ─── Layout Tab: Desalineación PPTO ───────────────────────────────────────────

def _render_desa_zona(categoria):
    """Detalle del presupuesto 2026 por nombre/cliente para una categoría (solo [PPTO 2026])."""
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    cat_in = _cat_in(categoria.replace("'", "''"))
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(f"""
            WITH zona_pesos AS (
                SELECT LTRIM(RTRIM(Zona)) AS zona,
                       SUM(CASE WHEN ANIOMES <= {_ANO_ACT * 100 + _MES_ACT}
                                THEN TRY_CAST([ META ] AS float) ELSE 0 END)
                       / NULLIF(SUM(TRY_CAST([ META ] AS float)), 0) AS peso_ytd
                FROM Metas_KAM
                WHERE ANIOMES >= {_ANO_ACT}01 AND ANIOMES <= {_ANO_ACT}12
                GROUP BY Zona
            ),
            base AS (
                SELECT ISNULL(NULLIF(LTRIM(RTRIM(p.NOMBRE)),''),'Sin descripci\u00f3n') AS nombre,
                       ISNULL(LTRIM(RTRIM(p.VENDEDOR_ACTUAL)),'(sin zona)')             AS zona,
                       ISNULL(LTRIM(RTRIM(p.[SITUACI\u00d3N ])),'(sin situaci\u00f3n)') AS situacion,
                       CASE WHEN LTRIM(RTRIM(ISNULL(p.RUT,''))) IN ('','0')
                            THEN NULL ELSE LTRIM(RTRIM(p.RUT)) END AS rut,
                       NULLIF(LTRIM(RTRIM(ISNULL(p.CODIGO,''))), '') AS codigo,
                       TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) AS ppto_linea,
                       {_sql_ppto_ytd()} AS ppto_ytd_linea
                FROM [PPTO 2026] p
                LEFT JOIN zona_pesos zp ON zp.zona = LTRIM(RTRIM(ISNULL(p.VENDEDOR_ACTUAL,'')))
                WHERE TRY_CAST(REPLACE(ISNULL(p.[PPTO 2026],'0'),',','.') AS float) > 0
                  AND ISNULL(LTRIM(RTRIM(p.[CATEGOR\u00cdA 2026])),'(sin cat)') {cat_in}
            ),
            venta_25 AS (
                SELECT LTRIM(RTRIM(RUT)) AS rut,
                       SUM(CASE WHEN MES <= {_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025_ytd
                FROM DW_TOTAL_FACTURA
                WHERE ANO = 2025 AND {_DW_FILTRO}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                GROUP BY LTRIM(RTRIM(RUT))
            )
            SELECT b.nombre, b.zona, b.situacion,
                   SUM(b.ppto_linea)        AS ppto_anual,
                   SUM(b.ppto_ytd_linea)    AS ppto_ytd,
                   COUNT(DISTINCT b.codigo) AS n_productos,
                   ISNULL(MAX(v25.venta_2025_ytd), 0) AS venta_2025_ytd
            FROM base b
            LEFT JOIN venta_25 v25 ON v25.rut = b.rut
            GROUP BY b.nombre, b.zona, b.situacion
            ORDER BY SUM(b.ppto_linea) DESC
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
    except Exception as e:
        print(f"[ERROR _render_desa_zona] {e}")
        return _empty_msg(f"Error cargando detalle: {e}")

    df = pd.DataFrame.from_records(rows, columns=cols)
    if df.empty:
        return _empty_msg("Sin datos para esta categoría.")

    for c in ["ppto_anual", "ppto_ytd", "n_productos", "venta_2025_ytd"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    tot_a = df["ppto_anual"].sum()
    df["pct_cat"] = df["ppto_anual"] / tot_a * 100 if tot_a > 0 else 0.0
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    tbl_header = {"backgroundColor": "#1a5276", "color": "white",
                  "fontWeight": "bold", "fontSize": "12px"}
    tbl_cell   = {"fontSize": "12px", "padding": "6px 10px",
                  "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"}

    table_data = []
    for _, r in df.iterrows():
        table_data.append({
            "Nombre":           r["nombre"],
            "Zona":             r["zona"],
            "Situación":        r["situacion"],
            "PPTO Anual":       _fmt_abs(r["ppto_anual"]),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(r["ppto_ytd"]),
            f"Venta 2025 YTD":  _fmt_abs(r["venta_2025_ytd"]) if r["venta_2025_ytd"] > 0 else "—",
            "% del Cat":        _fmt_pct(r["pct_cat"]),
            "# Productos":      int(r["n_productos"]),
        })
    table_data.append({
        "Nombre":           "── TOTAL",
        "Zona":             "",
        "Situación":        "",
        "PPTO Anual":       _fmt_abs(df["ppto_anual"].sum()),
        f"PPTO YTD ({mes_nombre})": _fmt_abs(df["ppto_ytd"].sum()),
        f"Venta 2025 YTD":  _fmt_abs(df["venta_2025_ytd"].sum()),
        "% del Cat":        "100.0%",
        "# Productos":      int(df["n_productos"].sum()),
    })

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"📋 Detalle por Nombre — {categoria}", "#1a5276"),
        dash_table.DataTable(
            data=table_data,
            columns=[{"name": c, "id": c} for c in list(table_data[0].keys())],
            style_table={"overflowX": "auto"},
            style_header=tbl_header,
            style_cell=tbl_cell,
            style_cell_conditional=[
                {"if": {"column_id": "Nombre"},    "textAlign": "left"},
                {"if": {"column_id": "Zona"},      "textAlign": "left"},
                {"if": {"column_id": "Situación"}, "textAlign": "left"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{Nombre} = "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#d6eaf8"},
            ],
            sort_action="native",
            filter_action="native",
            page_size=30,
        ),
    ])


def _build_desa_pivot(df):
    """Construye tabla pivotada: 1 fila por categoría, situaciones como columnas."""
    if df.empty:
        return [], ["Categoría"]

    situaciones = sorted(df["situacion"].unique().tolist())
    table_data = []

    for cat, grp in df.groupby("categoria", sort=False):
        row = {"Categoría": cat}
        total_ppto = grp["ppto_anual"].sum()
        total_ytd  = grp["ppto_ytd"].sum()
        row["PPTO Total"] = _fmt_abs(total_ppto)
        row["PPTO YTD"]   = _fmt_abs(total_ytd)
        for sit in situaciones:
            match = grp[grp["situacion"] == sit]
            val = match["ppto_anual"].sum() if not match.empty else 0
            pct = val / total_ppto * 100 if total_ppto > 0 and val > 0 else 0
            row[sit] = f"{_fmt_abs(val)}  ({pct:.0f}%)" if val > 0 else "—"
        table_data.append(row)

    # Fila TOTAL
    total_row = {"Categoría": "── TOTAL"}
    total_row["PPTO Total"] = _fmt_abs(df.drop_duplicates("categoria")["ppto_anual_cat"].sum())
    total_row["PPTO YTD"]   = _fmt_abs(df.drop_duplicates("categoria")["ppto_ytd_cat"].sum())
    for sit in situaciones:
        val = df[df["situacion"] == sit]["ppto_anual"].sum()
        total_row[sit] = _fmt_abs(val) if val > 0 else "—"
    table_data.append(total_row)

    cols_main = ["Categoría", "PPTO Total", "PPTO YTD"] + situaciones
    return table_data, cols_main


def layout_desalineacion():
    df = _df_desalineacion
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    if df.empty:
        error_detail = _last_desa_error or "Sin detalle disponible."
        return html.Div(style=CARD_STYLE, children=[
            _seccion("🎯 Análisis de Desalineación PPTO 2026"),
            _empty_msg("No se pudieron cargar datos."),
            html.Pre(error_detail, style={"fontSize": "11px", "color": "#c0392b",
                "background": "#fdf2f8", "padding": "10px", "borderRadius": "4px",
                "overflowX": "auto", "marginTop": "8px", "whiteSpace": "pre-wrap"}),
        ])

    # ── Totales globales ───────────────────────────────────────────────────────
    df_cat       = df.drop_duplicates(subset=["categoria"])
    tot_anual    = df_cat["ppto_anual_cat"].sum()
    tot_ytd      = df_cat["ppto_ytd_cat"].sum()
    n_clientes   = int(df["n_clientes"].sum())
    n_productos  = int(df["n_productos"].sum())

    # Distribución por situación (total cross-categorías)
    df_sit = df.groupby("situacion", as_index=False).agg(
        ppto_anual=("ppto_anual", "sum"),
        n_clientes=("n_clientes", "sum"),
        n_productos=("n_productos", "sum"),
    ).sort_values("ppto_anual", ascending=False)
    df_sit["pct"] = df_sit["ppto_anual"] / tot_anual * 100 if tot_anual > 0 else 0

    # Clasificación estratégica
    sit_licitado   = {"LICITACIÓN", "CENABAST", "CONVENIO MARCO"}
    sit_incremental = {"INCREMENTALES"}
    sit_riesgo     = {"NO LICITADO", "SIN PPTO 2026"}
    ppto_licitado   = df[df["situacion"].isin(sit_licitado)]["ppto_anual"].sum()
    ppto_increment  = df[df["situacion"].isin(sit_incremental)]["ppto_anual"].sum()
    ppto_riesgo     = df[df["situacion"].isin(sit_riesgo)]["ppto_anual"].sum()
    pct_licitado    = ppto_licitado / tot_anual * 100 if tot_anual > 0 else 0
    pct_increment   = ppto_increment / tot_anual * 100 if tot_anual > 0 else 0
    pct_riesgo      = ppto_riesgo / tot_anual * 100 if tot_anual > 0 else 0

    # ── Gráfico distribución por situación ────────────────────────────────────
    COLORES_SIT = {
        "LICITACIÓN":     "#2E75B6",
        "CONVENIO MARCO": "#27ae60",
        "CENABAST":       "#1abc9c",
        "INCREMENTALES":  "#e67e22",
        "NO LICITADO":    "#c0392b",
        "SIN PPTO 2026":  "#95a5a6",
    }
    fig_sit = go.Figure(go.Bar(
        x=df_sit["ppto_anual"],
        y=df_sit["situacion"],
        orientation="h",
        marker_color=[COLORES_SIT.get(s, "#7f8c8d") for s in df_sit["situacion"]],
        text=[f"{_fmt(v)}  ({p:.1f}%)" for v, p in zip(df_sit["ppto_anual"], df_sit["pct"])],
        textposition="outside",
        hovertemplate="<b>%{y}</b><br>PPTO: %{x:,.0f}<extra></extra>",
    ))
    fig_sit.update_layout(
        margin=dict(l=140, r=120, t=10, b=10), height=240,
        xaxis=dict(showticklabels=False, showgrid=False),
        yaxis=dict(tickfont=dict(size=12)),
        plot_bgcolor="white", paper_bgcolor="white",
        font=dict(family="Segoe UI, sans-serif", size=12),
    )

    # ── Gráfico PPTO por categoría (stacked por situación) ────────────────────
    cats  = sorted(df["categoria"].unique())
    sits  = df_sit["situacion"].tolist()
    fig_cat = go.Figure()
    for sit in sits:
        vals = []
        for cat in cats:
            row = df[(df["categoria"] == cat) & (df["situacion"] == sit)]
            vals.append(row["ppto_anual"].sum() if not row.empty else 0)
        fig_cat.add_trace(go.Bar(
            name=sit, x=cats, y=vals,
            marker_color=COLORES_SIT.get(sit, "#7f8c8d"),
            hovertemplate=f"<b>{sit}</b><br>%{{x}}: %{{y:,.0f}}<extra></extra>",
        ))
    fig_cat.update_layout(
        barmode="stack",
        margin=dict(l=10, r=10, t=10, b=10), height=260,
        legend=dict(orientation="h", y=-0.25, font=dict(size=11)),
        xaxis=dict(tickfont=dict(size=12)),
        yaxis=dict(showticklabels=False, showgrid=False),
        plot_bgcolor="white", paper_bgcolor="white",
        font=dict(family="Segoe UI, sans-serif", size=12),
    )

    # ── Tabla resumen por situación ────────────────────────────────────────────
    sit_table = []
    for _, r in df_sit.iterrows():
        tipo = ("🔵 Competitivo" if r["situacion"] in sit_licitado
                else "🟠 Incremental" if r["situacion"] in sit_incremental
                else "🔴 En riesgo")
        sit_table.append({
            "Situación":    r["situacion"],
            "Tipo":         tipo,
            "PPTO Anual":   _fmt_abs(r["ppto_anual"]),
            "% del Total":  _fmt_pct(r["pct"]),
            "# Clientes":   int(r["n_clientes"]),
            "# Productos":  int(r["n_productos"]),
        })

    # ── Pivot categoría × situación ────────────────────────────────────────────
    tbl_header = {"backgroundColor": "#1F3864", "color": "white",
                  "fontWeight": "bold", "fontSize": "12px"}
    tbl_cell   = {"fontSize": "12px", "padding": "6px 10px",
                  "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"}
    table_data, cols_main = _build_desa_pivot(df)

    return html.Div(children=[

        # ── Sección 1: KPIs + clasificación estratégica ───────────────────────
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"🎯 Análisis de Desalineación PPTO 2026 — YTD {mes_nombre} {_ANO_ACT}"),
            _kpi_row([
                _kpi_card("PPTO Anual 2026",           _fmt(tot_anual),   "#1F3864"),
                _kpi_card(f"PPTO YTD ({mes_nombre})",  _fmt(tot_ytd),     "#2E75B6"),
                _kpi_card("🔵 Competitivo (Lic+CM+CEN)", f"{_fmt(ppto_licitado)} ({pct_licitado:.0f}%)", "#2E75B6"),
                _kpi_card("🟠 Incremental",             f"{_fmt(ppto_increment)} ({pct_increment:.0f}%)", "#e67e22"),
                _kpi_card("🔴 En Riesgo (No Lic+S/P)",  f"{_fmt(ppto_riesgo)} ({pct_riesgo:.0f}%)",
                          "#c0392b" if pct_riesgo > 30 else "#e67e22"),
                _kpi_card("# Clientes",                str(n_clientes),   "#1a5276"),
                _kpi_card("# Productos",               str(n_productos),  "#145a32"),
            ]),
        ]),

        # ── Sección 2: Gráficos ───────────────────────────────────────────────
        html.Div(style={**CARD_STYLE, "display": "grid",
                        "gridTemplateColumns": "1fr 1fr", "gap": "24px"}, children=[
            html.Div(children=[
                _seccion("Distribución por Situación"),
                dcc.Graph(figure=fig_sit, config={"displayModeBar": False}),
            ]),
            html.Div(children=[
                _seccion("PPTO por Categoría × Situación"),
                dcc.Graph(figure=fig_cat, config={"displayModeBar": False}),
            ]),
        ]),

        # ── Sección 3: Resumen por situación ──────────────────────────────────
        html.Div(style=CARD_STYLE, children=[
            _seccion("Resumen por Situación"),
            dash_table.DataTable(
                data=sit_table,
                columns=[{"name": c, "id": c} for c in list(sit_table[0].keys())],
                style_table={"overflowX": "auto"},
                style_header=tbl_header,
                style_cell=tbl_cell,
                style_cell_conditional=[
                    {"if": {"column_id": "Situación"}, "textAlign": "left", "fontWeight": "600"},
                    {"if": {"column_id": "Tipo"},      "textAlign": "left"},
                    {"if": {"column_id": "# Clientes"},  "textAlign": "center"},
                    {"if": {"column_id": "# Productos"}, "textAlign": "center"},
                ],
                style_data_conditional=[
                    {"if": {"filter_query": '{Tipo} contains "riesgo"'},
                     "backgroundColor": "#fdf2f2"},
                    {"if": {"filter_query": '{Tipo} contains "Incremental"'},
                     "backgroundColor": "#fef9f0"},
                ],
                sort_action="native",
                page_size=10,
            ),
        ]),

        # ── Sección 4: Pivot categoría × situación con drill-down ─────────────
        html.Div(style=CARD_STYLE, children=[
            _seccion("Detalle por Categoría × Situación"),
            html.Div(style={"display": "flex", "alignItems": "center", "gap": "12px",
                            "marginBottom": "10px"}, children=[
                html.Label("Filtrar categoría:",
                           style={"fontSize": "12px", "color": "#555", "whiteSpace": "nowrap"}),
                dcc.Dropdown(
                    id="desa-cat-filtro",
                    options=[{"label": "Todas", "value": "__all__"}] +
                            [{"label": c, "value": c} for c in sorted(df["categoria"].unique())],
                    value="__all__",
                    clearable=False,
                    style={"width": "260px", "fontSize": "12px"},
                ),
                html.Span("← Haz clic en una fila para ver el detalle de clientes/productos",
                          style={"fontSize": "11px", "color": "#888"}),
            ]),
            dash_table.DataTable(
                id="desa-main-tabla",
                data=table_data,
                columns=[{"name": c, "id": c} for c in cols_main],
                style_table={"overflowX": "auto"},
                style_header=tbl_header,
                style_cell=tbl_cell,
                style_cell_conditional=[
                    {"if": {"column_id": "Categoría"},  "textAlign": "left", "fontWeight": "600"},
                    {"if": {"column_id": "PPTO Total"}, "fontWeight": "600"},
                ],
                style_data_conditional=[
                    {"if": {"filter_query": '{Categoría} = "── TOTAL"'},
                     "fontWeight": "700", "backgroundColor": "#e8edf7"},
                ],
                sort_action="native",
                row_selectable="single",
                selected_rows=[],
                page_size=20,
            ),
        ]),

        # Drill-down
        html.Div(id="desa-detalle-container"),
    ])


# ─── Layout Tab: Análisis de Precios ─────────────────────────────────────────

def _precios_kpis(df):
    """Construye fila KPI para el tab de precios."""
    ep = df["efecto_precio"].dropna()
    ev = df["efecto_volumen"].dropna()
    ganancia_precio = ep[ep > 0].sum()
    perdida_precio  = ep[ep < 0].sum()
    neto_precio     = ep.sum()
    neto_volumen    = ev.sum()
    n_alza          = (df["delta_precio_pct"] > 2).sum()
    n_baja          = (df["delta_precio_pct"] < -2).sum()

    precio_pond_2026 = (df["venta_2026"].sum() / df["cant_2026"].sum()
                        if df["cant_2026"].sum() > 0 else float("nan"))
    precio_pond_2025 = (df["venta_2025"].sum() / df["cant_2025"].sum()
                        if df["cant_2025"].sum() > 0 else float("nan"))
    delta_pond_pct = ((precio_pond_2026 - precio_pond_2025) / precio_pond_2025 * 100
                      if pd.notna(precio_pond_2025) and precio_pond_2025 > 0 else float("nan"))

    sem_neto = "🟢" if neto_precio >= 0 else "🔴"
    sem_delta = ("🟢" if pd.notna(delta_pond_pct) and delta_pond_pct > 0
                 else "🔴" if pd.notna(delta_pond_pct) and delta_pond_pct < 0 else "⚫")

    return _kpi_row([
        _kpi_card("💰 Ganancia por precio", _fmt_abs(ganancia_precio), "#1e7e34"),
        _kpi_card("📉 Pérdida por precio",  _fmt_abs(abs(perdida_precio)), "#c0392b"),
        _kpi_card(f"{sem_neto} Impacto neto precio", _fmt_abs(neto_precio),
                  "#1e7e34" if neto_precio >= 0 else "#c0392b"),
        _kpi_card("📦 Impacto por volumen", _fmt_abs(neto_volumen), "#2E75B6"),
        _kpi_card(f"{sem_delta} Precio prom ponderado",
                  f"{_fmt_abs(precio_pond_2026)} / {_fmt_abs(precio_pond_2025)}",
                  "#7b2d8b",
                  f"Δ {_fmt_pct(delta_pond_pct)} vs 2025"),
        _kpi_card("🟢 Productos precio al alza", f"{int(n_alza):,}", "#1e7e34",
                  "Δ% > +2%"),
        _kpi_card("🔴 Productos precio a la baja", f"{int(n_baja):,}", "#c0392b",
                  "Δ% < -2%"),
    ])


def _precios_tabla_categoria(df):
    """Tabla resumen por categoría con precio 2025, 2026, delta, efectos."""
    if df.empty:
        return _empty_msg()
    rows = []
    for cat, g in df.groupby("categoria", sort=False):
        p26 = g["venta_2026"].sum() / g["cant_2026"].sum() if g["cant_2026"].sum() > 0 else float("nan")
        p25 = g["venta_2025"].sum() / g["cant_2025"].sum() if g["cant_2025"].sum() > 0 else float("nan")
        dpct = (p26 - p25) / p25 * 100 if pd.notna(p25) and p25 > 0 and pd.notna(p26) else float("nan")
        ep = g["efecto_precio"].dropna().sum()
        ev = g["efecto_volumen"].dropna().sum()
        sem = "🟢" if pd.notna(dpct) and dpct > 0 else "🔴" if pd.notna(dpct) and dpct < 0 else "⚫"
        rows.append({
            "cat": cat,
            "Categoría": cat,
            "Precio 2025": _fmt_abs(p25) if pd.notna(p25) else "—",
            "Precio 2026": _fmt_abs(p26) if pd.notna(p26) else "—",
            "Δ Precio %": f"{sem} {_fmt_pct(dpct)}" if pd.notna(dpct) else "—",
            "Efecto precio": _fmt_abs(ep),
            "Efecto volumen": _fmt_abs(ev),
            "Venta 2026": _fmt_abs(g["venta_2026"].sum()),
            "Venta 2025": _fmt_abs(g["venta_2025"].sum()),
        })
    # Fila TOTAL
    tp26 = df["venta_2026"].sum() / df["cant_2026"].sum() if df["cant_2026"].sum() > 0 else float("nan")
    tp25 = df["venta_2025"].sum() / df["cant_2025"].sum() if df["cant_2025"].sum() > 0 else float("nan")
    tdpct = (tp26 - tp25) / tp25 * 100 if pd.notna(tp25) and tp25 > 0 and pd.notna(tp26) else float("nan")
    tsem = "🟢" if pd.notna(tdpct) and tdpct > 0 else "🔴" if pd.notna(tdpct) and tdpct < 0 else "⚫"
    rows.append({
        "cat": "── TOTAL",
        "Categoría": "── TOTAL",
        "Precio 2025": _fmt_abs(tp25) if pd.notna(tp25) else "—",
        "Precio 2026": _fmt_abs(tp26) if pd.notna(tp26) else "—",
        "Δ Precio %": f"{tsem} {_fmt_pct(tdpct)}" if pd.notna(tdpct) else "—",
        "Efecto precio":  _fmt_abs(df["efecto_precio"].dropna().sum()),
        "Efecto volumen": _fmt_abs(df["efecto_volumen"].dropna().sum()),
        "Venta 2026": _fmt_abs(df["venta_2026"].sum()),
        "Venta 2025": _fmt_abs(df["venta_2025"].sum()),
    })
    columns = ["Categoría", "Precio 2025", "Precio 2026",
               "Δ Precio %", "Efecto precio", "Efecto volumen",
               "Venta 2026", "Venta 2025"]
    tbl_header = {"backgroundColor": "#1F3864", "color": "white",
                  "fontWeight": "700", "fontSize": "12px"}
    tbl_cell   = {"fontSize": "12px", "padding": "6px 10px", "textAlign": "right"}
    return dash_table.DataTable(
        id="precios-cat-tabla",
        data=rows,
        columns=[{"name": c, "id": c} for c in columns],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": "Categoría"}, "textAlign": "left", "fontWeight": "600"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Δ Precio %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Δ Precio %} contains "🔴"'}, "color": "#c0392b"},
            {"if": {"filter_query": '{Categoría} = "── TOTAL"'},
             "fontWeight": "700", "backgroundColor": "#e8edf7"},
        ],
        sort_action="native",
        row_selectable="single",
        selected_rows=[],
        page_action="none",
    )


def _precios_tabla_zona(df):
    """Tabla por zona con precio ponderado 2025 vs 2026."""
    if df.empty:
        return _empty_msg()
    rows = []
    for zona, g in df.groupby("zona", sort=False):
        p26 = g["venta_2026"].sum() / g["cant_2026"].sum() if g["cant_2026"].sum() > 0 else float("nan")
        p25 = g["venta_2025"].sum() / g["cant_2025"].sum() if g["cant_2025"].sum() > 0 else float("nan")
        dpct = (p26 - p25) / p25 * 100 if pd.notna(p25) and p25 > 0 and pd.notna(p26) else float("nan")
        ep = g["efecto_precio"].dropna().sum()
        ev = g["efecto_volumen"].dropna().sum()
        sem = "🟢" if pd.notna(dpct) and dpct > 0 else "🔴" if pd.notna(dpct) and dpct < 0 else "⚫"
        rows.append({
            "zona": zona,
            "Zona": zona,
            "Precio 2025": _fmt_abs(p25) if pd.notna(p25) else "—",
            "Precio 2026": _fmt_abs(p26) if pd.notna(p26) else "—",
            "Δ Precio %": f"{sem} {_fmt_pct(dpct)}" if pd.notna(dpct) else "—",
            "Efecto precio": _fmt_abs(ep),
            "Efecto volumen": _fmt_abs(ev),
            "Venta 2026": _fmt_abs(g["venta_2026"].sum()),
            "Venta 2025": _fmt_abs(g["venta_2025"].sum()),
        })
    # Fila TOTAL
    tp26 = df["venta_2026"].sum() / df["cant_2026"].sum() if df["cant_2026"].sum() > 0 else float("nan")
    tp25 = df["venta_2025"].sum() / df["cant_2025"].sum() if df["cant_2025"].sum() > 0 else float("nan")
    tdpct = (tp26 - tp25) / tp25 * 100 if pd.notna(tp25) and tp25 > 0 and pd.notna(tp26) else float("nan")
    tsem = "🟢" if pd.notna(tdpct) and tdpct > 0 else "🔴" if pd.notna(tdpct) and tdpct < 0 else "⚫"
    rows.append({
        "zona": "── TOTAL",
        "Zona": "── TOTAL",
        "Precio 2025": _fmt_abs(tp25) if pd.notna(tp25) else "—",
        "Precio 2026": _fmt_abs(tp26) if pd.notna(tp26) else "—",
        "Δ Precio %": f"{tsem} {_fmt_pct(tdpct)}" if pd.notna(tdpct) else "—",
        "Efecto precio":  _fmt_abs(df["efecto_precio"].dropna().sum()),
        "Efecto volumen": _fmt_abs(df["efecto_volumen"].dropna().sum()),
        "Venta 2026": _fmt_abs(df["venta_2026"].sum()),
        "Venta 2025": _fmt_abs(df["venta_2025"].sum()),
    })
    tbl_header = {"backgroundColor": "#1F3864", "color": "white",
                  "fontWeight": "700", "fontSize": "12px"}
    tbl_cell   = {"fontSize": "12px", "padding": "6px 10px", "textAlign": "right"}
    return dash_table.DataTable(
        id="precios-zona-tabla",
        data=rows,
        columns=[{"name": c, "id": c} for c in ["Zona", "Precio 2025", "Precio 2026",
                                                  "Δ Precio %", "Efecto precio",
                                                  "Efecto volumen", "Venta 2026", "Venta 2025"]],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": "Zona"}, "textAlign": "left", "fontWeight": "600"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Δ Precio %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Δ Precio %} contains "🔴"'}, "color": "#c0392b"},
            {"if": {"filter_query": '{Zona} = "── TOTAL"'},
             "fontWeight": "700", "backgroundColor": "#e8edf7"},
        ],
        sort_action="native",
        row_selectable="single",
        selected_rows=[],
        page_action="none",
    )


def _precios_top_productos(df, n=20):
    """Dos tablas: top N productos con mayor ganancia y mayor pérdida de precio."""
    if df.empty:
        return _empty_msg()
    sub = df.dropna(subset=["efecto_precio"]).copy()
    sub["_ep"] = sub["efecto_precio"]

    top_gain = sub[sub["_ep"] > 0].nlargest(n, "_ep")
    top_loss = sub[sub["_ep"] < 0].nsmallest(n, "_ep")

    def _tbl(data, color, tid):
        rows = []
        for _, r in data.iterrows():
            dpct = r.get("delta_precio_pct", float("nan"))
            ref  = r.get("ref_precio", "real")
            p25_label = (
                f"📋 {_fmt_abs(r['precio_2025'])}" if ref == "ppto" and pd.notna(r.get("precio_2025"))
                else _fmt_abs(r["precio_2025"]) if pd.notna(r.get("precio_2025"))
                else "—"
            )
            rows.append({
                "Zona": r["zona"],
                "Código": r["codigo"],
                "Descripción": r["descripcion"],
                "Ref. precio": p25_label,
                "Precio 2026": _fmt_abs(r["precio_2026"]) if pd.notna(r.get("precio_2026")) else "—",
                "Δ %": _fmt_pct(dpct) if pd.notna(dpct) else "—",
                "Efecto precio": _fmt_abs(r["_ep"]),
                "Venta 2026": _fmt_abs(r["venta_2026"]),
            })
        # Fila TOTAL
        rows.append({
            "Zona": "", "Código": "", "Descripción": "── TOTAL",
            "Precio 2025": "—", "Precio 2026": "—", "Δ %": "—",
            "Efecto precio": _fmt_abs(data["_ep"].sum()),
            "Venta 2026":    _fmt_abs(data["venta_2026"].sum()),
        })
        tbl_header = {"backgroundColor": color, "color": "white",
                      "fontWeight": "700", "fontSize": "11px"}
        tbl_cell   = {"fontSize": "11px", "padding": "5px 8px", "textAlign": "right"}
        return dash_table.DataTable(
            id=tid,
            data=rows,
            columns=[{"name": c, "id": c} for c in
                     ["Zona", "Código", "Descripción", "Ref. precio",
                      "Precio 2026", "Δ %", "Efecto precio", "Venta 2026"]],
            style_header=tbl_header,
            style_cell=tbl_cell,
            style_cell_conditional=[
                {"if": {"column_id": "Descripción"}, "textAlign": "left"},
                {"if": {"column_id": "Zona"},        "textAlign": "left"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{Descripción} = "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            sort_action="native",
            page_action="none",
        )

    return html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr",
                            "gap": "24px"}, children=[
        html.Div(children=[
            _seccion(f"🟢 Top {n} — Mayor ganancia de precio", "#1e7e34"),
            _tbl(top_gain, "#1e7e34", "precios-top-gain-tbl"),
        ]),
        html.Div(children=[
            _seccion(f"🔴 Top {n} — Mayor pérdida de precio", "#c0392b"),
            _tbl(top_loss, "#c0392b", "precios-top-loss-tbl"),
        ]),
    ])


def _precios_grafico_categoria(df):
    """Gráfico de barras: efecto precio y efecto volumen por categoría."""
    if df.empty:
        return go.Figure()
    cats = []
    ep_vals = []
    ev_vals = []
    for cat, g in df.groupby("categoria", sort=False):
        cats.append(cat)
        ep_vals.append(g["efecto_precio"].dropna().sum())
        ev_vals.append(g["efecto_volumen"].dropna().sum())

    fig = go.Figure()
    fig.add_trace(go.Bar(
        name="Efecto precio",
        x=cats, y=ep_vals,
        marker_color=["#1e7e34" if v >= 0 else "#c0392b" for v in ep_vals],
        text=[_fmt(v) for v in ep_vals],
        textposition="outside",
    ))
    fig.add_trace(go.Bar(
        name="Efecto volumen",
        x=cats, y=ev_vals,
        marker_color="#2E75B6",
        opacity=0.6,
        text=[_fmt(v) for v in ev_vals],
        textposition="outside",
    ))
    fig.update_layout(
        barmode="group",
        height=320,
        margin={"t": 20, "b": 40, "l": 60, "r": 20},
        legend={"orientation": "h", "yanchor": "bottom", "y": 1.02, "x": 0},
        plot_bgcolor="white",
        paper_bgcolor="white",
        font={"size": 11},
        yaxis={"tickformat": "$,.0f", "gridcolor": "#eee"},
        xaxis={"tickfont": {"size": 11}},
    )
    return fig


def layout_precios():
    df = _df_precios
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("💲 Análisis de Precios"),
            _empty_msg("Sin datos de precios. Presione Actualizar con VPN activa."),
        ])

    cats = sorted(df["categoria"].dropna().unique().tolist())
    zonas = sorted(df["zona"].dropna().unique().tolist())

    cat_opts  = [{"label": "Todas las categorías", "value": "__all__"}] + \
                [{"label": c, "value": c} for c in cats]
    zona_opts = [{"label": "Todas las zonas", "value": "__all__"}] + \
                [{"label": z, "value": z} for z in zonas]

    return html.Div([
        # KPIs globales
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"💲 Análisis de Precios — YTD {mes_nombre} 2026 vs 2025"),
            _precios_kpis(df),
        ]),

        # Glosario de columnas
        html.Div(style={**CARD_STYLE, "padding": "12px 16px"}, children=[
            html.Details(children=[
                html.Summary(style={"fontWeight": "700", "fontSize": "12px",
                                    "color": "#1F3864", "cursor": "pointer"},
                             children="📖 Glosario de columnas — ¿qué significa cada métrica?"),
                html.Div(style={"display": "grid", "gridTemplateColumns": "repeat(3,1fr)",
                                "gap": "10px", "marginTop": "10px", "fontSize": "11px",
                                "color": "#444"}, children=[
                    html.Div([html.B("Precio 2025 / 2026"), html.Br(),
                              "Precio promedio ponderado = Venta total ÷ Unidades vendidas. "
                              "Refleja el precio real de venta, no el precio de lista."]),
                    html.Div([html.B("Δ Precio %"), html.Br(),
                              "Variación porcentual del precio entre 2025 y 2026. "
                              "🟢 sube / 🔴 baja. Refleja si se vendió más caro o más barato."]),
                    html.Div([html.B("Efecto precio"), html.Br(),
                              "Impacto en ingresos atribuible al cambio de precio: "
                              "(Precio 2026 − Precio 2025) × Unidades 2026. "
                              "Positivo = ganancia por precio, negativo = pérdida."]),
                    html.Div([html.B("Efecto volumen"), html.Br(),
                              "Impacto en ingresos atribuible al cambio de unidades: "
                              "(Unidades 2026 − Unidades 2025) × Precio 2025. "
                              "Positivo = se vendió más, negativo = se vendió menos."]),
                    html.Div([html.B("Venta 2026 / 2025"), html.Br(),
                              "Ingreso neto total (sin IVA) acumulado de enero al mes actual, "
                              "para cada año. Base de comparación YTD."]),
                    html.Div([html.B("Efecto precio + volumen"), html.Br(),
                              "Suma de ambos efectos ≈ variación total de venta entre años. "
                              "Un efecto precio negativo con volumen positivo indica que "
                              "se vendió más cantidad pero a menor precio."]),
                    html.Div([html.B("Ref. precio (📋 = PPTO)"), html.Br(),
                              "Precio de referencia para el cálculo del delta. "
                              "Sin ícono = precio real promedio 2025. "
                              "📋 = producto nuevo en 2026 sin historial 2025 — "
                              "se usa el precio presupuestado 2026 como base de comparación."]),
                ]),
            ]),
        ]),

        # Filtros
        html.Div(style={**CARD_STYLE, "display": "flex", "gap": "16px",
                        "alignItems": "center", "padding": "12px 16px"}, children=[
            html.Span("Filtros:", style={"fontWeight": "600", "fontSize": "13px",
                                         "color": "#1F3864", "whiteSpace": "nowrap"}),
            dcc.Dropdown(id="precio-cat-filtro", options=cat_opts, value="__all__",
                         placeholder="Categoría", clearable=False,
                         style={"flex": "1", "fontSize": "12px"}),
            dcc.Dropdown(id="precio-zona-filtro", options=zona_opts, value="__all__",
                         placeholder="Zona", clearable=False,
                         style={"flex": "1", "fontSize": "12px"}),
        ]),

        # Tabla por categoría + tabla por zona
        html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr",
                         "gap": "16px", "marginBottom": "16px"}, children=[
            html.Div(style=CARD_STYLE, children=[
                _seccion("🗂️ Resumen por Categoría"),
                html.Div(id="precio-resumen-cat"),
            ]),
            html.Div(style=CARD_STYLE, children=[
                _seccion("🏢 Resumen por Zona"),
                html.Div(id="precio-resumen-zona"),
            ]),
        ]),

        # Top productos: ganancias vs pérdidas
        html.Div(style=CARD_STYLE, children=[
            _seccion("📦 Top Productos por Impacto de Precio"),
            html.Div(id="precio-top-productos"),
        ]),

        # Drill-down al hacer click en categoría
        html.Div(id="precio-drill-cat"),

        # Drill-down al hacer click en zona
        html.Div(id="precio-drill-zona"),
    ])


# ─── App Dash ─────────────────────────────────────────────────────────────────

app = dash.Dash(__name__, suppress_callback_exceptions=True)
app.title = "Análisis Comercial LBF"
server = app.server
server.secret_key = "LBF_ANALYTICS_2026_KEY"


# ─── Login layout ──────────────────────────────────────────────────────────────


# ─── Tab PPTO vs Venta ────────────────────────────────────────────────────────

def layout_ppto_vs_venta():
    df = _df_ppto_vs_venta
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))

    if df.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("🔬 PPTO vs Venta"),
            _empty_msg("Sin datos. Presione 🔄 Actualizar con VPN activa."),
        ])

    zonas = sorted([z for z in df["VENDEDOR_ACTUAL"].dropna().unique().tolist() if str(z).strip()])
    cats  = sorted([c for c in df["CATEGORIA_2026"].dropna().unique().tolist() if str(c).strip()])

    return html.Div([
        # Filtros + título
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"🔬 PPTO vs Venta — YTD {mes_nombre} 2026"),
            html.Div(style={"display": "flex", "gap": "20px", "flexWrap": "wrap",
                            "alignItems": "flex-end", "marginTop": "8px"}, children=[
                html.Div([
                    html.Label("Zona:", style={"fontSize": "12px", "fontWeight": "600",
                                               "marginBottom": "4px", "display": "block"}),
                    dcc.Dropdown(
                        id="pvv-zona-filter",
                        options=[{"label": "Todas las zonas", "value": "TODAS"}] +
                                [{"label": z, "value": z} for z in zonas],
                        value="TODAS",
                        clearable=False,
                        style={"width": "240px", "fontSize": "12px"},
                    ),
                ]),
                html.Div([
                    html.Label("Categoría:", style={"fontSize": "12px", "fontWeight": "600",
                                                    "marginBottom": "4px", "display": "block"}),
                    dcc.Dropdown(
                        id="pvv-cat-filter",
                        options=[{"label": "Todas las categorías", "value": "TODAS"}] +
                                [{"label": c, "value": c} for c in cats],
                        value="TODAS",
                        clearable=False,
                        style={"width": "200px", "fontSize": "12px"},
                    ),
                ]),
            ]),
        ]),
        # KPIs (poblados por callback)
        html.Div(id="pvv-kpi-container"),
        # Tabla principal por categoría (poblada por callback)
        html.Div(id="pvv-tabla-container"),
        # Drill-down de productos (poblado por callback al hacer clic en fila)
        html.Div(id="pvv-detalle-container"),
        # Análisis ejecutivo de caída (poblado por callback)
        html.Div(id="pvv-analisis-caida"),
    ])


# ─── Tab Resumen Ejecutivo ─────────────────────────────────────────────────────

def _generate_pptx():
    """Genera un archivo PPTX con resumen ejecutivo de 2 slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io, tempfile

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        AZUL       = RGBColor(0x1F, 0x38, 0x64)
        AZUL_LIGHT = RGBColor(0x2E, 0x75, 0xB6)
        VERDE      = RGBColor(0x1e, 0x7e, 0x34)
        ROJO       = RGBColor(0xc0, 0x39, 0x2b)
        GRIS       = RGBColor(0xF4, 0xF6, 0xFA)
        BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)

        def add_textbox(slide, left, top, width, height, text, bold=False,
                        fontsize=11, color=None, bg=None, align=PP_ALIGN.LEFT):
            txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            if bg:
                txb.fill.solid()
                txb.fill.fore_color.rgb = bg
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.bold = bold
            run.font.size = Pt(fontsize)
            if color:
                run.font.color.rgb = color
            return txb

        def add_rect(slide, left, top, width, height, color):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
        df_cat = _df_categoria
        df_ud  = _df_resumen_ud_grupo
        df_cai = _df_clientes_caida

        # ── SLIDE 1: PPTO vs Venta por Categoría ──────────────────────────────
        blank_layout = prs.slide_layouts[6]
        sl1 = prs.slides.add_slide(blank_layout)

        # Header
        add_rect(sl1, 0, 0, 13.33, 0.7, AZUL)
        add_textbox(sl1, 0.15, 0.1, 9, 0.5,
                    f"Análisis Presupuesto 2026 — Resumen Ejecutivo YTD {mes_nombre}",
                    bold=True, fontsize=16, color=BLANCO)
        add_textbox(sl1, 9.5, 0.12, 3.5, 0.45,
                    f"Datos al {_last_update}",
                    fontsize=9, color=RGBColor(0xAE, 0xC6, 0xE8), align=PP_ALIGN.RIGHT)

        # KPIs globales
        if not df_cat.empty:
            ppto_total  = df_cat["ppto_total"].sum()   if "ppto_total"     in df_cat.columns else df_cat["ppto_anual"].sum()
            ppto_ytd    = df_cat["total_ppto_ytd"].sum() if "total_ppto_ytd" in df_cat.columns else df_cat["ppto_ytd"].sum()
            venta_ytd   = df_cat["venta_2026_ytd"].sum()
            venta_25    = df_cat["venta_2025_ytd"].sum()
            alcance     = venta_ytd / ppto_ytd * 100 if ppto_ytd > 0 else 0
            var_25      = (venta_ytd - venta_25) / venta_25 * 100 if venta_25 > 0 else 0
            ppto_cod    = df_cat["ppto_anual"].sum()
            ppto_incr   = df_cat["ppto_incr_anual"].sum()  if "ppto_incr_anual"   in df_cat.columns else 0
            ppto_nuevos = df_cat["ppto_nuevos_anual"].sum() if "ppto_nuevos_anual" in df_cat.columns else 0

            kpis = [
                ("PPTO Total 2026", _fmt(ppto_total)),
                ("PPTO c/Código",   _fmt(ppto_cod)),
                ("PPTO Incremental",_fmt(ppto_incr)),
                ("PPTO Prod. Nuevos",_fmt(ppto_nuevos)),
                (f"Venta YTD {mes_nombre}", _fmt(venta_ytd)),
                ("Alcance %",       f"{alcance:.1f}%"),
                ("vs 2025",         f"{var_25:+.1f}%"),
            ]
            kpi_w = 13.0 / len(kpis)
            for i, (lbl, val) in enumerate(kpis):
                x = 0.15 + i * kpi_w
                col = VERDE if (lbl == "Alcance %" and alcance >= 80) else (ROJO if lbl == "Alcance %" and alcance < 50 else AZUL_LIGHT)
                add_rect(sl1, x, 0.8, kpi_w - 0.08, 0.75, col)
                add_textbox(sl1, x + 0.05, 0.82, kpi_w - 0.1, 0.28, lbl,
                            fontsize=7.5, color=BLANCO, align=PP_ALIGN.CENTER)
                add_textbox(sl1, x + 0.05, 1.1, kpi_w - 0.1, 0.38, val,
                            bold=True, fontsize=13, color=BLANCO, align=PP_ALIGN.CENTER)

        # Tabla por Categoría
        add_textbox(sl1, 0.15, 1.72, 6, 0.3, "PPTO vs Venta por Categoría",
                    bold=True, fontsize=10, color=AZUL)

        if not df_cat.empty:
            headers = ["Categoría", "PPTO c/Cód", "PPTO Incr.", "Prod. Nuevos", "PPTO Total", "Venta 2026", "Alcance %"]
            col_ws  = [1.6, 1.1, 1.0, 1.1, 1.1, 1.1, 0.85]
            row_h   = 0.28
            t_top   = 2.05
            t_left  = 0.15

            # Header row
            x = t_left
            for h, w in zip(headers, col_ws):
                add_rect(sl1, x, t_top, w - 0.02, row_h, AZUL)
                add_textbox(sl1, x + 0.03, t_top + 0.03, w - 0.06, row_h - 0.04,
                            h, bold=True, fontsize=7, color=BLANCO)
                x += w

            rows_data = []
            for _, r in df_cat.iterrows():
                pt  = r.get("ppto_total", r.get("ppto_anual", 0))
                pyt = r.get("total_ppto_ytd", r.get("ppto_ytd", 0))
                alc = r["venta_2026_ytd"] / pyt * 100 if pyt > 0 else 0
                rows_data.append([
                    str(r["categoria"])[:18],
                    _fmt(r["ppto_anual"]),
                    _fmt(r.get("ppto_incr_anual", 0)),
                    _fmt(r.get("ppto_nuevos_anual", 0)),
                    _fmt(pt),
                    _fmt(r["venta_2026_ytd"]),
                    f"{alc:.1f}%",
                ])
            # Total row
            pt_tot  = df_cat["ppto_total"].sum() if "ppto_total" in df_cat.columns else df_cat["ppto_anual"].sum()
            pyt_tot = df_cat["total_ppto_ytd"].sum() if "total_ppto_ytd" in df_cat.columns else df_cat["ppto_ytd"].sum()
            alc_tot = df_cat["venta_2026_ytd"].sum() / pyt_tot * 100 if pyt_tot > 0 else 0
            rows_data.append([
                "── TOTAL",
                _fmt(df_cat["ppto_anual"].sum()),
                _fmt(df_cat["ppto_incr_anual"].sum() if "ppto_incr_anual" in df_cat.columns else 0),
                _fmt(df_cat["ppto_nuevos_anual"].sum() if "ppto_nuevos_anual" in df_cat.columns else 0),
                _fmt(pt_tot),
                _fmt(df_cat["venta_2026_ytd"].sum()),
                f"{alc_tot:.1f}%",
            ])

            for ri, row in enumerate(rows_data):
                y = t_top + row_h + ri * row_h
                is_total = (ri == len(rows_data) - 1)
                bg = RGBColor(0xE8, 0xED, 0xF7) if is_total else (GRIS if ri % 2 == 0 else BLANCO)
                x = t_left
                for ci, (val, w) in enumerate(zip(row, col_ws)):
                    add_rect(sl1, x, y, w - 0.02, row_h, bg)
                    add_textbox(sl1, x + 0.03, y + 0.03, w - 0.06, row_h - 0.04,
                                val, bold=is_total, fontsize=7,
                                color=AZUL if is_total else RGBColor(0x33, 0x33, 0x33))
                    x += w

        # Gráfico barras categoría (usando plotly → imagen)
        try:
            import plotly.graph_objects as go
            import plotly.io as pio
            if not df_cat.empty:
                cats = df_cat["categoria"].tolist()
                fig = go.Figure()
                fig.add_bar(name="Venta 2026", x=cats,
                            y=df_cat["venta_2026_ytd"].tolist(),
                            marker_color="#2E75B6")
                pyt_vals = df_cat["total_ppto_ytd"].tolist() if "total_ppto_ytd" in df_cat.columns else df_cat["ppto_ytd"].tolist()
                fig.add_bar(name="PPTO YTD", x=cats, y=pyt_vals,
                            marker_color="#1F3864", opacity=0.6)
                fig.update_layout(
                    barmode="group", margin=dict(l=30, r=10, t=20, b=60),
                    height=260, width=480, paper_bgcolor="white",
                    plot_bgcolor="white", font_size=9,
                    legend=dict(orientation="h", y=-0.25)
                )
                img_bytes = pio.to_image(fig, format="png", scale=2)
                img_stream = io.BytesIO(img_bytes)
                sl1.shapes._spTree  # access tree
                from pptx.util import Inches as In
                sl1.shapes.add_picture(img_stream, In(7.0), In(1.72), In(6.1), In(2.6))
        except Exception as eg:
            print(f"[WARN] No se pudo generar gráfico slide 1: {eg}")

        # ── SLIDE 2: UD_GRUPO + Caída ──────────────────────────────────────────
        sl2 = prs.slides.add_slide(blank_layout)

        add_rect(sl2, 0, 0, 13.33, 0.7, AZUL)
        add_textbox(sl2, 0.15, 0.1, 9, 0.5,
                    f"Variación por UD_GRUPO y Análisis de Caída — YTD {mes_nombre} 2026 vs 2025",
                    bold=True, fontsize=15, color=BLANCO)

        # Tabla UD_GRUPO
        add_textbox(sl2, 0.15, 0.82, 6, 0.3, "Venta por UD_GRUPO: 2026 vs 2025",
                    bold=True, fontsize=10, color=AZUL)

        if not df_ud.empty:
            hdrs2 = ["UD_GRUPO", "Venta 2026", "Venta 2025", "Diferencia", "Var %"]
            cws2  = [2.5, 1.4, 1.4, 1.4, 1.0]
            row_h = 0.26
            t_top2 = 1.15
            t_left2 = 0.15

            x = t_left2
            for h, w in zip(hdrs2, cws2):
                add_rect(sl2, x, t_top2, w - 0.02, row_h, AZUL)
                add_textbox(sl2, x + 0.03, t_top2 + 0.03, w - 0.06, row_h - 0.04,
                            h, bold=True, fontsize=7, color=BLANCO)
                x += w

            ud_rows = []
            for _, r in df_ud.iterrows():
                vp = r.get("var_pct", float("nan"))
                vp_str = f"{vp:+.1f}%" if not (isinstance(vp, float) and math.isnan(vp)) else "—"
                ud_rows.append([
                    str(r["ud_grupo"])[:28],
                    _fmt(r["venta_2026_ytd"]),
                    _fmt(r["venta_2025_ytd"]),
                    _fmt(r["diferencia"]),
                    vp_str,
                ])
            # Total
            tot26 = df_ud["venta_2026_ytd"].sum()
            tot25 = df_ud["venta_2025_ytd"].sum()
            tot_d = tot26 - tot25
            tot_vp = f"{(tot_d/tot25*100):+.1f}%" if tot25 > 0 else "—"
            ud_rows.append(["── TOTAL", _fmt(tot26), _fmt(tot25), _fmt(tot_d), tot_vp])

            for ri, row in enumerate(ud_rows):
                y = t_top2 + row_h + ri * row_h
                is_total = (ri == len(ud_rows) - 1)
                bg = RGBColor(0xE8, 0xED, 0xF7) if is_total else (GRIS if ri % 2 == 0 else BLANCO)
                x = t_left2
                for ci, (val, w) in enumerate(zip(row, cws2)):
                    add_rect(sl2, x, y, w - 0.02, row_h, bg)
                    dif_neg = ci == 3 and val.startswith("-") if ci == 3 else False
                    txt_col = ROJO if dif_neg else (AZUL if is_total else RGBColor(0x33, 0x33, 0x33))
                    add_textbox(sl2, x + 0.03, y + 0.03, w - 0.06, row_h - 0.04,
                                val, bold=is_total, fontsize=7, color=txt_col)
                    x += w

        # Gráfico UD_GRUPO
        try:
            import plotly.graph_objects as go
            import plotly.io as pio
            if not df_ud.empty:
                grupos = df_ud["ud_grupo"].tolist()
                fig2 = go.Figure()
                fig2.add_bar(name="2026", x=grupos,
                             y=df_ud["venta_2026_ytd"].tolist(), marker_color="#2E75B6")
                fig2.add_bar(name="2025", x=grupos,
                             y=df_ud["venta_2025_ytd"].tolist(), marker_color="#A9C4E4", opacity=0.8)
                fig2.update_layout(
                    barmode="group", margin=dict(l=30, r=10, t=20, b=80),
                    height=260, width=480, paper_bgcolor="white",
                    plot_bgcolor="white", font_size=9,
                    legend=dict(orientation="h", y=-0.35),
                    xaxis=dict(tickangle=-35)
                )
                img_bytes2 = pio.to_image(fig2, format="png", scale=2)
                sl2.shapes.add_picture(io.BytesIO(img_bytes2), Inches(7.0), Inches(0.82), Inches(6.1), Inches(2.7))
        except Exception as eg2:
            print(f"[WARN] No se pudo generar gráfico slide 2: {eg2}")

        # Sección caída
        add_textbox(sl2, 0.15, 4.65, 6, 0.3, "Top Clientes en Caída (2026 vs 2025)",
                    bold=True, fontsize=10, color=ROJO)

        if not df_cai.empty:
            hdrs3 = ["Cliente", "Zona", "Venta 2026", "Venta 2025", "Caída"]
            cws3  = [3.2, 1.6, 1.4, 1.4, 1.1]
            row_h3 = 0.25
            t_top3 = 4.97

            x = 0.15
            for h, w in zip(hdrs3, cws3):
                add_rect(sl2, x, t_top3, w - 0.02, row_h3, ROJO)
                add_textbox(sl2, x + 0.03, t_top3 + 0.03, w - 0.06, row_h3 - 0.04,
                            h, bold=True, fontsize=7, color=BLANCO)
                x += w

            top_cai = df_cai.nsmallest(8, "diferencia")
            for ri, (_, r) in enumerate(top_cai.iterrows()):
                y = t_top3 + row_h3 + ri * row_h3
                bg = GRIS if ri % 2 == 0 else BLANCO
                vals = [
                    str(r.get("nombre", r.get("RUT", "")))[:35],
                    str(r.get("zona", ""))[:18],
                    _fmt(r["venta_2026_ytd"]),
                    _fmt(r["venta_2025_ytd"]),
                    _fmt(r["diferencia"]),
                ]
                x = 0.15
                for ci, (val, w) in enumerate(zip(vals, cws3)):
                    add_rect(sl2, x, y, w - 0.02, row_h3, bg)
                    add_textbox(sl2, x + 0.03, y + 0.03, w - 0.06, row_h3 - 0.04,
                                val, fontsize=7, color=ROJO if ci == 4 else RGBColor(0x33, 0x33, 0x33))
                    x += w

        # Guardar a buffer
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    except Exception as e:
        print(f"[ERROR _generate_pptx] {e}")
        import traceback; traceback.print_exc()
        return None


def layout_resumen():
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    df_cat = _df_categoria
    df_ud  = _df_resumen_ud_grupo
    df_cai = _df_clientes_caida

    if df_cat.empty and df_ud.empty:
        return html.Div(style=CARD_STYLE, children=[
            _seccion("📊 Resumen Ejecutivo"),
            _empty_msg("Sin datos. Presione Actualizar con VPN activa."),
        ])

    # ── KPIs globales ────────────────────────────────────────────────────────
    ppto_total  = df_cat["ppto_total"].sum()       if not df_cat.empty and "ppto_total"     in df_cat.columns else 0
    ppto_cod    = df_cat["ppto_anual"].sum()        if not df_cat.empty else 0
    ppto_incr   = df_cat["ppto_incr_anual"].sum()   if not df_cat.empty and "ppto_incr_anual"   in df_cat.columns else 0
    ppto_nuevos = df_cat["ppto_nuevos_anual"].sum() if not df_cat.empty and "ppto_nuevos_anual" in df_cat.columns else 0
    ppto_ytd    = df_cat["total_ppto_ytd"].sum()    if not df_cat.empty and "total_ppto_ytd"    in df_cat.columns else 0
    venta_ytd   = df_cat["venta_2026_ytd"].sum()    if not df_cat.empty else 0
    venta_25    = df_cat["venta_2025_ytd"].sum()    if not df_cat.empty else 0
    alcance     = venta_ytd / ppto_ytd * 100 if ppto_ytd > 0 else 0
    var_25      = (venta_ytd - venta_25) / venta_25 * 100 if venta_25 > 0 else float("nan")
    gap         = venta_ytd - ppto_ytd

    kpi_cards = _kpi_row([
        _kpi_card("PPTO Total 2026",    _fmt(ppto_total),  "#1F3864"),
        _kpi_card("PPTO c/Código",      _fmt(ppto_cod),    "#2E75B6"),
        _kpi_card("PPTO Incremental",   _fmt(ppto_incr),   "#2E75B6"),
        _kpi_card("Prod. Nuevos",       _fmt(ppto_nuevos), "#2E75B6"),
        _kpi_card(f"Venta YTD {mes_nombre}", _fmt(venta_ytd), "#1F3864"),
        _kpi_card("PPTO YTD",          _fmt(ppto_ytd),    "#555"),
        _kpi_card("Alcance %",
                  f"{_sem(alcance)} {alcance:.1f}%",
                  "#1e7e34" if alcance >= 80 else ("#e67e22" if alcance >= 50 else "#c0392b")),
        _kpi_card(f"Gap vs PPTO YTD",  _fmt(gap),
                  "#1e7e34" if gap >= 0 else "#c0392b"),
        _kpi_card("Var vs 2025",
                  f"{var_25:+.1f}%" if not (isinstance(var_25, float) and math.isnan(var_25)) else "—",
                  "#1e7e34" if (not math.isnan(var_25) and var_25 >= 0) else "#c0392b"),
    ])

    # ── Tabla por Categoría ──────────────────────────────────────────────────
    tbl_header = {"backgroundColor": "#1F3864", "color": "white",
                  "fontWeight": "700", "fontSize": "12px", "padding": "6px 8px"}
    tbl_cell   = {"fontSize": "12px", "padding": "5px 8px", "fontFamily": "Segoe UI, sans-serif"}

    if not df_cat.empty:
        cat_rows = []
        for _, r in df_cat.iterrows():
            pt  = r.get("ppto_total", r.get("ppto_anual", 0))
            pyt = r.get("total_ppto_ytd", r.get("ppto_ytd", 0))
            alc = r["venta_2026_ytd"] / pyt * 100 if pyt > 0 else float("nan")
            cat_rows.append({
                "Categoría":     r["categoria"],
                "PPTO c/Cód":   _fmt(r["ppto_anual"]),
                "PPTO Incr.":   _fmt(r.get("ppto_incr_anual", 0)),
                "Prod. Nuevos": _fmt(r.get("ppto_nuevos_anual", 0)),
                "PPTO Total":   _fmt(pt),
                "Venta 2026":   _fmt(r["venta_2026_ytd"]),
                "vs 2025":      _fmt(r.get("venta_2025_ytd", 0)),
                "Alcance %":    f"{_sem(alc)} {alc:.1f}%" if not (isinstance(alc, float) and math.isnan(alc)) else "—",
            })
        # Total
        pt_tot  = df_cat["ppto_total"].sum() if "ppto_total" in df_cat.columns else df_cat["ppto_anual"].sum()
        pyt_tot = df_cat["total_ppto_ytd"].sum() if "total_ppto_ytd" in df_cat.columns else df_cat["ppto_ytd"].sum()
        alc_tot = venta_ytd / pyt_tot * 100 if pyt_tot > 0 else float("nan")
        cat_rows.append({
            "Categoría":     "── TOTAL",
            "PPTO c/Cód":   _fmt(ppto_cod),
            "PPTO Incr.":   _fmt(ppto_incr),
            "Prod. Nuevos": _fmt(ppto_nuevos),
            "PPTO Total":   _fmt(pt_tot),
            "Venta 2026":   _fmt(venta_ytd),
            "vs 2025":      _fmt(venta_25),
            "Alcance %":    f"{_sem(alc_tot)} {alc_tot:.1f}%" if not (isinstance(alc_tot, float) and math.isnan(alc_tot)) else "—",
        })

        tabla_cat = dash_table.DataTable(
            data=cat_rows,
            columns=[{"name": c, "id": c} for c in
                     ["Categoría", "PPTO c/Cód", "PPTO Incr.", "Prod. Nuevos",
                      "PPTO Total", "Venta 2026", "vs 2025", "Alcance %"]],
            style_header=tbl_header,
            style_cell=tbl_cell,
            style_cell_conditional=[
                {"if": {"column_id": "Categoría"}, "textAlign": "left", "fontWeight": "600"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{Categoría} = "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            page_action="none",
        )
    else:
        tabla_cat = _empty_msg("Sin datos de categoría.")

    # ── Gráfico barras por Categoría ─────────────────────────────────────────
    if not df_cat.empty:
        cats = df_cat["categoria"].tolist()
        pyt_vals = df_cat["total_ppto_ytd"].tolist() if "total_ppto_ytd" in df_cat.columns else df_cat["ppto_ytd"].tolist()
        fig_cat = go.Figure()
        fig_cat.add_bar(name=f"Venta YTD {mes_nombre}", x=cats,
                        y=df_cat["venta_2026_ytd"].tolist(), marker_color="#2E75B6")
        fig_cat.add_bar(name="PPTO YTD", x=cats, y=pyt_vals,
                        marker_color="#1F3864", opacity=0.55)
        fig_cat.update_layout(
            barmode="group", height=280, margin=dict(l=40, r=20, t=20, b=40),
            plot_bgcolor="white", paper_bgcolor="white",
            legend=dict(orientation="h", y=-0.2),
            yaxis=dict(tickformat="$,.0f"),
        )
        grafico_cat = dcc.Graph(figure=fig_cat, config={"displayModeBar": False})
    else:
        grafico_cat = _empty_msg()

    # ── Tabla UD_GRUPO ───────────────────────────────────────────────────────
    if not df_ud.empty:
        ud_rows = []
        for _, r in df_ud.iterrows():
            vp = r.get("var_pct", float("nan"))
            ud_rows.append({
                "UD_GRUPO":    r["ud_grupo"],
                "Venta 2026":  _fmt(r["venta_2026_ytd"]),
                "Venta 2025":  _fmt(r["venta_2025_ytd"]),
                "Diferencia":  _fmt(r["diferencia"]),
                "Var %":       f"{vp:+.1f}%" if not (isinstance(vp, float) and math.isnan(vp)) else "—",
            })
        tot26 = df_ud["venta_2026_ytd"].sum()
        tot25 = df_ud["venta_2025_ytd"].sum()
        tot_d = tot26 - tot25
        ud_rows.append({
            "UD_GRUPO":   "── TOTAL",
            "Venta 2026": _fmt(tot26),
            "Venta 2025": _fmt(tot25),
            "Diferencia": _fmt(tot_d),
            "Var %":      f"{(tot_d/tot25*100):+.1f}%" if tot25 > 0 else "—",
        })
        tabla_ud = dash_table.DataTable(
            data=ud_rows,
            columns=[{"name": c, "id": c} for c in
                     ["UD_GRUPO", "Venta 2026", "Venta 2025", "Diferencia", "Var %"]],
            style_header=tbl_header,
            style_cell=tbl_cell,
            style_cell_conditional=[
                {"if": {"column_id": "UD_GRUPO"}, "textAlign": "left", "fontWeight": "600"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{UD_GRUPO} = "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            page_action="none",
        )
    else:
        tabla_ud = _empty_msg("Sin datos de UD_GRUPO.")

    # ── Gráfico UD_GRUPO ─────────────────────────────────────────────────────
    if not df_ud.empty:
        grupos = df_ud["ud_grupo"].tolist()
        fig_ud = go.Figure()
        fig_ud.add_bar(name="2026", x=grupos,
                       y=df_ud["venta_2026_ytd"].tolist(), marker_color="#2E75B6")
        fig_ud.add_bar(name="2025", x=grupos,
                       y=df_ud["venta_2025_ytd"].tolist(), marker_color="#A9C4E4", opacity=0.85)
        fig_ud.update_layout(
            barmode="group", height=300, margin=dict(l=40, r=20, t=20, b=80),
            plot_bgcolor="white", paper_bgcolor="white",
            legend=dict(orientation="h", y=-0.3),
            yaxis=dict(tickformat="$,.0f"),
            xaxis=dict(tickangle=-35),
        )
        grafico_ud = dcc.Graph(figure=fig_ud, config={"displayModeBar": False})
    else:
        grafico_ud = _empty_msg()

    # ── Sección Caída ────────────────────────────────────────────────────────
    if not df_cai.empty:
        total_caida = df_cai["diferencia"].sum()
        n_caida     = len(df_cai)
        cai_rows = []
        for _, r in df_cai.head(20).iterrows():
            cp = r.get("caida_pct", float("nan"))
            cai_rows.append({
                "Cliente":    str(r.get("nombre", r.get("RUT", "")))[:45],
                "Zona":       r.get("zona", ""),
                "Venta 2026": _fmt(r["venta_2026_ytd"]),
                "Venta 2025": _fmt(r["venta_2025_ytd"]),
                "Caída $":    _fmt(r["diferencia"]),
                "Caída %":    f"{cp:.1f}%" if not (isinstance(cp, float) and math.isnan(cp)) else "—",
            })
        tabla_caida = dash_table.DataTable(
            data=cai_rows,
            columns=[{"name": c, "id": c} for c in
                     ["Cliente", "Zona", "Venta 2026", "Venta 2025", "Caída $", "Caída %"]],
            style_header={**tbl_header, "backgroundColor": "#c0392b"},
            style_cell=tbl_cell,
            style_cell_conditional=[
                {"if": {"column_id": "Cliente"}, "textAlign": "left"},
                {"if": {"column_id": "Zona"},    "textAlign": "left"},
            ],
            style_data_conditional=[
                {"if": {"column_id": "Caída $"}, "color": "#c0392b", "fontWeight": "600"},
            ],
            page_action="none",
            style_table={"maxHeight": "320px", "overflowY": "auto"},
        )
        caida_kpis = _kpi_row([
            _kpi_card("Clientes en caída", str(n_caida), "#c0392b"),
            _kpi_card("Total caída $", _fmt(total_caida), "#c0392b"),
        ])
        seccion_caida = html.Div([caida_kpis, tabla_caida])
    else:
        seccion_caida = _empty_msg("Sin datos de caída.")

    return html.Div([
        # KPIs
        html.Div(style=CARD_STYLE, children=[
            html.Div(style={"display": "flex", "justifyContent": "space-between",
                            "alignItems": "center", "marginBottom": "12px"}, children=[
                _seccion(f"📊 Resumen Ejecutivo — YTD {mes_nombre} 2026"),
                html.A(
                    html.Button("📥 Descargar PPTX",
                                id="resumen-btn-pptx",
                                style={"background": "#1F3864", "color": "white",
                                       "border": "none", "borderRadius": "6px",
                                       "padding": "8px 16px", "cursor": "pointer",
                                       "fontSize": "12px", "fontWeight": "600"}),
                    id="resumen-pptx-link",
                    download=f"Resumen_LBF_{mes_nombre}2026.pptx",
                    href="",
                    style={"textDecoration": "none"},
                ),
            ]),
            kpi_cards,
        ]),

        # Por Categoría
        html.Div(style=CARD_STYLE, children=[
            _seccion("🗂️ PPTO vs Venta por Categoría"),
            html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr", "gap": "20px"}, children=[
                html.Div(tabla_cat),
                html.Div(grafico_cat),
            ]),
        ]),

        # Por UD_GRUPO
        html.Div(style=CARD_STYLE, children=[
            _seccion("🏭 Venta 2026 vs 2025 por UD_GRUPO"),
            html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr", "gap": "20px"}, children=[
                html.Div(tabla_ud),
                html.Div(grafico_ud),
            ]),
        ]),

        # Caída
        html.Div(style=CARD_STYLE, children=[
            _seccion("📉 Clientes en Caída (Top 20)", "#c0392b"),
            seccion_caida,
        ]),

        dcc.Download(id="resumen-download-pptx"),
    ])


# ─── Main app layout ───────────────────────────────────────────────────────────

def _main_layout():
    cache_label = f"Datos: {_last_update}" if _last_update != "—" else ""
    bd_label    = "🔒 Modo caché" if not CONN_STR else ""
    return html.Div(style={"fontFamily": "Segoe UI, Roboto, sans-serif",
                            "backgroundColor": "#f4f6fa", "minHeight": "100vh"}, children=[
        # Navbar
        html.Div(style={
            "background": "linear-gradient(90deg, #0d1f3c, #1F3864)",
            "padding": "12px 24px",
            "display": "flex", "alignItems": "center", "gap": "16px",
            "boxShadow": "0 2px 12px rgba(0,0,0,0.3)",
        }, children=[
            html.Div(style={"display": "flex", "alignItems": "center", "gap": "10px",
                            "flex": "1"}, children=[
                html.Span("◈", style={"color": "#4da3e8", "fontSize": "16px"}),
                html.H1("Análisis Presupuesto 2026",
                        style={"color": "white", "margin": "0",
                               "fontSize": "17px", "fontWeight": "700",
                               "letterSpacing": "-0.3px"}),
            ]),
            html.Button("🔄 Actualizar", id="btn-refresh",
                        style={"background": "#2E75B6", "color": "white", "border": "none",
                               "borderRadius": "6px", "padding": "6px 14px",
                               "cursor": "pointer", "fontSize": "12px", "fontWeight": "600"}),
            html.Div(id="last-update", children=cache_label,
                     style={"color": "#aec6e8", "fontSize": "11px"}),
            html.Div(bd_label, style={"color": "#f39c12", "fontSize": "11px"}),
        ]),
        # Tabs
        dcc.Tabs(id="main-tabs", value="tab-categoria",
                 style={"margin": "0", "overflowX": "auto"},
                 colors={"border": "#e0e0e0", "primary": "#1F3864", "background": "#f4f6fa"},
                 children=[
            dcc.Tab(label="🗂️ Categoría", value="tab-categoria",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
            dcc.Tab(label="🏢 Zona", value="tab-zona",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
            dcc.Tab(label="🎯 Desalineación", value="tab-desalineacion",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
            dcc.Tab(label="💲 Precios", value="tab-precios",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
            dcc.Tab(label="📊 Resumen", value="tab-resumen",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
            dcc.Tab(label="🔬 PPTO/Venta", value="tab-ppto-vs-venta",
                    style={"padding": "6px 8px", "fontSize": "11px"},
                    selected_style={"padding": "6px 8px", "fontSize": "11px",
                                    "fontWeight": "700", "color": "#1F3864"}),
        ]),
        html.Div(id="tab-content", style={"padding": "16px"}),
        dcc.Store(id="refresh-store"),
    ])


# ─── Root layout ───────────────────────────────────────────────────────────────
# Cargar caché antes de construir el layout para que los globals tengan datos
_load_from_cache()

app.layout = _main_layout()


# ─── Callbacks ─────────────────────────────────────────────────────────────────


@app.callback(
    Output("resumen-download-pptx", "data"),
    Input("resumen-btn-pptx", "n_clicks"),
    prevent_initial_call=True,
)
def download_pptx(n):
    if not n:
        raise PreventUpdate
    import base64
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    content = _generate_pptx()
    if content is None:
        raise PreventUpdate
    b64 = base64.b64encode(content).decode()
    return {
        "content": b64,
        "filename": f"Resumen_LBF_{mes_nombre}{_ANO_ACT}.pptx",
        "base64": True,
        "type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }


@app.callback(
    Output("refresh-store", "data"),
    Output("last-update", "children"),
    Output("btn-refresh", "disabled"),
    Input("btn-refresh", "n_clicks"),
    prevent_initial_call=True,
)
def do_refresh(n):
    if not n:
        raise PreventUpdate
    if not CONN_STR:
        return n, "⚠️ Sin acceso a BD (requiere VPN + db_config.py)", False
    try:
        _reload_all_data()
        return n, f"✅ Actualizado: {_last_update}", False
    except Exception as e:
        return n, f"⚠️ Error al actualizar: {e}", False


@app.callback(
    Output("tab-content", "children"),
    Input("main-tabs", "value"),
    Input("refresh-store", "data"),
)
def render_tab(tab, _refresh):
    if tab == "tab-categoria":
        return layout_categoria()
    elif tab == "tab-zona":
        return layout_zona()
    elif tab == "tab-desalineacion":
        return layout_desalineacion()
    elif tab == "tab-precios":
        return layout_precios()
    elif tab == "tab-resumen":
        return layout_resumen()
    elif tab == "tab-ppto-vs-venta":
        return layout_ppto_vs_venta()
    return _empty_msg("Tab no encontrado.")




@app.callback(
    Output("cat-detalle-container", "children"),
    Input("cat-tabla", "selected_rows"),
    State("cat-tabla", "data"),
    prevent_initial_call=True,
)
def update_cat_detalle(selected_rows, data):
    if not selected_rows or not data:
        raise PreventUpdate
    row_idx = selected_rows[0]
    if row_idx >= len(data):
        raise PreventUpdate
    cat = data[row_idx].get("cat", "")
    if not cat or cat == "── TOTAL":
        raise PreventUpdate
    return _build_cat_detalle(cat)


@app.callback(
    Output("cat-det-cli", "options"),
    Output("cat-det-cli", "value"),
    Input("cat-det-zona", "value"),
    State("cat-det-store", "data"),
    prevent_initial_call=True,
)
def update_cat_det_cli_opts(zona, store_data):
    if not store_data:
        raise PreventUpdate
    df = pd.DataFrame(store_data)
    if zona and zona != "TODAS":
        df = df[df["zona"] == zona]
    df_ppto = df[df["origen"] == "PPTO"] if "origen" in df.columns else df
    clientes = df_ppto[["rut", "nombre_cliente"]].drop_duplicates().sort_values("nombre_cliente")
    opts = [{"label": "Todos los clientes", "value": "TODOS"}] + \
           [{"label": f"{r['nombre_cliente']} ({r['rut']})", "value": r["rut"]}
            for _, r in clientes.iterrows() if r["rut"]]
    return opts, "TODOS"


@app.callback(
    Output("cat-det-tabla", "children"),
    Input("cat-det-zona", "value"),
    Input("cat-det-cli", "value"),
    Input("cat-det-buscar", "value"),
    Input("cat-det-store", "data"),
)
def update_cat_det_tabla(zona, cli, buscar, store_data):
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    if not store_data:
        return _empty_msg()
    df = pd.DataFrame(store_data)
    for c in ["ppto_anual", "ppto_ytd", "cant_ppto", "venta_2026_ytd", "es_incremental"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)

    # Separar incrementales ANTES de aplicar filtros
    df_incr = df[df["es_incremental"] == 1].copy()
    df = df[df["es_incremental"] == 0].copy()

    # Guardar copia sin filtros para el resumen por zona
    df_all_zonas   = df.copy()
    df_incr_all    = df_incr.copy()

    if zona and zona != "TODAS":
        df = df[df["zona"] == zona]
        df_incr = df_incr[df_incr["zona"] == zona]
    if cli and cli != "TODOS":
        df = df[df["rut"] == cli]
    if buscar and buscar.strip():
        term = buscar.strip().lower()
        mask = (df["CODIGO"].str.lower().str.contains(term, na=False) |
                df["descripcion"].str.lower().str.contains(term, na=False))
        df = df[mask]
    if df.empty and df_incr.empty:
        return _empty_msg("Sin datos para los filtros seleccionados.")

    tbl_style_header = {"backgroundColor": "#2E75B6", "color": "white",
                        "fontWeight": "bold", "fontSize": "12px"}
    tbl_style_cell   = {"fontSize": "12px", "padding": "6px 10px",
                        "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"}

    # ── Tabla principal ────────────────────────────────────────────────────────
    table_data = []
    for _, r in df.iterrows():
        sin_ppto = r.get("origen", "PPTO") == "SIN_PPTO"
        cumpl  = (r["venta_2026_ytd"] / r["ppto_ytd"] * 100) if r["ppto_ytd"] > 0 else float("nan")
        proyec = (r["venta_2026_ytd"] / _MES_ACT * 12) if _MES_ACT > 0 else 0
        gap    = r["venta_2026_ytd"] - r["ppto_ytd"]
        table_data.append({
            "zona":        r["zona"],
            "cliente":     r.get("nombre_cliente", ""),
            "codigo":      r["CODIGO"],
            "descripcion": r["descripcion"],
            "sem":         _sem(cumpl) if not sin_ppto else "⚪",
            "PPTO Anual":  _fmt_abs(r["ppto_anual"]),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(r["ppto_ytd"]),
            "Venta 2026 YTD": _fmt_abs(r["venta_2026_ytd"]),
            "Cumpl %":     _fmt_pct(cumpl),
            "Proyec. Año": _fmt_abs(proyec),
            "Gap $":       _fmt_abs(gap),
        })

    tot_ppto_anual = df["ppto_anual"].sum()
    tot_ppto_ytd   = df["ppto_ytd"].sum()
    tot_venta      = df["venta_2026_ytd"].sum()
    tot_cumpl      = (tot_venta / tot_ppto_ytd * 100) if tot_ppto_ytd > 0 else float("nan")
    tot_proyec     = (tot_venta / _MES_ACT * 12) if _MES_ACT > 0 else 0
    tot_gap        = tot_venta - tot_ppto_ytd
    table_data.append({
        "zona": "", "cliente": "", "codigo": "", "descripcion": "── TOTAL", "sem": "",
        "PPTO Anual":     _fmt_abs(tot_ppto_anual),
        f"PPTO YTD ({mes_nombre})": _fmt_abs(tot_ppto_ytd),
        "Venta 2026 YTD": _fmt_abs(tot_venta),
        "Cumpl %":        _fmt_pct(tot_cumpl),
        "Proyec. Año":    _fmt_abs(tot_proyec),
        "Gap $":          _fmt_abs(tot_gap),
    })

    display_cols = ["zona", "cliente", "codigo", "descripcion", "sem",
                    "PPTO Anual", f"PPTO YTD ({mes_nombre})",
                    "Venta 2026 YTD", "Cumpl %", "Proyec. Año", "Gap $"]

    tabla_principal = dash_table.DataTable(
        data=table_data,
        columns=[{"name": c, "id": c} for c in display_cols],
        style_table={"overflowX": "auto"},
        style_header=tbl_style_header,
        style_cell=tbl_style_cell,
        style_cell_conditional=[
            {"if": {"column_id": "zona"}, "textAlign": "left"},
            {"if": {"column_id": "cliente"}, "textAlign": "left", "maxWidth": "180px",
             "overflow": "hidden", "textOverflow": "ellipsis"},
            {"if": {"column_id": "codigo"}, "textAlign": "left"},
            {"if": {"column_id": "descripcion"}, "textAlign": "left", "maxWidth": "240px",
             "overflow": "hidden", "textOverflow": "ellipsis"},
            {"if": {"column_id": "sem"}, "textAlign": "center", "width": "40px"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{descripcion} = "── TOTAL"'},
             "fontWeight": "700", "backgroundColor": "#ddeeff"},
        ],
        page_size=100,
        sort_action="native",
    )

    # ── Sección: PPTO con 0 ventas — agrupado por zona ────────────────────────
    df_ppto = df[df["origen"] == "PPTO"] if "origen" in df.columns else df
    df_sv = df_ppto[df_ppto["venta_2026_ytd"] == 0].copy()
    sv_ppto_anual = df_sv["ppto_anual"].sum()
    sv_ppto_ytd   = df_sv["ppto_ytd"].sum()
    pct_del_gap   = abs(sv_ppto_ytd / tot_gap * 100) if tot_gap != 0 else float("nan")

    if df_sv.empty:
        sin_venta_section = html.Div(
            "✅ Todos los productos presupuestados tienen al menos una venta en 2026.",
            style={"color": "#27ae60", "fontSize": "13px", "padding": "12px 0", "fontWeight": "600"}
        )
    else:
        # Tabla agrupada por zona
        sv_zona_rows = []
        for zona_i, g in df_sv.groupby("zona", sort=True):
            sv_zona_rows.append({
                "_zona": zona_i,
                "Zona": zona_i,
                "# Productos": len(g),
                "PPTO Anual": _fmt_abs(g["ppto_anual"].sum()),
                f"PPTO YTD ({mes_nombre})": _fmt_abs(g["ppto_ytd"].sum()),
            })
        sv_zona_rows.append({
            "_zona": "",
            "Zona": f"── TOTAL ({len(df_sv)} productos)",
            "# Productos": len(df_sv),
            "PPTO Anual": _fmt_abs(sv_ppto_anual),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(sv_ppto_ytd),
        })
        sv_zona_cols = ["Zona", "# Productos", "PPTO Anual", f"PPTO YTD ({mes_nombre})"]
        # Guardar detalle por zona como store
        sv_detail_store = {zona_i: [
            {"cliente": r.get("nombre_cliente",""), "codigo": r["CODIGO"],
             "descripcion": r["descripcion"],
             "PPTO Anual": _fmt_abs(r["ppto_anual"]),
             f"PPTO YTD ({mes_nombre})": _fmt_abs(r["ppto_ytd"])}
            for _, r in g.iterrows()
        ] for zona_i, g in df_sv.groupby("zona", sort=True)}

        sin_venta_section = html.Div(children=[
            _kpi_row([
                _kpi_card("Prods sin venta 2026", str(len(df_sv)), "#c0392b"),
                _kpi_card("PPTO Anual sin vender", _fmt(sv_ppto_anual), "#c0392b"),
                _kpi_card("PPTO YTD sin vender", _fmt(sv_ppto_ytd), "#c0392b"),
                _kpi_card("% del Gap YTD", _fmt_pct(pct_del_gap),
                          "#c0392b" if not math.isnan(pct_del_gap) else "#888"),
            ]),
            html.P("💡 Haz clic en una zona para ver los productos sin venta.",
                   style={"fontSize": "11px", "color": "#666", "marginBottom": "6px"}),
            dcc.Store(id="sv-detail-store", data=sv_detail_store),
            dash_table.DataTable(
                id="sv-zona-tabla",
                data=[{k: v for k, v in r.items() if k != "_zona"} for r in sv_zona_rows],
                columns=[{"name": c, "id": c} for c in sv_zona_cols],
                style_table={"overflowX": "auto"},
                style_header={"backgroundColor": "#922b21", "color": "white",
                              "fontWeight": "bold", "fontSize": "12px"},
                style_cell=tbl_style_cell,
                style_cell_conditional=[
                    {"if": {"column_id": "Zona"}, "textAlign": "left"},
                    {"if": {"column_id": "# Productos"}, "textAlign": "center"},
                ],
                style_data_conditional=[
                    {"if": {"filter_query": '{Zona} contains "── TOTAL"'},
                     "fontWeight": "700", "backgroundColor": "#fadbd8"},
                ],
                page_action="none", sort_action="native", row_selectable="single",
            ),
            html.Div(id="sv-zona-detalle"),
        ])

    # ── Sección: PPTO incremental agrupado por zona ───────────────────────────
    if df_incr.empty:
        incr_section = html.Div()
    else:
        incr_zona_rows = []
        incr_detail_store = {}
        for zona_i, g in df_incr.groupby("zona", sort=True):
            incr_zona_rows.append({
                "Zona": zona_i,
                "# Líneas": len(g),
                "PPTO Anual": _fmt_abs(g["ppto_anual"].sum()),
                f"PPTO YTD ({mes_nombre})": _fmt_abs(g["ppto_ytd"].sum()),
            })
            incr_detail_store[zona_i] = [
                {"Descripción": r["descripcion"],
                 "Cliente": r.get("nombre_cliente","") or "—",
                 "PPTO Anual": _fmt_abs(r["ppto_anual"]),
                 f"PPTO YTD ({mes_nombre})": _fmt_abs(r["ppto_ytd"])}
                for _, r in g.iterrows()
            ]
        incr_zona_rows.append({
            "Zona": "── TOTAL",
            "# Líneas": len(df_incr),
            "PPTO Anual": _fmt_abs(df_incr["ppto_anual"].sum()),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(df_incr["ppto_ytd"].sum()),
        })
        incr_zona_cols = ["Zona", "# Líneas", "PPTO Anual", f"PPTO YTD ({mes_nombre})"]
        incr_section = html.Div(children=[
            _kpi_row([
                _kpi_card("PPTO Incremental Anual", _fmt(df_incr["ppto_anual"].sum()), "#e67e22"),
                _kpi_card(f"PPTO Incremental YTD ({mes_nombre})", _fmt(df_incr["ppto_ytd"].sum()), "#e67e22"),
                _kpi_card("# Líneas", str(len(df_incr)), "#7b2d8b",
                          "Sin código — no comparables con venta real"),
            ]),
            html.P("💡 Haz clic en una zona para ver el detalle.",
                   style={"fontSize": "11px", "color": "#666", "marginBottom": "6px"}),
            dcc.Store(id="incr-detail-store", data=incr_detail_store),
            dash_table.DataTable(
                id="incr-zona-tabla",
                data=incr_zona_rows,
                columns=[{"name": c, "id": c} for c in incr_zona_cols],
                style_table={"overflowX": "auto"},
                style_header={"backgroundColor": "#e67e22", "color": "white",
                              "fontWeight": "bold", "fontSize": "12px"},
                style_cell=tbl_style_cell,
                style_cell_conditional=[
                    {"if": {"column_id": "Zona"}, "textAlign": "left"},
                    {"if": {"column_id": "# Líneas"}, "textAlign": "center"},
                ],
                style_data_conditional=[
                    {"if": {"filter_query": '{Zona} = "── TOTAL"'},
                     "fontWeight": "700", "backgroundColor": "#e8edf7"},
                ],
                page_action="none", sort_action="native", row_selectable="single",
            ),
            html.Div(id="incr-zona-detalle"),
        ])

    # ── Resumen por Zona ──────────────────────────────────────────────────────────
    zonas_all = sorted(set(df_all_zonas["zona"].dropna().unique().tolist()) |
                       set(df_incr_all["zona"].dropna().unique().tolist() if not df_incr_all.empty else []))
    zona_resumen_rows = []
    for zona_i in zonas_all:
        df_z    = df_all_zonas[df_all_zonas["zona"] == zona_i]
        dincr_z = df_incr_all[df_incr_all["zona"] == zona_i] if not df_incr_all.empty else pd.DataFrame()
        ppto_cod = df_z[df_z["origen"] == "PPTO"]["ppto_anual"].sum() if "origen" in df_z.columns else df_z["ppto_anual"].sum()
        ppto_ytd_cod = df_z[df_z["origen"] == "PPTO"]["ppto_ytd"].sum() if "origen" in df_z.columns else df_z["ppto_ytd"].sum()
        if not dincr_z.empty:
            mask_nv       = dincr_z["descripcion"].str.upper().str.strip() == "PRODUCTOS NUEVOS"
            ppto_incr_z   = dincr_z[~mask_nv]["ppto_anual"].sum()
            ppto_nuevos_z = dincr_z[mask_nv]["ppto_anual"].sum()
            ppto_incr_ytd_z   = dincr_z[~mask_nv]["ppto_ytd"].sum()
            ppto_nuevos_ytd_z = dincr_z[mask_nv]["ppto_ytd"].sum()
        else:
            ppto_incr_z = ppto_nuevos_z = ppto_incr_ytd_z = ppto_nuevos_ytd_z = 0
        ppto_total_z   = ppto_cod + ppto_incr_z + ppto_nuevos_z
        total_ytd_z    = ppto_ytd_cod + ppto_incr_ytd_z + ppto_nuevos_ytd_z
        venta_z        = df_z["venta_2026_ytd"].sum()
        alcance_z      = (venta_z / total_ytd_z * 100) if total_ytd_z > 0 else float("nan")
        zona_resumen_rows.append({
            "Zona": zona_i,
            "PPTO 2026": _fmt_abs(ppto_total_z),
            "PPTO c/Cód": _fmt_abs(ppto_cod),
            "PPTO Incr.": _fmt_abs(ppto_incr_z),
            "PPTO Prod. Nuevos": _fmt_abs(ppto_nuevos_z),
            f"PPTO YTD ({mes_nombre})": _fmt_abs(total_ytd_z),
            "Venta+Guías YTD": _fmt_abs(venta_z),
            "Alcance %": _fmt_pct(alcance_z),
        })
    if not df_incr_all.empty:
        mask_nv_all       = df_incr_all["descripcion"].str.upper().str.strip() == "PRODUCTOS NUEVOS"
        tot_incr_z        = df_incr_all[~mask_nv_all]["ppto_anual"].sum()
        tot_nuevos_z      = df_incr_all[mask_nv_all]["ppto_anual"].sum()
        tot_incr_ytd      = df_incr_all[~mask_nv_all]["ppto_ytd"].sum()
        tot_nuevos_ytd    = df_incr_all[mask_nv_all]["ppto_ytd"].sum()
    else:
        tot_incr_z = tot_nuevos_z = tot_incr_ytd = tot_nuevos_ytd = 0
    tot_cod_z     = df_all_zonas[df_all_zonas["origen"] == "PPTO"]["ppto_anual"].sum() if "origen" in df_all_zonas.columns else df_all_zonas["ppto_anual"].sum()
    tot_ytd_cod   = df_all_zonas[df_all_zonas["origen"] == "PPTO"]["ppto_ytd"].sum() if "origen" in df_all_zonas.columns else df_all_zonas["ppto_ytd"].sum()
    tot_total_z   = tot_cod_z + tot_incr_z + tot_nuevos_z
    tot_total_ytd = tot_ytd_cod + tot_incr_ytd + tot_nuevos_ytd
    tot_venta_z   = df_all_zonas["venta_2026_ytd"].sum()
    tot_alcance_z = (tot_venta_z / tot_total_ytd * 100) if tot_total_ytd > 0 else float("nan")
    zona_resumen_rows.append({
        "Zona": "── TOTAL",
        "PPTO 2026": _fmt_abs(tot_total_z),
        "PPTO c/Cód": _fmt_abs(tot_cod_z),
        "PPTO Incr.": _fmt_abs(tot_incr_z),
        "PPTO Prod. Nuevos": _fmt_abs(tot_nuevos_z),
        f"PPTO YTD ({mes_nombre})": _fmt_abs(tot_total_ytd),
        "Venta+Guías YTD": _fmt_abs(tot_venta_z),
        "Alcance %": _fmt_pct(tot_alcance_z),
    })
    zona_resumen_cols = ["Zona", "PPTO 2026", "PPTO c/Cód", "PPTO Incr.", "PPTO Prod. Nuevos",
                         f"PPTO YTD ({mes_nombre})", "Venta+Guías YTD", "Alcance %"]
    zona_resumen_section = dash_table.DataTable(
        data=zona_resumen_rows,
        columns=[{"name": c, "id": c} for c in zona_resumen_cols],
        style_table={"overflowX": "auto"},
        style_header={"backgroundColor": "#1F3864", "color": "white", "fontWeight": "bold", "fontSize": "12px"},
        style_cell={"fontSize": "12px", "padding": "6px 10px", "textAlign": "right", "fontFamily": "Segoe UI, sans-serif"},
        style_cell_conditional=[{"if": {"column_id": "Zona"}, "textAlign": "left"}],
        style_data_conditional=[{"if": {"filter_query": '{Zona} = "── TOTAL"'}, "fontWeight": "700", "backgroundColor": "#e8edf7"}],
        page_action="none",
        sort_action="native",
    )

    return html.Div(children=[
        html.H4("📊 Resumen por Zona",
                style={"color": "#1F3864", "fontSize": "13px", "fontWeight": "700", "marginBottom": "8px"}),
        html.P("PPTO c/Cód: presupuesto con código de producto. PPTO Incr.: presupuesto sin código (otros). "
               "PPTO Prod. Nuevos: solo desde Mayo 2026.",
               style={"fontSize": "11px", "color": "#888", "marginBottom": "8px"}),
        zona_resumen_section,
        html.Hr(style={"margin": "20px 0", "borderColor": "#ddd"}),
        tabla_principal,
        html.Hr(style={"margin": "20px 0", "borderColor": "#ddd"}),
        html.H4("🔴 Productos con Presupuesto y Sin Venta 2026",
                style={"color": "#922b21", "fontSize": "13px", "fontWeight": "700",
                       "marginBottom": "12px"}),
        sin_venta_section,
        html.Hr(style={"margin": "20px 0", "borderColor": "#ddd"}),
        html.H4("⚠️ Presupuesto Incremental — Sin código de producto",
                style={"color": "#e67e22", "fontSize": "13px", "fontWeight": "700",
                       "marginBottom": "12px"}),
        html.P("Estas filas no tienen código de producto en el PPTO 2026. "
               "No se comparan con venta real — se muestran como referencia.",
               style={"fontSize": "11px", "color": "#888", "marginBottom": "8px"}),
        incr_section,
    ])


@app.callback(
    Output("sv-zona-detalle", "children"),
    Input("sv-zona-tabla", "selected_rows"),
    State("sv-zona-tabla", "data"),
    State("sv-detail-store", "data"),
    prevent_initial_call=True,
)
def sv_zona_drill(sel, tbl_data, detail_store):
    if not sel or not detail_store:
        raise PreventUpdate
    zona = tbl_data[sel[0]]["Zona"]
    if "── TOTAL" in zona:
        raise PreventUpdate
    rows = detail_store.get(zona, [])
    if not rows:
        return _empty_msg(f"Sin detalle para {zona}")
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    cols = ["cliente", "codigo", "descripcion", "PPTO Anual", f"PPTO YTD ({mes_nombre})"]
    return html.Div(style={"marginTop": "12px", "padding": "12px",
                           "backgroundColor": "#fff8f8", "borderRadius": "6px",
                           "border": "1px solid #e74c3c"}, children=[
        html.H5(f"Detalle sin venta — {zona}",
                style={"color": "#922b21", "fontSize": "12px", "fontWeight": "700", "marginBottom": "8px"}),
        dash_table.DataTable(
            data=rows, columns=[{"name": c, "id": c} for c in cols],
            style_table={"overflowX": "auto"},
            style_header={"backgroundColor": "#922b21", "color": "white", "fontWeight": "bold", "fontSize": "12px"},
            style_cell={"fontSize": "12px", "padding": "5px 8px", "textAlign": "right",
                        "fontFamily": "Segoe UI, sans-serif"},
            style_cell_conditional=[
                {"if": {"column_id": "cliente"},     "textAlign": "left", "maxWidth": "160px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
                {"if": {"column_id": "codigo"},      "textAlign": "left"},
                {"if": {"column_id": "descripcion"}, "textAlign": "left", "maxWidth": "240px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
            ],
            page_action="none", sort_action="native",
        ),
    ])


@app.callback(
    Output("incr-zona-detalle", "children"),
    Input("incr-zona-tabla", "selected_rows"),
    State("incr-zona-tabla", "data"),
    State("incr-detail-store", "data"),
    prevent_initial_call=True,
)
def incr_zona_drill(sel, tbl_data, detail_store):
    if not sel or not detail_store:
        raise PreventUpdate
    zona = tbl_data[sel[0]]["Zona"]
    if "── TOTAL" in zona:
        raise PreventUpdate
    rows = detail_store.get(zona, [])
    if not rows:
        return _empty_msg(f"Sin detalle para {zona}")
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    cols = ["Descripción", "Cliente", "PPTO Anual", f"PPTO YTD ({mes_nombre})"]
    return html.Div(style={"marginTop": "12px", "padding": "12px",
                           "backgroundColor": "#fff8ee", "borderRadius": "6px",
                           "border": "1px solid #e67e22"}, children=[
        html.H5(f"Detalle incremental — {zona}",
                style={"color": "#e67e22", "fontSize": "12px", "fontWeight": "700", "marginBottom": "8px"}),
        dash_table.DataTable(
            data=rows, columns=[{"name": c, "id": c} for c in cols],
            style_table={"overflowX": "auto"},
            style_header={"backgroundColor": "#e67e22", "color": "white", "fontWeight": "bold", "fontSize": "12px"},
            style_cell={"fontSize": "12px", "padding": "5px 8px", "textAlign": "right",
                        "fontFamily": "Segoe UI, sans-serif"},
            style_cell_conditional=[
                {"if": {"column_id": "Descripción"}, "textAlign": "left", "maxWidth": "260px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
                {"if": {"column_id": "Cliente"}, "textAlign": "left"},
            ],
            page_action="none", sort_action="native",
        ),
    ])


@app.callback(
    Output("cat-top25-caida", "children"),
    Output("cat-top25-alza",  "children"),
    Input("cat-det-categoria", "data"),
    prevent_initial_call=True,
)
def update_cat_top25(categoria):
    if not categoria:
        raise PreventUpdate
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    try:
        conn = get_conn()
        cur = conn.cursor()
        cat_safe = categoria.replace("'","''")
        cat_in   = _cat_in(cat_safe)
        cur.execute(f"""
            WITH
            v26 AS (
                SELECT RUT, MAX(LEFT(ISNULL(NOMBRE,RUT),50)) AS nombre,
                       SUM(CASE WHEN MES<={_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2026
                FROM DW_TOTAL_FACTURA
                WHERE ANO={_ANO_ACT} AND {_DW_FILTRO}
                  AND ISNULL(CATEGORIA,'') {cat_in}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                GROUP BY RUT
            ),
            v25 AS (
                SELECT RUT,
                       SUM(CASE WHEN MES<={_MES_ACT} THEN CAST(VENTA AS float) ELSE 0 END) AS venta_2025
                FROM DW_TOTAL_FACTURA
                WHERE ANO=2025 AND {_DW_FILTRO}
                  AND ISNULL(CATEGORIA,'') {cat_in}
                  AND ISNULL(RUT,'') NOT IN ('','0')
                GROUP BY RUT
            )
            SELECT a.RUT, a.nombre, a.venta_2026,
                   ISNULL(b.venta_2025,0) AS venta_2025,
                   a.venta_2026 - ISNULL(b.venta_2025,0) AS diferencia
            FROM v26 a LEFT JOIN v25 b ON b.RUT=a.RUT
            WHERE ISNULL(b.venta_2025,0) > 0
        """)
        rows = cur.fetchall()
        cols = [d[0].strip() for d in cur.description]
        conn.close()
        df = pd.DataFrame.from_records(rows, columns=cols)
        for c in ["venta_2026","venta_2025","diferencia"]:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
        df["dif_pct"] = df.apply(
            lambda r: r["diferencia"]/r["venta_2025"]*100 if r["venta_2025"]>0 else float("nan"), axis=1)
    except Exception as e:
        print(f"[ERROR cat_top25] {e}")
        msg = _empty_msg(f"Error cargando datos: {e}")
        return msg, msg

    if df.empty:
        msg = _empty_msg("Sin datos para esta categoría.")
        return msg, msg

    tbl_style_header = {"fontWeight":"bold","fontSize":"12px","color":"white"}
    tbl_style_cell   = {"fontSize":"12px","padding":"5px 8px","textAlign":"right",
                        "fontFamily":"Segoe UI, sans-serif"}
    cond_left = [{"if":{"column_id":"rut"},"textAlign":"left","maxWidth":"90px"},
                 {"if":{"column_id":"nombre"},"textAlign":"left","maxWidth":"180px",
                  "overflow":"hidden","textOverflow":"ellipsis"}]

    def _tbl(df_sub, tbl_id, hdr_color):
        rows_d = []
        for _, r in df_sub.iterrows():
            rows_d.append({
                "rut": r["RUT"], "nombre": r["nombre"],
                f"Venta {mes_nombre} 2025": _fmt_abs(r["venta_2025"]),
                f"Venta {mes_nombre} 2026": _fmt_abs(r["venta_2026"]),
                "Diferencia $": _fmt_abs(r["diferencia"]),
                "Var %": _fmt_pct(r["dif_pct"]),
            })
        tot_v25 = df_sub["venta_2025"].sum()
        tot_v26 = df_sub["venta_2026"].sum()
        tot_dif = df_sub["diferencia"].sum()
        tot_pct = (tot_dif / tot_v25 * 100) if tot_v25 > 0 else float("nan")
        rows_d.append({
            "rut": "", "nombre": f"── TOTAL ({len(df_sub)} clientes)",
            f"Venta {mes_nombre} 2025": _fmt_abs(tot_v25),
            f"Venta {mes_nombre} 2026": _fmt_abs(tot_v26),
            "Diferencia $": _fmt_abs(tot_dif),
            "Var %": _fmt_pct(tot_pct),
        })
        return dash_table.DataTable(
            id=tbl_id, data=rows_d,
            columns=[{"name":c,"id":c} for c in ["rut","nombre",
                      f"Venta {mes_nombre} 2025",f"Venta {mes_nombre} 2026",
                      "Diferencia $","Var %"]],
            style_table={"overflowX":"auto"},
            style_header={**tbl_style_header,"backgroundColor":hdr_color},
            style_cell=tbl_style_cell,
            style_cell_conditional=cond_left,
            style_data_conditional=[
                {"if": {"filter_query": '{nombre} contains "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            page_action="none", sort_action="native", row_selectable="single",
        )

    return (_tbl(df.sort_values("diferencia").head(25),        "cat-top25-caida-tbl", "#922b21"),
            _tbl(df.sort_values("diferencia",ascending=False).head(25), "cat-top25-alza-tbl",  "#1e8449"))


@app.callback(
    Output("cat-top25-cli-detalle", "children"),
    Input("cat-top25-caida-tbl", "selected_rows"),
    Input("cat-top25-alza-tbl",  "selected_rows"),
    State("cat-top25-caida-tbl", "data"),
    State("cat-top25-alza-tbl",  "data"),
    prevent_initial_call=True,
)
def cat_top25_cli_drill(sel_c, sel_a, data_c, data_a):
    from dash import ctx
    triggered = ctx.triggered_id
    if triggered == "cat-top25-caida-tbl" and sel_c and data_c:
        row = data_c[sel_c[0]]
    elif triggered == "cat-top25-alza-tbl" and sel_a and data_a:
        row = data_a[sel_a[0]]
    else:
        raise PreventUpdate
    rut    = row["rut"]
    nombre = row["nombre"]
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    df = _load_detalle_cliente(rut)
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[_empty_msg(f"Sin detalle para {nombre}")])
    tbl_hdr = {"backgroundColor":"#2E75B6","color":"white","fontWeight":"bold","fontSize":"12px"}
    tbl_cell = {"fontSize":"12px","padding":"6px 10px","textAlign":"right","fontFamily":"Segoe UI, sans-serif"}
    df_cat = df.groupby("categoria",as_index=False).agg(venta_2025=("venta_2025","sum"),venta_2026=("venta_2026","sum"))
    df_cat["diferencia"] = df_cat["venta_2026"] - df_cat["venta_2025"]
    df_cat["dif_pct"] = df_cat.apply(lambda r: r["diferencia"]/r["venta_2025"]*100 if r["venta_2025"]>0 else float("nan"),axis=1)
    tot_v25=df["venta_2025"].sum(); tot_v26=df["venta_2026"].sum(); tot_dif=tot_v26-tot_v25
    tot_pct=(tot_dif/tot_v25*100) if tot_v25>0 else float("nan")
    cat_rows=[{"categoria":r["categoria"],f"Venta {mes_nombre} 2025":_fmt_abs(r["venta_2025"]),
               f"Venta {mes_nombre} 2026":_fmt_abs(r["venta_2026"]),"Diferencia $":_fmt_abs(r["diferencia"]),"Var %":_fmt_pct(r["dif_pct"])}
              for _,r in df_cat.iterrows()]
    cat_rows.append({"categoria":"── TOTAL",f"Venta {mes_nombre} 2025":_fmt_abs(tot_v25),
                     f"Venta {mes_nombre} 2026":_fmt_abs(tot_v26),"Diferencia $":_fmt_abs(tot_dif),"Var %":_fmt_pct(tot_pct)})
    prod_rows=[{"codigo":r["codigo"],"descripcion":r["descripcion"],"categoria":r["categoria"],
                f"Venta {mes_nombre} 2025":_fmt_abs(r["venta_2025"]),f"Venta {mes_nombre} 2026":_fmt_abs(r["venta_2026"]),
                "Diferencia $":_fmt_abs(r["diferencia"]),"Var %":_fmt_pct(r["dif_pct"])}
               for _,r in df.iterrows()]
    prod_rows.append({"codigo":"","descripcion":"── TOTAL","categoria":"",
                      f"Venta {mes_nombre} 2025":_fmt_abs(tot_v25),f"Venta {mes_nombre} 2026":_fmt_abs(tot_v26),
                      "Diferencia $":_fmt_abs(tot_dif),"Var %":_fmt_pct(tot_pct)})
    total_cond=[{"if":{"filter_query":'{descripcion} = "── TOTAL" || {categoria} = "── TOTAL"'},"fontWeight":"700","backgroundColor":"#e8edf7"}]
    cat_cols=["categoria",f"Venta {mes_nombre} 2025",f"Venta {mes_nombre} 2026","Diferencia $","Var %"]
    prod_cols=["codigo","descripcion","categoria",f"Venta {mes_nombre} 2025",f"Venta {mes_nombre} 2026","Diferencia $","Var %"]
    return html.Div(style=CARD_STYLE, children=[
        html.H4(f"📋 Detalle: {nombre} ({rut})",style={"color":"#1F3864","fontSize":"14px","fontWeight":"700","marginBottom":"16px"}),
        html.Div(style={"display":"grid","gridTemplateColumns":"1fr 1fr","gap":"24px"},children=[
            html.Div(children=[
                html.H5("Por Categoría",style={"fontSize":"12px","fontWeight":"700","marginBottom":"8px","color":"#2E75B6"}),
                dash_table.DataTable(data=cat_rows,columns=[{"name":c,"id":c} for c in cat_cols],
                    style_table={"overflowX":"auto"},style_header=tbl_hdr,style_cell=tbl_cell,
                    style_cell_conditional=[{"if":{"column_id":"categoria"},"textAlign":"left"}],
                    style_data_conditional=total_cond,page_action="none",sort_action="native"),
            ]),
            html.Div(children=[
                html.H5("Por Producto",style={"fontSize":"12px","fontWeight":"700","marginBottom":"8px","color":"#2E75B6"}),
                dash_table.DataTable(data=prod_rows,columns=[{"name":c,"id":c} for c in prod_cols],
                    style_table={"overflowX":"auto"},style_header=tbl_hdr,style_cell=tbl_cell,
                    style_cell_conditional=[{"if":{"column_id":"codigo"},"textAlign":"left"},
                                            {"if":{"column_id":"descripcion"},"textAlign":"left","maxWidth":"200px","overflow":"hidden","textOverflow":"ellipsis"},
                                            {"if":{"column_id":"categoria"},"textAlign":"left"}],
                    style_data_conditional=total_cond,page_action="none",sort_action="native"),
            ]),
        ]),
    ])


@app.callback(
    Output("zona-top25-caida", "children"),
    Output("zona-top25-alza", "children"),
    Input("main-tabs", "value"),
    Input("refresh-store", "data"),
)
def update_zona_top25(tab, _):
    if tab != "tab-zona":
        raise PreventUpdate
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    df = _load_top25_clientes()
    if df.empty:
        msg = _empty_msg("Sin datos disponibles. Presiona 🔄 para actualizar.")
        return msg, msg

    tbl_style_header = {"backgroundColor": "#1F3864", "color": "white", "fontWeight": "bold", "fontSize": "12px"}
    tbl_style_cell   = {"fontSize": "12px", "padding": "6px 10px", "textAlign": "right",
                        "fontFamily": "Segoe UI, sans-serif"}

    def _build_tbl(df_sub, tbl_id, header_color):
        rows = []
        for _, r in df_sub.iterrows():
            rows.append({
                "rut":        r["RUT"],
                "nombre":     r["nombre"],
                "zona":       r["zona"],
                f"Venta {mes_nombre} 2025": _fmt_abs(r["venta_2025_ytd"]),
                f"Venta {mes_nombre} 2026": _fmt_abs(r["venta_2026_ytd"]),
                "Diferencia $": _fmt_abs(r["diferencia"]),
                "Var %":       _fmt_pct(r["dif_pct"]),
            })
        tot_v25 = df_sub["venta_2025_ytd"].sum()
        tot_v26 = df_sub["venta_2026_ytd"].sum()
        tot_dif = df_sub["diferencia"].sum()
        tot_pct = (tot_dif / tot_v25 * 100) if tot_v25 > 0 else float("nan")
        rows.append({
            "rut": "", "nombre": f"── TOTAL ({len(df_sub)} clientes)", "zona": "",
            f"Venta {mes_nombre} 2025": _fmt_abs(tot_v25),
            f"Venta {mes_nombre} 2026": _fmt_abs(tot_v26),
            "Diferencia $": _fmt_abs(tot_dif),
            "Var %": _fmt_pct(tot_pct),
        })
        cols = ["rut", "nombre", "zona",
                f"Venta {mes_nombre} 2025", f"Venta {mes_nombre} 2026",
                "Diferencia $", "Var %"]
        return dash_table.DataTable(
            id=tbl_id,
            data=rows,
            columns=[{"name": c, "id": c} for c in cols],
            style_table={"overflowX": "auto"},
            style_header={**tbl_style_header, "backgroundColor": header_color},
            style_cell=tbl_style_cell,
            style_cell_conditional=[
                {"if": {"column_id": "rut"},    "textAlign": "left", "maxWidth": "90px"},
                {"if": {"column_id": "nombre"}, "textAlign": "left", "maxWidth": "180px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
                {"if": {"column_id": "zona"},   "textAlign": "left", "maxWidth": "120px",
                 "overflow": "hidden", "textOverflow": "ellipsis"},
            ],
            style_data_conditional=[
                {"if": {"filter_query": '{nombre} contains "── TOTAL"'},
                 "fontWeight": "700", "backgroundColor": "#e8edf7"},
            ],
            page_action="none",
            sort_action="native",
            row_selectable="single",
        )

    df_caida = df.sort_values("diferencia").head(25)
    df_alza  = df.sort_values("diferencia", ascending=False).head(25)

    return (_build_tbl(df_caida, "zona-tabla-caida", "#922b21"),
            _build_tbl(df_alza,  "zona-tabla-alza",  "#1e8449"))


@app.callback(
    Output("zona-cli-detalle", "children"),
    Input("zona-tabla-caida", "selected_rows"),
    Input("zona-tabla-alza",  "selected_rows"),
    State("zona-tabla-caida", "data"),
    State("zona-tabla-alza",  "data"),
    prevent_initial_call=True,
)
def zona_cli_drill(sel_caida, sel_alza, data_caida, data_alza):
    from dash import ctx
    triggered = ctx.triggered_id
    if triggered == "zona-tabla-caida" and sel_caida:
        row = data_caida[sel_caida[0]]
    elif triggered == "zona-tabla-alza" and sel_alza:
        row = data_alza[sel_alza[0]]
    else:
        raise PreventUpdate

    rut    = row["rut"]
    nombre = row["nombre"]
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    df = _load_detalle_cliente(rut)
    if df.empty:
        return html.Div(style=CARD_STYLE, children=[_empty_msg(f"Sin detalle para {nombre}")])

    tbl_style_header = {"backgroundColor": "#2E75B6", "color": "white", "fontWeight": "bold", "fontSize": "12px"}
    tbl_style_cell   = {"fontSize": "12px", "padding": "6px 10px", "textAlign": "right",
                        "fontFamily": "Segoe UI, sans-serif"}

    # Por Producto
    prod_rows = []
    for _, r in df.iterrows():
        prod_rows.append({
            "codigo":      r["codigo"],
            "descripcion": r["descripcion"],
            "categoria":   r["categoria"],
            f"Venta {mes_nombre} 2025": _fmt_abs(r["venta_2025"]),
            f"Venta {mes_nombre} 2026": _fmt_abs(r["venta_2026"]),
            "Diferencia $": _fmt_abs(r["diferencia"]),
            "Var %": _fmt_pct(r["dif_pct"]),
        })
    tot_v25 = df["venta_2025"].sum(); tot_v26 = df["venta_2026"].sum()
    tot_dif = tot_v26 - tot_v25
    tot_pct = (tot_dif / tot_v25 * 100) if tot_v25 > 0 else float("nan")
    prod_rows.append({
        "codigo": "", "descripcion": "── TOTAL", "categoria": "",
        f"Venta {mes_nombre} 2025": _fmt_abs(tot_v25),
        f"Venta {mes_nombre} 2026": _fmt_abs(tot_v26),
        "Diferencia $": _fmt_abs(tot_dif), "Var %": _fmt_pct(tot_pct),
    })

    prod_cols = ["codigo", "descripcion", "categoria",
                 f"Venta {mes_nombre} 2025", f"Venta {mes_nombre} 2026",
                 "Diferencia $", "Var %"]

    # Por Categoría
    df_cat = df.groupby("categoria", as_index=False).agg(
        venta_2025=("venta_2025", "sum"), venta_2026=("venta_2026", "sum")
    )
    df_cat["diferencia"] = df_cat["venta_2026"] - df_cat["venta_2025"]
    df_cat["dif_pct"]    = df_cat.apply(
        lambda r: r["diferencia"] / r["venta_2025"] * 100 if r["venta_2025"] > 0 else float("nan"), axis=1
    )
    cat_rows = []
    for _, r in df_cat.iterrows():
        cat_rows.append({
            "categoria": r["categoria"],
            f"Venta {mes_nombre} 2025": _fmt_abs(r["venta_2025"]),
            f"Venta {mes_nombre} 2026": _fmt_abs(r["venta_2026"]),
            "Diferencia $": _fmt_abs(r["diferencia"]),
            "Var %": _fmt_pct(r["dif_pct"]),
        })
    cat_rows.append({
        "categoria": "── TOTAL",
        f"Venta {mes_nombre} 2025": _fmt_abs(tot_v25),
        f"Venta {mes_nombre} 2026": _fmt_abs(tot_v26),
        "Diferencia $": _fmt_abs(tot_dif), "Var %": _fmt_pct(tot_pct),
    })
    cat_cols = ["categoria", f"Venta {mes_nombre} 2025", f"Venta {mes_nombre} 2026",
                "Diferencia $", "Var %"]

    total_row_cond = [{"if": {"filter_query": '{descripcion} = "── TOTAL" || {categoria} = "── TOTAL"'},
                       "fontWeight": "700", "backgroundColor": "#e8edf7"}]

    return html.Div(style=CARD_STYLE, children=[
        html.H4(f"📋 Detalle cliente: {nombre} ({rut})",
                style={"color": "#1F3864", "fontSize": "14px", "fontWeight": "700", "marginBottom": "16px"}),
        html.Div(style={"display": "grid", "gridTemplateColumns": "1fr 1fr", "gap": "24px"}, children=[
            html.Div(children=[
                html.H5("Por Categoría", style={"fontSize": "12px", "fontWeight": "700",
                                                 "marginBottom": "8px", "color": "#2E75B6"}),
                dash_table.DataTable(
                    data=cat_rows, columns=[{"name": c, "id": c} for c in cat_cols],
                    style_table={"overflowX": "auto"}, style_header=tbl_style_header,
                    style_cell=tbl_style_cell,
                    style_cell_conditional=[{"if": {"column_id": "categoria"}, "textAlign": "left"}],
                    style_data_conditional=total_row_cond,
                    page_action="none", sort_action="native",
                ),
            ]),
            html.Div(children=[
                html.H5("Por Producto", style={"fontSize": "12px", "fontWeight": "700",
                                                "marginBottom": "8px", "color": "#2E75B6"}),
                dash_table.DataTable(
                    data=prod_rows, columns=[{"name": c, "id": c} for c in prod_cols],
                    style_table={"overflowX": "auto"}, style_header=tbl_style_header,
                    style_cell=tbl_style_cell,
                    style_cell_conditional=[
                        {"if": {"column_id": "codigo"},      "textAlign": "left"},
                        {"if": {"column_id": "descripcion"}, "textAlign": "left", "maxWidth": "220px",
                         "overflow": "hidden", "textOverflow": "ellipsis"},
                        {"if": {"column_id": "categoria"},   "textAlign": "left"},
                    ],
                    style_data_conditional=total_row_cond,
                    page_action="none", sort_action="native",
                ),
            ]),
        ]),
    ])


@app.callback(
    Output("zona-kpi-row", "children"),
    Output("zona-tabla-container", "children"),
    Input("zona-cat-filtro", "value"),
    prevent_initial_call=True,
)
def update_zona_tabla(cat_filtro):
    if not cat_filtro or cat_filtro == "TODAS":
        df = _df_zona.copy()
    else:
        if _df_zona_cat.empty:
            return html.Div(), _empty_msg("Sin datos de categoría disponibles.")
        dfc = _df_zona_cat[_df_zona_cat["categoria"] == cat_filtro].copy()
        if dfc.empty:
            return html.Div(), _empty_msg("Sin datos para la categoría seleccionada.")
        df = dfc.groupby("zona", as_index=False).agg(
            ppto_anual=("ppto_anual", "sum"),
            ppto_ytd=("ppto_ytd", "sum"),
            venta_2026_ytd=("venta_2026_ytd", "sum"),
            venta_2025_ytd=("venta_2025_ytd", "sum"),
        )
        df["cumpl_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] / r["ppto_ytd"] * 100) if r["ppto_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["gap"] = df["venta_2026_ytd"] - df["ppto_ytd"]
        df["crec_pct"] = df.apply(
            lambda r: (r["venta_2026_ytd"] - r["venta_2025_ytd"]) / r["venta_2025_ytd"] * 100
            if r["venta_2025_ytd"] > 0 else float("nan"),
            axis=1
        )
        df["proyeccion_anual"] = df["venta_2026_ytd"] / _MES_ACT * 12 if _MES_ACT > 0 else 0
        df["gap_vs_ppto_anual"] = df["proyeccion_anual"] - df["ppto_anual"]
        df = df.sort_values("ppto_anual", ascending=False)

    df = df[df["ppto_anual"] > 0].copy()
    return _build_zona_content(df)


@app.callback(
    Output("zona-top25-container", "children"),
    Output("zona-caida-container", "children"),
    Input("zona-selector", "value"),
    prevent_initial_call=True,
)
def update_zona_top25_selector(zona):
    if not zona:
        raise PreventUpdate
    return _build_top25_clientes(zona), _build_caida_clientes(zona)


@app.callback(
    Output("desa-main-tabla", "data"),
    Output("desa-main-tabla", "columns"),
    Output("desa-main-tabla", "selected_rows"),
    Input("desa-cat-filtro", "value"),
    prevent_initial_call=True,
)
def filter_desa_tabla(cat_filtro):
    df = _df_desalineacion
    if df.empty:
        raise PreventUpdate
    if cat_filtro and cat_filtro != "__all__":
        df = df[df["categoria"] == cat_filtro]
    table_data, cols = _build_desa_pivot(df)
    columns = [{"name": c, "id": c} for c in cols]
    return table_data, columns, []


@app.callback(
    Output("desa-detalle-container", "children"),
    Input("desa-main-tabla", "selected_rows"),
    State("desa-main-tabla", "data"),
    prevent_initial_call=True,
)
def render_desa_detalle(selected_rows, table_data):
    if not selected_rows or not table_data:
        raise PreventUpdate
    row = table_data[selected_rows[0]]
    cat = row.get("Categoría", "")
    if not cat:
        raise PreventUpdate
    return _render_desa_zona(cat)


# ─── Callbacks Tab Precios ────────────────────────────────────────────────────

_PRECIOS_CATS_EXCLUIR = {"Activo fijo", "DES", ""}

def _precios_filtrar(cat_filtro, zona_filtro):
    """Aplica filtros de categoría y zona al DataFrame de precios."""
    df = _df_precios.copy() if not _df_precios.empty else pd.DataFrame()
    if df.empty:
        return df
    # Excluir categorías no analizables
    df = df[~df["categoria"].isin(_PRECIOS_CATS_EXCLUIR)]
    if cat_filtro and cat_filtro != "__all__":
        df = df[df["categoria"] == cat_filtro]
    if zona_filtro and zona_filtro != "__all__":
        df = df[df["zona"] == zona_filtro]
    return df


@app.callback(
    Output("precio-resumen-cat",   "children"),
    Output("precio-resumen-zona",  "children"),
    Output("precio-top-productos", "children"),
    Input("precio-cat-filtro",  "value"),
    Input("precio-zona-filtro", "value"),
    prevent_initial_call=False,
)
def update_precios_filtros(cat_filtro, zona_filtro):
    df = _precios_filtrar(cat_filtro, zona_filtro)
    if df.empty:
        msg = _empty_msg("Sin datos para los filtros seleccionados.")
        return msg, msg, msg
    tcat = _precios_tabla_categoria(df)
    tzon = _precios_tabla_zona(df)
    ttop = _precios_top_productos(df)
    return tcat, tzon, ttop


@app.callback(
    Output("precio-drill-zona", "children"),
    Input("precios-zona-tabla", "selected_rows"),
    State("precios-zona-tabla", "data"),
    State("precio-cat-filtro",  "value"),
    prevent_initial_call=True,
)
def precio_drill_zona(selected_rows, data, cat_filtro):
    if not selected_rows or not data:
        raise PreventUpdate
    row = data[selected_rows[0]]
    zona = row.get("Zona", "")
    if not zona:
        raise PreventUpdate

    df = _df_precios.copy() if not _df_precios.empty else pd.DataFrame()
    if df.empty:
        raise PreventUpdate
    df = df[df["zona"] == zona]
    if cat_filtro and cat_filtro != "__all__":
        df = df[df["categoria"] == cat_filtro]

    # Tabla por cliente
    rows_cli = []
    for (rut, nombre), g in df.groupby(["rut", "nombre_cliente"], sort=False):
        p26 = g["venta_2026"].sum() / g["cant_2026"].sum() if g["cant_2026"].sum() > 0 else float("nan")
        p25 = g["venta_2025"].sum() / g["cant_2025"].sum() if g["cant_2025"].sum() > 0 else float("nan")
        dpct = (p26 - p25) / p25 * 100 if pd.notna(p25) and p25 > 0 and pd.notna(p26) else float("nan")
        ep = g["efecto_precio"].dropna().sum()
        ev = g["efecto_volumen"].dropna().sum()
        sem = "🟢" if pd.notna(dpct) and dpct > 0 else "🔴" if pd.notna(dpct) and dpct < 0 else "⚫"
        rows_cli.append({
            "rut": rut,
            "Cliente": nombre or rut,
            "Precio 2025": _fmt_abs(p25) if pd.notna(p25) else "—",
            "Precio 2026": _fmt_abs(p26) if pd.notna(p26) else "—",
            "Δ Precio %": f"{sem} {_fmt_pct(dpct)}" if pd.notna(dpct) else "—",
            "Efecto precio": _fmt_abs(ep),
            "Efecto volumen": _fmt_abs(ev),
            "Venta 2026": _fmt_abs(g["venta_2026"].sum()),
            "Venta 2025": _fmt_abs(g["venta_2025"].sum()),
            "# Productos": len(g["codigo"].unique()),
        })
    rows_cli.sort(key=lambda r: abs(
        float(r["Efecto precio"].replace("$","").replace(".","").replace("-","") or "0")
    ), reverse=True)

    tbl_header = {"backgroundColor": "#2E75B6", "color": "white",
                  "fontWeight": "700", "fontSize": "11px"}
    tbl_cell   = {"fontSize": "11px", "padding": "5px 8px", "textAlign": "right"}
    cli_cols   = ["Cliente", "Precio 2025", "Precio 2026", "Δ Precio %",
                  "Efecto precio", "Efecto volumen", "Venta 2026", "Venta 2025", "# Productos"]

    tabla_cli = dash_table.DataTable(
        data=rows_cli,
        columns=[{"name": c, "id": c} for c in cli_cols],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": "Cliente"}, "textAlign": "left"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Δ Precio %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Δ Precio %} contains "🔴"'}, "color": "#c0392b"},
        ],
        sort_action="native",
        page_action="none",
    )

    # Top productos de la zona
    top_prod = _precios_top_productos(df, n=10)

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"🔍 Detalle zona: {zona}", "#2E75B6"),
        _seccion("Análisis por Cliente", "#555"),
        tabla_cli,
        html.Div(style={"marginTop": "20px"}, children=[
            _seccion("Productos destacados en esta zona", "#555"),
            top_prod,
        ]),
    ])


@app.callback(
    Output("precio-drill-cat", "children"),
    Input("precios-cat-tabla", "selected_rows"),
    State("precios-cat-tabla", "data"),
    State("precio-zona-filtro", "value"),
    prevent_initial_call=True,
)
def precio_drill_cat(selected_rows, data, zona_filtro):
    if not selected_rows or not data:
        raise PreventUpdate
    row = data[selected_rows[0]]
    cat = row.get("Categoría", "")
    if not cat or cat == "── TOTAL":
        raise PreventUpdate

    df = _df_precios.copy() if not _df_precios.empty else pd.DataFrame()
    if df.empty:
        raise PreventUpdate
    df = df[~df["categoria"].isin(_PRECIOS_CATS_EXCLUIR)]
    df = df[df["categoria"] == cat]
    if zona_filtro and zona_filtro != "__all__":
        df = df[df["zona"] == zona_filtro]

    # Tabla por zona dentro de la categoría
    rows_zona = []
    for zona, g in df.groupby("zona", sort=False):
        p26 = g["venta_2026"].sum() / g["cant_2026"].sum() if g["cant_2026"].sum() > 0 else float("nan")
        p25 = g["venta_2025"].sum() / g["cant_2025"].sum() if g["cant_2025"].sum() > 0 else float("nan")
        dpct = (p26 - p25) / p25 * 100 if pd.notna(p25) and p25 > 0 and pd.notna(p26) else float("nan")
        ep = g["efecto_precio"].dropna().sum()
        ev = g["efecto_volumen"].dropna().sum()
        sem = "🟢" if pd.notna(dpct) and dpct > 0 else "🔴" if pd.notna(dpct) and dpct < 0 else "⚫"
        rows_zona.append({
            "Zona": zona,
            "Precio 2025": _fmt_abs(p25) if pd.notna(p25) else "—",
            "Precio 2026": _fmt_abs(p26) if pd.notna(p26) else "—",
            "Δ Precio %": f"{sem} {_fmt_pct(dpct)}" if pd.notna(dpct) else "—",
            "Efecto precio": _fmt_abs(ep),
            "Efecto volumen": _fmt_abs(ev),
            "Venta 2026": _fmt_abs(g["venta_2026"].sum()),
            "Venta 2025": _fmt_abs(g["venta_2025"].sum()),
            "# Clientes": g["rut"].nunique(),
            "# Productos": g["codigo"].nunique(),
        })
    rows_zona.sort(key=lambda r: -abs(
        float(r["Efecto precio"].replace("$","").replace(".","").replace("-","") or "0")
    ))
    # Fila TOTAL
    tp26 = df["venta_2026"].sum() / df["cant_2026"].sum() if df["cant_2026"].sum() > 0 else float("nan")
    tp25 = df["venta_2025"].sum() / df["cant_2025"].sum() if df["cant_2025"].sum() > 0 else float("nan")
    tdpct = (tp26 - tp25) / tp25 * 100 if pd.notna(tp25) and tp25 > 0 and pd.notna(tp26) else float("nan")
    tsem = "🟢" if pd.notna(tdpct) and tdpct > 0 else "🔴" if pd.notna(tdpct) and tdpct < 0 else "⚫"
    rows_zona.append({
        "Zona": "── TOTAL",
        "Precio 2025": _fmt_abs(tp25) if pd.notna(tp25) else "—",
        "Precio 2026": _fmt_abs(tp26) if pd.notna(tp26) else "—",
        "Δ Precio %": f"{tsem} {_fmt_pct(tdpct)}" if pd.notna(tdpct) else "—",
        "Efecto precio":  _fmt_abs(df["efecto_precio"].dropna().sum()),
        "Efecto volumen": _fmt_abs(df["efecto_volumen"].dropna().sum()),
        "Venta 2026": _fmt_abs(df["venta_2026"].sum()),
        "Venta 2025": _fmt_abs(df["venta_2025"].sum()),
        "# Clientes": "—", "# Productos": "—",
    })

    tbl_header = {"backgroundColor": "#7b2d8b", "color": "white",
                  "fontWeight": "700", "fontSize": "11px"}
    tbl_cell   = {"fontSize": "11px", "padding": "5px 8px", "textAlign": "right"}
    tabla_zona = dash_table.DataTable(
        data=rows_zona,
        columns=[{"name": c, "id": c} for c in
                 ["Zona", "Precio 2025", "Precio 2026", "Δ Precio %",
                  "Efecto precio", "Efecto volumen", "Venta 2026", "Venta 2025",
                  "# Clientes", "# Productos"]],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": "Zona"}, "textAlign": "left", "fontWeight": "600"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Δ Precio %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Δ Precio %} contains "🔴"'}, "color": "#c0392b"},
            {"if": {"filter_query": '{Zona} = "── TOTAL"'},
             "fontWeight": "700", "backgroundColor": "#e8edf7"},
        ],
        sort_action="native",
        page_action="none",
    )

    top_prod = _precios_top_productos(df, n=10)

    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"🔍 Detalle categoría: {cat}", "#7b2d8b"),
        _seccion("Desglose por Zona", "#555"),
        tabla_zona,
        html.Div(style={"marginTop": "20px"}, children=[
            _seccion("Productos destacados en esta categoría", "#555"),
            top_prod,
        ]),
    ])


# ─── Callbacks PPTO vs Venta ───────────────────────────────────────────────────

_PVV_TIPO_TRAZ   = "PRESUPUESTO_TRAZABLE"
_PVV_TIPO_INCR   = "INCREMENTALES_SIN_PRODUCTOS_NUEVOS"
_PVV_TIPO_NUEVOS = "INCREMENTALES_CON_PRODUCTOS_NUEVOS"
_PVV_TIPO_LABELS = {
    _PVV_TIPO_TRAZ:   "Trazable",
    _PVV_TIPO_INCR:   "Incr. s/PN",
    _PVV_TIPO_NUEVOS: "Incr. c/PN",
}


@app.callback(
    Output("pvv-kpi-container", "children"),
    Output("pvv-tabla-container", "children"),
    Output("pvv-analisis-caida", "children"),
    Input("pvv-zona-filter", "value"),
    Input("pvv-cat-filter", "value"),
)
def pvv_update_main(zona, cat):
    df = _df_ppto_vs_venta
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    if df.empty:
        msg = _empty_msg("Sin datos. Presione 🔄 Actualizar con VPN activa.")
        return msg, msg, html.Div()

    filtrado = zona != "TODAS" or cat != "TODAS"
    df = df.copy()

    if zona and zona != "TODAS":
        df = df[df["VENDEDOR_ACTUAL"] == zona]
    if cat and cat != "TODAS":
        df = df[df["CATEGORIA_2026"] == cat]

    tbl_header = {"backgroundColor": "#1F3864", "color": "white",
                  "fontWeight": "700", "fontSize": "12px", "padding": "6px 8px"}
    tbl_cell   = {"fontSize": "12px", "padding": "5px 8px", "fontFamily": "Segoe UI, sans-serif"}

    # ── KPIs — globales desde vista cuando no hay filtro, calculados desde df si hay filtro ─
    if not filtrado and _pvv_kpis:
        ppto_total  = _pvv_kpis["ppto_total"]
        venta_ytd   = _pvv_kpis["venta_ytd"]
        meta_ytd    = _pvv_kpis["meta_ytd"]
        alcance_ytd = _pvv_kpis["alcance_ytd"] * 100
        cumpl_ppto  = _pvv_kpis["cumpl_ppto"]  * 100
        gap_meta    = _pvv_kpis["gap_meta"]
        venta_25    = _pvv_kpis["venta_ytd_25"]
        var_abs_25  = _pvv_kpis["var_abs_25"]
    else:
        ppto_total  = df["PPTO_2026"].sum()
        venta_ytd   = df["VENTA_2026"].sum()
        venta_25    = df["VENTA_2025"].sum()
        meta_ytd    = None
        alcance_ytd = venta_ytd / ppto_total * 100 if ppto_total > 0 else 0
        cumpl_ppto  = alcance_ytd
        gap_meta    = None
        var_abs_25  = venta_ytd - venta_25

    def _sem_pct(v):
        return "🟢" if v >= 80 else ("🟡" if v >= 50 else "🔴")

    kpi_items = [
        _kpi_card("PPTO TT 2026",             _fmt(ppto_total), "#1F3864"),
        _kpi_card(f"Meta YTD {mes_nombre}",
                  _fmt(meta_ytd) if meta_ytd is not None else "—", "#2E75B6"),
        _kpi_card(f"Venta YTD {mes_nombre}",  _fmt(venta_ytd),  "#2E75B6"),
        _kpi_card("Cumplim. YTD",
                  f"{_sem_pct(alcance_ytd)} {alcance_ytd:.1f}%",
                  "#1e7e34" if alcance_ytd >= 80 else ("#e67e22" if alcance_ytd >= 50 else "#c0392b")),
        _kpi_card("Diff Meta vs Venta",
                  _fmt(gap_meta) if gap_meta is not None else "—",
                  "#1e7e34" if (gap_meta is not None and gap_meta >= 0) else "#c0392b"),
        _kpi_card("Cumplim. PPTO",
                  f"{_sem_pct(cumpl_ppto)} {cumpl_ppto:.1f}%",
                  "#1e7e34" if cumpl_ppto >= 80 else ("#e67e22" if cumpl_ppto >= 50 else "#c0392b")),
        _kpi_card("Var vs 2025",
                  _fmt(var_abs_25),
                  "#1e7e34" if var_abs_25 >= 0 else "#c0392b"),
    ]
    kpis_div = _kpi_row(kpi_items)

    # ── Tabla por CATEGORIA_2026 (solo categorías con PPTO > 0) ──────────────
    _PVV_SIN_CLI = "PRESUPUESTO_SIN_CLIENTE"
    df_ppto = df[df["CATEGORIA_2026"].isin(["EQM", "EVA", "MAH", "SQ"])].copy()

    cat_rows = []
    for cat_name, gdf in df_ppto.groupby("CATEGORIA_2026"):
        traz     = gdf.loc[gdf["TIPO_ANALISIS"] == _PVV_TIPO_TRAZ,   "PPTO_2026"].sum()
        prod_new = gdf.loc[gdf["TIPO_ANALISIS"] == _PVV_TIPO_NUEVOS, "PPTO_2026"].sum()
        incr     = gdf.loc[gdf["TIPO_ANALISIS"] == _PVV_TIPO_INCR,   "PPTO_2026"].sum()
        sin_cli  = gdf.loc[gdf["TIPO_ANALISIS"] == _PVV_SIN_CLI,     "PPTO_2026"].sum()
        ppto_t   = gdf["PPTO_2026"].sum()
        v26      = gdf["VENTA_2026"].sum()
        v25      = gdf["VENTA_2025"].sum()
        alc      = v26 / ppto_t * 100 if ppto_t > 0 else 0
        cat_rows.append({
            "Categoría":  str(cat_name),
            "Trazable":   _fmt(traz),
            "Prod. Nuevo":_fmt(prod_new),
            "Incremental":_fmt(incr),
            "Sin Cliente":_fmt(sin_cli),
            "PPTO 2026":  _fmt(ppto_t),
            "Venta 2026": _fmt(v26),
            "Alcance %":  f"{_sem_pct(alc)} {alc:.1f}%",
            "Venta 2025": _fmt(v25),
            "_ppto_num":  ppto_t,
        })
    cat_rows.sort(key=lambda r: -r["_ppto_num"])

    # Fila total
    ttraz    = df_ppto.loc[df_ppto["TIPO_ANALISIS"] == _PVV_TIPO_TRAZ,   "PPTO_2026"].sum()
    tpn      = df_ppto.loc[df_ppto["TIPO_ANALISIS"] == _PVV_TIPO_NUEVOS, "PPTO_2026"].sum()
    tincr    = df_ppto.loc[df_ppto["TIPO_ANALISIS"] == _PVV_TIPO_INCR,   "PPTO_2026"].sum()
    tsin_cli = df_ppto.loc[df_ppto["TIPO_ANALISIS"] == _PVV_SIN_CLI,     "PPTO_2026"].sum()
    tppto    = df_ppto["PPTO_2026"].sum()
    tv26     = df_ppto["VENTA_2026"].sum()
    tv25     = df_ppto["VENTA_2025"].sum()
    talc     = tv26 / tppto * 100 if tppto > 0 else 0
    cat_rows.append({
        "Categoría":  "── TOTAL",
        "Trazable":   _fmt(ttraz),
        "Prod. Nuevo":_fmt(tpn),
        "Incremental":_fmt(tincr),
        "Sin Cliente":_fmt(tsin_cli),
        "PPTO 2026":  _fmt(tppto),
        "Venta 2026": _fmt(tv26),
        "Alcance %":  f"{talc:.1f}%",
        "Venta 2025": _fmt(tv25),
        "_ppto_num":  0,
    })

    cols_tabla = ["Categoría", "Trazable", "Prod. Nuevo", "Incremental",
                  "Sin Cliente", "PPTO 2026", "Venta 2026", "Alcance %", "Venta 2025"]
    tabla = dash_table.DataTable(
        id="pvv-cat-tabla",
        data=cat_rows,
        columns=[{"name": c, "id": c} for c in cols_tabla],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": "Categoría"}, "textAlign": "left", "fontWeight": "600"},
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Categoría} = "── TOTAL"'},
             "fontWeight": "700", "backgroundColor": "#e8edf7"},
            {"if": {"filter_query": '{Alcance %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Alcance %} contains "🟡"'}, "color": "#e67e22"},
            {"if": {"filter_query": '{Alcance %} contains "🔴"'}, "color": "#c0392b"},
        ],
        row_selectable="single",
        selected_rows=[],
        page_action="none",
        sort_action="native",
        style_table={"overflowX": "auto"},
        tooltip_header={
            "Trazable":    "PPTO trazable: productos en presupuesto original con cliente identificado",
            "Prod. Nuevo": "PPTO incremental con productos nuevos",
            "Incremental": "PPTO incremental sin productos nuevos",
            "Sin Cliente": "PPTO sin cliente asignado",
        },
        tooltip_delay=0,
        tooltip_duration=None,
    )

    nota = "" if not filtrado else " (filtrado — KPIs globales se muestran sin filtro)"
    try:
        analisis_caida = _pvv_analisis_caida_inner(zona, cat)
    except Exception as e:
        import traceback
        err = traceback.format_exc()
        print(f"[ERROR _pvv_analisis_caida_inner] {e}\n{err}")
        analisis_caida = html.Div(style=CARD_STYLE, children=[
            _seccion("📉 Análisis de Caída", "#c0392b"),
            html.Pre(str(err)[:2000], style={"fontSize": "10px", "color": "red", "whiteSpace": "pre-wrap"}),
        ])
    return (
        html.Div(style=CARD_STYLE, children=[
            _seccion(f"📊 Indicadores YTD {mes_nombre} 2026{nota}"),
            kpis_div,
        ]),
        html.Div(style=CARD_STYLE, children=[
            _seccion("🗂️ PPTO por Categoría — clic en fila para ver productos"),
            html.P(
                "Trazable: con cliente | Prod. Nuevo: incr. c/PN | Incremental: incr. s/PN | Sin Cliente: sin RUT asignado",
                style={"fontSize": "11px", "color": "#666", "marginBottom": "10px"}),
            tabla,
        ]),
        analisis_caida,
    )


@app.callback(
    Output("pvv-detalle-container", "children"),
    Input("pvv-cat-tabla", "selected_rows"),
    State("pvv-cat-tabla", "data"),
    State("pvv-zona-filter", "value"),
    prevent_initial_call=True,
)
def pvv_update_detalle(selected_rows, data, zona):
    if not selected_rows or not data:
        raise PreventUpdate
    row = data[selected_rows[0]]
    cat = row.get("Categoría", "")
    if not cat or cat == "── TOTAL":
        raise PreventUpdate

    df = _df_ppto_vs_venta
    if df.empty:
        raise PreventUpdate

    df = df.copy()
    for c in ["VENDEDOR_ACTUAL", "CATEGORIA_2026", "TIPO_ANALISIS", "ESTADO_ANALISIS"]:
        if c in df.columns:
            df[c] = df[c].astype(str).str.strip()

    df = df[df["CATEGORIA_2026"] == cat]
    if zona and zona != "TODAS":
        df = df[df["VENDEDOR_ACTUAL"] == zona]
    if df.empty:
        return _empty_msg(f"Sin productos para categoría {cat}.")

    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    tbl_header = {"backgroundColor": "#2E75B6", "color": "white",
                  "fontWeight": "700", "fontSize": "11px", "padding": "5px 8px"}
    tbl_cell   = {"fontSize": "11px", "padding": "5px 8px", "fontFamily": "Segoe UI, sans-serif"}

    # Agrupar por producto para consolidar zonas
    grp = df.groupby(["CODIGO", "DESCRIPCION", "TIPO_ANALISIS", "ESTADO_ANALISIS"]).agg(
        ppto=("PPTO_2026", "sum"),
        venta_26=("VENTA_2026", "sum"),
        venta_25=("VENTA_2025", "sum"),
    ).reset_index().sort_values("ppto", ascending=False)

    prod_rows = []
    for _, r in grp.head(60).iterrows():
        ppto_v = r["ppto"]
        v26    = r["venta_26"]
        v25    = r["venta_25"]
        alc    = v26 / ppto_v * 100 if ppto_v > 0 else float("nan")
        prod_rows.append({
            "Código":      str(r["CODIGO"])[:20],
            "Descripción": str(r["DESCRIPCION"])[:65],
            "Tipo":        _PVV_TIPO_LABELS.get(str(r["TIPO_ANALISIS"]), str(r["TIPO_ANALISIS"])[:20]),
            "Estado":      str(r["ESTADO_ANALISIS"])[:30],
            "PPTO 2026":   _fmt(ppto_v),
            "Venta 2026":  _fmt(v26),
            "Alcance %":   f"{_sem(alc)} {alc:.1f}%" if not (isinstance(alc, float) and math.isnan(alc)) else "—",
            "Venta 2025":  _fmt(v25),
        })

    tabla_prod = dash_table.DataTable(
        data=prod_rows,
        columns=[{"name": c, "id": c} for c in
                 ["Código", "Descripción", "Tipo", "Estado",
                  "PPTO 2026", "Venta 2026", "Alcance %", "Venta 2025"]],
        style_header=tbl_header,
        style_cell=tbl_cell,
        style_cell_conditional=[
            {"if": {"column_id": c}, "textAlign": "left"}
            for c in ["Código", "Descripción", "Tipo", "Estado"]
        ],
        style_data_conditional=[
            {"if": {"filter_query": '{Alcance %} contains "🟢"'}, "color": "#1e7e34"},
            {"if": {"filter_query": '{Alcance %} contains "🟡"'}, "color": "#e67e22"},
            {"if": {"filter_query": '{Alcance %} contains "🔴"'}, "color": "#c0392b"},
            {"if": {"filter_query": '{Estado} contains "FUERA"'},  "backgroundColor": "#fdf3f8"},
            {"if": {"filter_query": '{Estado} contains "SIN VENTA"'}, "backgroundColor": "#fff5f5"},
        ],
        page_size=20,
        sort_action="native",
        style_table={"maxHeight": "420px", "overflowY": "auto", "overflowX": "auto"},
    )

    zona_label = zona if zona and zona != "TODAS" else "Todas las zonas"
    return html.Div(style=CARD_STYLE, children=[
        _seccion(f"🔍 Detalle — {cat} | {zona_label} | Top 60 productos", "#2E75B6"),
        tabla_prod,
    ])


def _pvv_analisis_caida_inner(zona, cat):
    mes_nombre = _MESES_NOMBRE.get(_MES_ACT, str(_MES_ACT))
    TH = {"backgroundColor": "#1F3864", "color": "white",
          "fontWeight": "700", "fontSize": "11px", "padding": "5px 8px"}
    TC = {"fontSize": "11px", "padding": "5px 7px", "fontFamily": "Segoe UI, sans-serif"}

    # ── Datos de precios con filtros ──────────────────────────────────────────
    df_p = _df_precios.copy() if not _df_precios.empty else pd.DataFrame()
    if not df_p.empty:
        if zona and zona != "TODAS":
            df_p = df_p[df_p["zona"] == zona]
        if cat and cat != "TODAS":
            df_p = df_p[df_p["categoria"] == cat]

    # ── Totales base para waterfall ───────────────────────────────────────────
    dfpv = _df_ppto_vs_venta.copy()
    if not dfpv.empty:
        if zona and zona != "TODAS":
            dfpv = dfpv[dfpv["VENDEDOR_ACTUAL"] == zona]
        if cat and cat != "TODAS":
            dfpv = dfpv[dfpv["CATEGORIA_2026"] == cat]
        dfpv = dfpv[dfpv["CATEGORIA_2026"].isin(["EQM", "EVA", "MAH", "SQ"])]
        v26_total = dfpv["VENTA_2026"].sum()
        v25_total = dfpv["VENTA_2025"].sum()
    elif _pvv_kpis:
        v26_total = _pvv_kpis.get("venta_ytd", 0)
        v25_total = _pvv_kpis.get("venta_ytd_25", 0)
    else:
        v26_total = v25_total = 0
    delta_total = v26_total - v25_total
    pct_caida   = delta_total / v25_total * 100 if v25_total > 0 else 0

    if not df_p.empty and "efecto_precio" in df_p.columns:
        ep_total = df_p["efecto_precio"].dropna().sum()
        ev_total = df_p["efecto_volumen"].dropna().sum()
    else:
        ep_total = ev_total = float("nan")

    # ── 1. Waterfall bridge ───────────────────────────────────────────────────
    wf_x, wf_y, wf_col, wf_base = [], [], [], []
    running = v25_total

    def _add_bar(label, val, running_in):
        wf_x.append(label)
        wf_y.append(abs(val))
        wf_base.append(running_in + min(val, 0))
        wf_col.append("#c0392b" if val < 0 else "#1e7e34")
        return running_in + val

    wf_x.append(f"Venta {mes_nombre}\n2025"); wf_y.append(v25_total)
    wf_base.append(0); wf_col.append("#1F3864")

    if not math.isnan(ep_total) and abs(ep_total) > 1e5:
        running = _add_bar("Efecto\nPrecio", ep_total, running)
    if not math.isnan(ev_total) and abs(ev_total) > 1e5:
        running = _add_bar("Efecto\nVolumen", ev_total, running)
    otros = delta_total - (
        (ep_total if not math.isnan(ep_total) else 0) +
        (ev_total if not math.isnan(ev_total) else 0)
    )
    if abs(otros) > 1e5:
        running = _add_bar("Mix /\nNuevos", otros, running)

    wf_x.append(f"Venta {mes_nombre}\n2026"); wf_y.append(v26_total)
    wf_base.append(0); wf_col.append("#2E75B6")

    fig_wf = go.Figure()
    for i, (x, y, col, base) in enumerate(zip(wf_x, wf_y, wf_col, wf_base)):
        is_total = i == 0 or i == len(wf_x) - 1
        fig_wf.add_bar(
            x=[x], y=[y], base=[base if not is_total else 0],
            marker_color=col, opacity=0.88,
            text=[_fmt(wf_y[i] if is_total else (wf_y[i] if wf_col[i] == "#1e7e34"
                                                  else -wf_y[i]))],
            textposition="outside", textfont={"size": 11},
            showlegend=False,
        )
    fig_wf.update_layout(
        barmode="stack", height=300,
        margin=dict(l=50, r=50, t=10, b=30),
        plot_bgcolor="white", paper_bgcolor="white",
        yaxis=dict(tickformat="$,.0f", gridcolor="#eee"),
        xaxis=dict(tickfont={"size": 11}),
    )

    ep_fmt = _fmt(ep_total) if not math.isnan(ep_total) else "—"
    ev_fmt = _fmt(ev_total) if not math.isnan(ev_total) else "—"
    kpi_wf = _kpi_row([
        _kpi_card(f"Venta {mes_nombre} 2025", _fmt(v25_total), "#1F3864"),
        _kpi_card(f"Venta {mes_nombre} 2026", _fmt(v26_total), "#2E75B6"),
        _kpi_card("Variación $", _fmt(delta_total),
                  "#1e7e34" if delta_total >= 0 else "#c0392b"),
        _kpi_card("Variación %", f"{pct_caida:+.1f}%",
                  "#1e7e34" if pct_caida >= 0 else "#c0392b"),
        _kpi_card("Efecto Precio",  ep_fmt,
                  "#c0392b" if (not math.isnan(ep_total) and ep_total < 0) else "#1e7e34"),
        _kpi_card("Efecto Volumen", ev_fmt,
                  "#c0392b" if (not math.isnan(ev_total) and ev_total < 0) else "#1e7e34"),
    ])

    seccion_wf = html.Div(style=CARD_STYLE, children=[
        _seccion(f"📉 Puente de Caída — YTD {mes_nombre} 2026 vs 2025", "#c0392b"),
        kpi_wf,
        html.Div(style={"marginTop": "14px"}, children=[
            dcc.Graph(figure=fig_wf, config={"displayModeBar": False}),
        ]),
    ])

    # ── 2. Top clientes en caída (desde PPTO_VS_VENTA) ───────────────────────
    if not dfpv.empty:
        df_cli = dfpv.groupby(["RUT", "NOMBRE"]).agg(
            zona=("VENDEDOR_ACTUAL", "first"),
            v26=("VENTA_2026", "sum"),
            v25=("VENTA_2025", "sum"),
        ).reset_index()
        df_cli["caida"] = df_cli["v26"] - df_cli["v25"]
        df_cli["caida_pct"] = df_cli.apply(
            lambda r: r["caida"] / r["v25"] * 100 if r["v25"] > 0 else float("nan"), axis=1
        )
        df_cli_neg = df_cli[df_cli["caida"] < 0].sort_values("caida").head(20)
    else:
        df_cli_neg = pd.DataFrame()

    if df_cli_neg.empty:
        seccion_cli = html.Div(style=CARD_STYLE, children=[
            _seccion("🏴 Clientes en Caída", "#8B0000"),
            _empty_msg("Sin datos de clientes en caída para este filtro."),
        ])
    else:
        n_cli = len(df_cli_neg)
        total_caida_cli = df_cli_neg["caida"].sum()
        nombres  = [str(n)[:35] for n in df_cli_neg["NOMBRE"].tolist()]
        caidas   = df_cli_neg["caida"].tolist()

        fig_cli = go.Figure(go.Bar(
            x=caidas, y=nombres, orientation="h",
            marker_color="#c0392b", opacity=0.82,
            text=[_fmt(v) for v in caidas],
            textposition="outside", textfont={"size": 10},
        ))
        fig_cli.update_layout(
            height=max(300, n_cli * 26 + 50),
            margin=dict(l=10, r=110, t=10, b=20),
            plot_bgcolor="white", paper_bgcolor="white",
            xaxis=dict(tickformat="$,.0f", gridcolor="#eee",
                       zeroline=True, zerolinecolor="#aaa", zerolinewidth=1.5),
            yaxis=dict(autorange="reversed"),
        )

        cli_rows = []
        for _, r in df_cli_neg.iterrows():
            pct = r["caida_pct"]
            cli_rows.append({
                "Zona":       str(r["zona"])[:22],
                "Cliente":    str(r["NOMBRE"])[:40],
                "Venta 2025": _fmt(r["v25"]),
                "Venta 2026": _fmt(r["v26"]),
                "Caída $":    _fmt(r["caida"]),
                "Caída %":    f"{pct:.1f}%" if not (isinstance(pct, float) and math.isnan(pct)) else "—",
            })

        tabla_cli = dash_table.DataTable(
            data=cli_rows,
            columns=[{"name": c, "id": c} for c in
                     ["Zona", "Cliente", "Venta 2025", "Venta 2026", "Caída $", "Caída %"]],
            style_header=TH,
            style_cell=TC,
            style_cell_conditional=[
                {"if": {"column_id": c}, "textAlign": "left"} for c in ["Zona", "Cliente"]
            ],
            style_data_conditional=[
                {"if": {"column_id": c}, "color": "#c0392b", "fontWeight": "700"}
                for c in ["Caída $", "Caída %"]
            ],
            sort_action="native", page_size=20,
            style_table={"maxHeight": "400px", "overflowY": "auto"},
        )

        seccion_cli = html.Div(style=CARD_STYLE, children=[
            _seccion(f"🏴 Top {n_cli} Clientes en Caída — {_fmt(total_caida_cli)} total", "#8B0000"),
            html.Div(style={"display": "grid", "gridTemplateColumns": "1.3fr 1fr",
                            "gap": "20px", "alignItems": "start"}, children=[
                dcc.Graph(figure=fig_cli, config={"displayModeBar": False}),
                tabla_cli,
            ]),
        ])

    # ── 3. Efecto precio por categoría + top productos afectados ─────────────
    if not df_p.empty and "efecto_precio" in df_p.columns:
        # Gráfico efecto precio+volumen por categoría
        cat_ep = (
            df_p[df_p["categoria"].isin(["EQM", "EVA", "MAH", "SQ"])]
            .groupby("categoria")
            .agg(ep=("efecto_precio", "sum"), ev=("efecto_volumen", "sum"))
            .reset_index()
            .sort_values("ep")
        )

        fig_cat = go.Figure()
        fig_cat.add_bar(
            name="Efecto Precio", x=cat_ep["categoria"].tolist(),
            y=cat_ep["ep"].tolist(),
            marker_color=["#c0392b" if v < 0 else "#1e7e34" for v in cat_ep["ep"]],
            opacity=0.85, text=[_fmt(v) for v in cat_ep["ep"]],
            textposition="outside",
        )
        fig_cat.add_bar(
            name="Efecto Volumen", x=cat_ep["categoria"].tolist(),
            y=cat_ep["ev"].tolist(),
            marker_color=["#e67e22" if v < 0 else "#2E75B6" for v in cat_ep["ev"]],
            opacity=0.75, text=[_fmt(v) for v in cat_ep["ev"]],
            textposition="outside",
        )
        fig_cat.update_layout(
            barmode="group", height=270,
            margin=dict(l=50, r=20, t=10, b=50),
            plot_bgcolor="white", paper_bgcolor="white",
            yaxis=dict(tickformat="$,.0f", gridcolor="#eee",
                       zeroline=True, zerolinecolor="#999"),
            legend=dict(orientation="h", y=-0.3),
        )

        # Top 20 productos con mayor efecto precio negativo
        df_top_ep = (
            df_p[df_p["efecto_precio"].notna() & (df_p["efecto_precio"] < 0)]
            .sort_values("efecto_precio")
            .head(20)
        )

        prod_rows = []
        for _, r in df_top_ep.iterrows():
            dpct = r.get("delta_precio_pct", float("nan"))
            prod_rows.append({
                "Cat":       str(r.get("categoria", ""))[:8],
                "Producto":  str(r.get("descripcion", r.get("codigo", "")))[:45],
                "Cliente":   str(r.get("nombre_cliente", ""))[:28],
                "P. 2025":   _fmt_abs(r["precio_2025"]) if pd.notna(r.get("precio_2025")) else "—",
                "P. 2026":   _fmt_abs(r["precio_2026"]) if pd.notna(r.get("precio_2026")) else "—",
                "Δ Precio %": f"{dpct:.1f}%" if pd.notna(dpct) else "—",
                "Ef. Precio": _fmt(r["efecto_precio"]),
                "Venta 26":  _fmt(r["venta_2026"]),
            })

        tabla_ep = dash_table.DataTable(
            data=prod_rows,
            columns=[{"name": c, "id": c} for c in
                     ["Cat", "Producto", "Cliente", "P. 2025", "P. 2026",
                      "Δ Precio %", "Ef. Precio", "Venta 26"]],
            style_header={**TH, "backgroundColor": "#7b2d8b"},
            style_cell=TC,
            style_cell_conditional=[
                {"if": {"column_id": c}, "textAlign": "left"}
                for c in ["Cat", "Producto", "Cliente"]
            ],
            style_data_conditional=[
                {"if": {"column_id": c}, "color": "#c0392b", "fontWeight": "700"}
                for c in ["Δ Precio %", "Ef. Precio"]
            ],
            sort_action="native", page_size=15,
            style_table={"maxHeight": "400px", "overflowY": "auto", "overflowX": "auto"},
        )

        seccion_precios = html.Div(style=CARD_STYLE, children=[
            _seccion("💸 Efecto Precio y Volumen por Categoría + Top Productos Afectados", "#7b2d8b"),
            html.P(
                "Efecto Precio = Δprecio × cant. vendida 2026 | "
                "Efecto Volumen = Δcantidad × precio 2025 (solo productos con historial real)",
                style={"fontSize": "11px", "color": "#666", "marginBottom": "10px"}),
            dcc.Graph(figure=fig_cat, config={"displayModeBar": False}),
            html.Div(style={"marginTop": "14px"}, children=[
                _seccion("Top 20 Productos con Mayor Impacto Negativo de Precio", "#555"),
                tabla_ep,
            ]),
        ])
    else:
        seccion_precios = html.Div(style=CARD_STYLE, children=[
            _seccion("💸 Efecto Precio", "#7b2d8b"),
            _empty_msg("Sin datos de precios. Presione 🔄 Actualizar."),
        ])

    return html.Div([seccion_wf, seccion_cli, seccion_precios])


# ─── Inicialización ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("=" * 60, flush=True)
    # _load_from_cache() ya se llamó al nivel de módulo.
    # Recargar desde BD si no hay caché completo o faltan los KPIs de PPTO/Venta
    _needs_reload = _df_categoria.empty or _df_ppto_vs_venta.empty or not _pvv_kpis
    if _needs_reload and CONN_STR:
        print("[INFO] Caché incompleto — cargando desde BD...")
        _reload_all_data()
    elif _needs_reload:
        print("[WARN] Sin caché completo y sin acceso a BD. Presione Actualizar con VPN.")
    print(f"[INFO] Iniciando app en http://localhost:8052", flush=True)
    app.run(host="0.0.0.0", port=8052, debug=False, use_reloader=False)
